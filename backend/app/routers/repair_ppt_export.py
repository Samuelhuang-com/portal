"""
報修模組 PPT 匯出 — 獨立 Router
======================================
路由前綴：/repair/ppt-export
  POST /export  → 觸發匯出，回傳 .pptx StreamingResponse

支援模組：
  module="dazhi"  → 大直工務部
  module="luqun"  → 盧群商場工務報修

投影片順序（固定）：
  0. Cover（封面，更新年月）
  1. 3.1 報修統計
  2. 3.2 結案時間
  3. 3.3 柏拉圖分析（matplotlib 圖）
  4. 3.3 報修類型各類別（表格）
  5. Dashboard 報修類型分布（matplotlib 圓餅）
  6. 3.4 本月客房報修表（≤20 列/頁）
  7. 報修金額統計
  8. 未完成附表－飯店（≤20 列/頁）
"""

import logging
import os
from datetime import datetime
from io import BytesIO
from typing import Optional
from urllib.parse import quote

from fastapi import APIRouter, Body, Depends, HTTPException
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
from sqlalchemy.orm import Session

from app.core.database import get_db
from app.dependencies import get_current_user
from app.models.user import User

logger = logging.getLogger("repair_ppt_export")

router = APIRouter(prefix="/repair/ppt-export", tags=["報修 PPT 匯出"])

TEMPLATE_PATH = os.path.abspath(
    os.path.join(os.path.dirname(__file__), "..", "static",
                 "pptx_templates", "hotel_report_template.pptx")
)

# dazhi (hotel) module uses IHG Fion template
IHG_TEMPLATE_PATH = os.path.abspath(
    os.path.join(os.path.dirname(__file__), "..", "static",
                 "pptx_templates", "hotel_report_template_ihg.pptx")
)

ROWS_PER_SLIDE = 14   # 每頁最多行數（12pt 字 × 26pt 列高，確保不超出頁面）


# ═══════════════════════════════════════════════════════
# Request schema
# ═══════════════════════════════════════════════════════

class RepairPptBody(BaseModel):
    module: str   # "dazhi" | "luqun"
    year:   int
    month:  int   # 報表月份（客房報修表 / 未完成附表用）


# ═══════════════════════════════════════════════════════
# Internal helpers — matplotlib charts
# ═══════════════════════════════════════════════════════

def _cjk_font() -> Optional[str]:
    """回傳可用的 CJK 字型名稱，找不到回傳 None。"""
    try:
        import matplotlib.font_manager as fm
        _PRIORITY = [
            "Microsoft JhengHei", "Microsoft JhengHei UI",
            "PMingLiU", "MingLiU", "DFKai-SB",
            "Microsoft YaHei", "Microsoft YaHei UI", "SimHei", "SimSun",
            "PingFang TC", "PingFang SC", "Heiti TC",
            "Noto Sans CJK TC", "Noto Sans CJK SC", "WenQuanYi Micro Hei",
        ]
        available = {f.name for f in fm.fontManager.ttflist}
        found = next((fn for fn in _PRIORITY if fn in available), None)
        if found:
            return found
        keywords = ("jhenghei", "mingliu", "dfkai", "yahei", "simhei",
                    "pingfang", "heiti", "noto", "cjk", "wenquanyi",
                    "chinese", "gothic", "mincho")
        return next(
            (f.name for f in fm.fontManager.ttflist
             if any(k in f.name.lower() for k in keywords)),
            None,
        )
    except Exception:
        return None


def _setup_mpl():
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    font = _cjk_font()
    if font:
        plt.rcParams["font.family"] = font
    plt.rcParams["axes.unicode_minus"] = False
    return plt


def _make_pareto_chart(type_rows: list) -> Optional[BytesIO]:
    """柏拉圖：Bar（深/淺藍）+ 累計折線 + 80% 虛線，回傳 PNG BytesIO。"""
    try:
        plt = _setup_mpl()
        # 排序、過濾 0 件
        rows = sorted([r for r in type_rows if r.get("row_total", 0) > 0],
                      key=lambda r: -r["row_total"])
        if not rows:
            return None
        labels  = [r["type"] for r in rows]
        counts  = [r["row_total"] for r in rows]
        total   = sum(counts)
        running = 0
        cum_pct = []
        for c in counts:
            running += c
            cum_pct.append(round(running / total * 100) if total else 0)

        C_BRAND  = "#1B3A5C"
        C_LIGHT  = "#4BA8E8"
        C_RED    = "#c0392b"

        fig, ax1 = plt.subplots(figsize=(10.5, 3.8), dpi=130)
        fig.patch.set_facecolor("white")

        bar_colors = [C_BRAND if p <= 80 else C_LIGHT for p in cum_pct]
        bars = ax1.bar(range(len(labels)), counts, color=bar_colors, width=0.65)
        ax1.set_xticks(range(len(labels)))
        ax1.set_xticklabels(labels, fontsize=8, rotation=30, ha="right")
        ax1.set_ylabel("件數", fontsize=8, color=C_BRAND)
        ax1.tick_params(axis="y", labelsize=8)
        ax1.set_ylim(0, max(counts) * 1.3)
        # Bar 數字標籤
        for bar, v in zip(bars, counts):
            ax1.text(bar.get_x() + bar.get_width() / 2,
                     bar.get_height() + max(counts) * 0.02,
                     str(v), ha="center", va="bottom", fontsize=7, color=C_BRAND)

        ax2 = ax1.twinx()
        ax2.plot(range(len(labels)), cum_pct, color=C_RED,
                 marker="o", markersize=4, linewidth=2, zorder=5)
        ax2.set_ylabel("累計 %", fontsize=8, color=C_RED)
        ax2.set_ylim(0, 110)
        ax2.tick_params(axis="y", labelsize=8, colors=C_RED)
        # 折線數字標籤
        for i, p in enumerate(cum_pct):
            ax2.text(i, p + 3, f"{p}%", ha="center", va="bottom",
                     fontsize=7, color=C_RED, fontweight="bold")
        # 80% 虛線
        ax2.axhline(y=80, color=C_RED, linestyle="--", linewidth=1.2, alpha=0.7)
        ax2.text(len(labels) - 0.5, 82, "80%", color=C_RED, fontsize=8)

        ax1.grid(axis="y", linestyle="--", linewidth=0.4, color="#dddddd", zorder=0)
        ax1.set_title("報修類型柏拉圖分析", fontsize=10, fontweight="bold",
                      color=C_BRAND, pad=6)

        # 圖例色塊
        from matplotlib.patches import Patch
        legend_elements = [
            Patch(facecolor=C_BRAND, label="累計 ≤ 80%（重點）"),
            Patch(facecolor=C_LIGHT, label="累計 > 80%（次要）"),
        ]
        ax1.legend(handles=legend_elements, fontsize=7,
                   loc="upper right", framealpha=0.7)

        fig.tight_layout(pad=0.8)
        buf = BytesIO()
        fig.savefig(buf, format="png", dpi=130, bbox_inches="tight")
        buf.seek(0)
        plt.close(fig)
        return buf
    except Exception as e:
        logger.warning("Pareto chart error: %s", e, exc_info=True)
        return None


def _make_pie_chart(type_rows: list) -> Optional[BytesIO]:
    """
    類型分布圓餅圖，回傳 PNG BytesIO。
    百分比標籤直接取 service 計算的 cum_pct（= row_total / len(year_cases) * 100），
    與前端 3.3 報修類型 TAB「年度佔比」欄口徑完全一致。
    """
    try:
        plt = _setup_mpl()
        rows = [r for r in type_rows if r.get("row_total", 0) > 0]
        if not rows:
            return None
        rows.sort(key=lambda r: -r["row_total"])
        labels   = [r["type"]                    for r in rows]
        sizes    = [r["row_total"]               for r in rows]
        cum_pcts = [r.get("cum_pct", 0)          for r in rows]   # 與 3.3 TAB 同口徑
        COLORS = [
            "#1B3A5C", "#4BA8E8", "#F5A623", "#2ECC71", "#E74C3C",
            "#9B59B6", "#1ABC9C", "#E67E22", "#3498DB", "#34495E",
            "#F39C12", "#27AE60", "#8E44AD", "#D35400",
        ]
        colors = [COLORS[i % len(COLORS)] for i in range(len(labels))]

        LABEL_THRESHOLD = 5.0   # cum_pct >= 5% 才顯示外部標籤

        fig, ax = plt.subplots(figsize=(9, 4.5), dpi=130)
        fig.patch.set_facecolor("white")

        # 外部類別名稱：僅大 slice 顯示
        outer_labels = [
            lb if pct >= LABEL_THRESHOLD else ""
            for lb, pct in zip(labels, cum_pcts)
        ]
        # autopct=None：改為事後用 cum_pct 手動標注
        wedges, texts = ax.pie(
            sizes,
            labels=outer_labels,
            colors=colors,
            pctdistance=0.72,
            labeldistance=1.13,
            startangle=90,
            textprops={"fontsize": 7.5},
        )
        for txt in texts:
            txt.set_fontsize(7.5)

        # 手動在 slice 內側標注 cum_pct（與 TAB 年度佔比一致）
        import numpy as np
        _start = 90.0
        for wedge, pct in zip(wedges, cum_pcts):
            if pct < LABEL_THRESHOLD:
                continue
            angle = (wedge.theta1 + wedge.theta2) / 2
            x = 0.72 * np.cos(np.radians(angle))
            y = 0.72 * np.sin(np.radians(angle))
            ax.text(x, y, f"{pct}%", ha="center", va="center",
                    fontsize=7, color="white", fontweight="bold")

        # 圖例：顯示件數 + 年度佔比（與 TAB 欄位對齊）
        ax.legend(
            wedges,
            [f"{lb}（{sz}件，{pct}%）" for lb, sz, pct in zip(labels, sizes, cum_pcts)],
            loc="center left",
            bbox_to_anchor=(1.0, 0.5),
            fontsize=7.5,
            framealpha=0.8,
        )
        ax.set_title("報修類型分布", fontsize=10, fontweight="bold",
                     color="#1B3A5C", pad=6)
        ax.axis("equal")
        fig.tight_layout(pad=0.5)
        buf = BytesIO()
        fig.savefig(buf, format="png", dpi=130, bbox_inches="tight")
        buf.seek(0)
        plt.close(fig)
        return buf
    except Exception as e:
        logger.warning("Pie chart error: %s", e, exc_info=True)
        return None


# ═══════════════════════════════════════════════════════
# Internal helpers — slide builders
# ═══════════════════════════════════════════════════════

def _sanitize(val) -> str:
    s = str(val) if val is not None else "—"
    return (s
            .replace("_x000D_", " ")
            .replace('\r\n', ' ').replace('\r', ' ')
            .replace('\n', ' ').replace('\t', ' '))



def _fmt_rate(v) -> str:
    return f"{v}%" if v is not None else "—"


def _fmt_money(v) -> str:
    if v is None or v == 0:
        return "—"
    return f"${v:,.0f}"


def _add_chart_slide(prs, template_idx: int,
                     title: str, subtitle: str,
                     chart_buf: Optional[BytesIO],
                     now_str: str, SW: float, SH: float,
                     title_fn=None):
    """Clone template, set title, embed chart image."""
    from app.routers.hotel_overview import (
        _clone_template_slide, _set_slide_title, _pptx_txt,
    )
    from pptx.util import Inches
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    _fn = title_fn or _set_slide_title

    C_BLACK = RGBColor(0x00, 0x00, 0x00)
    C_GRAY  = RGBColor(0x88, 0x88, 0x88)

    slide = _clone_template_slide(prs, template_idx)
    _fn(slide, title, "", now_str, SW, SH)
    if subtitle:
        _chart_sub_y = 0.592 if title_fn is not None else 0.52
        _pptx_txt(slide, subtitle, 0.35, _chart_sub_y, SW - 4.5, 0.22,
                  size=10, bold=True, color=C_BLACK)

    if chart_buf:
        IMG_W = SW - 1.0
        IMG_H = SH - 1.8
        IMG_X = 0.5
        IMG_Y = 1.05
        slide.shapes.add_picture(chart_buf, Inches(IMG_X), Inches(IMG_Y),
                                 Inches(IMG_W), Inches(IMG_H))
    else:
        _pptx_txt(slide, "（本期暫無圖表資料）",
                  2.0, 3.5, 9.0, 1.0, size=14, color=C_GRAY, italic=True)
    return slide


def _set_cell_bottom_border(cell, width_pt: float = 2.5, color_hex: str = "1B3A5C"):
    """在 pptx 表格儲存格底部加粗線（lxml XML 操作）。
    OOXML 規範：a:lnB 必須在 tcPr 內的 fill 元素之前，
    所以用 insert(0, ...) 插到最前面而非 SubElement 追加末尾。
    """
    from lxml import etree
    from pptx.oxml.ns import qn as _qn
    width_emu = int(width_pt * 12700)
    tc = cell._tc

    # 取得或建立 tcPr（必須在 txBody 之後）
    tcPr = tc.find(_qn('a:tcPr'))
    if tcPr is None:
        # 插在 txBody 之後
        txBody = tc.find(_qn('a:txBody'))
        idx = list(tc).index(txBody) + 1 if txBody is not None else len(tc)
        tcPr = etree.Element(_qn('a:tcPr'))
        tc.insert(idx, tcPr)

    # 移除已有 lnB
    for existing in tcPr.findall(_qn('a:lnB')):
        tcPr.remove(existing)

    # 建構 lnB 元素
    lnB = etree.Element(_qn('a:lnB'))
    lnB.set('w', str(width_emu))
    lnB.set('cap', 'flat')
    lnB.set('cmpd', 'sng')
    solidFill = etree.SubElement(lnB, _qn('a:solidFill'))
    srgbClr   = etree.SubElement(solidFill, _qn('a:srgbClr'))
    srgbClr.set('val', color_hex)

    # 插到 tcPr 最前面（border 元素必須先於 fill 元素，符合 OOXML schema 順序）
    tcPr.insert(0, lnB)


def _add_table_slides(prs, template_idx: int,
                      title: str, subtitle: str,
                      columns: list, rows: list,
                      now_str: str, SW: float, SH: float,
                      max_rows: int = ROWS_PER_SLIDE,
                      sep_after_rows: list = None,
                      title_fn=None):
    """Clone 1+ slides, build generic table (≤ max_rows per page).
    sep_after_rows: 0-based data row indices (not counting header) that get a thick bottom border.
    """
    from app.routers.hotel_overview import (
        _clone_template_slide, _set_slide_title,
        _pptx_cell, _pptx_header_row, _pptx_txt,
    )
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    _fn = title_fn or _set_slide_title

    C_DARK    = RGBColor(0x1B, 0x3A, 0x5C)
    C_ROW_ALT = RGBColor(0xEE, 0xF5, 0xFB)
    C_GRAY    = RGBColor(0x88, 0x88, 0x88)
    _ALIGN    = {"right": PP_ALIGN.RIGHT, "center": PP_ALIGN.CENTER}

    TABLE_Y = 0.95
    TABLE_H = SH - TABLE_Y - 0.45
    TABLE_W = SW - 0.8

    pages = [rows[i:i + max_rows] for i in range(0, max(len(rows), 1), max_rows)]
    total_pages = len(pages)

    for pg_idx, pg_rows in enumerate(pages):
        slide = _clone_template_slide(prs, template_idx)
        pg_sub = (subtitle
                  if total_pages == 1
                  else f"{subtitle}　第 {pg_idx + 1} 頁／共 {total_pages} 頁　（共 {len(rows)} 筆）")
        _fn(slide, title, "", now_str, SW, SH)
        if pg_sub:
            _sub_y = 0.592 if title_fn is not None else 0.52
            _pptx_txt(slide, pg_sub, 0.35, _sub_y, SW - 4.5, 0.22,
                      size=10, bold=True, color=RGBColor(0x00, 0x00, 0x00))

        if not pg_rows or not columns:
            _pptx_txt(slide, "（本期暫無資料）",
                      2.0, 3.5, 9.0, 1.0, size=14, color=C_GRAY, italic=True)
            continue

        n_cols = len(columns)
        n_rows = len(pg_rows) + 1
        tbl = slide.shapes.add_table(
            n_rows, n_cols,
            Inches(0.4), Inches(TABLE_Y), Inches(TABLE_W), Inches(TABLE_H)
        ).table

        used = 0.0
        for ci, col in enumerate(columns):
            w = col.get("width", TABLE_W / n_cols)
            if ci == n_cols - 1:
                w = max(TABLE_W - used, 0.3)
            tbl.columns[ci].width = Inches(w)
            used += w

        for ci, col in enumerate(columns):
            _pptx_cell(tbl, 0, ci, col.get("label", col["key"]), bold=True)
        _pptx_header_row(tbl, n_cols, size=12)

        for ri, row in enumerate(pg_rows, 1):
            # 支援列級別顏色覆寫（_row_bg / _row_fg），不覆寫則用預設交替色
            _override_bg = row.get("_row_bg")
            _override_fg = row.get("_row_fg")
            bg  = _override_bg if _override_bg is not None else (C_ROW_ALT if ri % 2 == 0 else None)
            fg  = _override_fg if _override_fg is not None else C_DARK
            for ci, col in enumerate(columns):
                raw = _sanitize(row.get(col["key"], ""))
                align = _ALIGN.get(col.get("align", "left"))
                # 第一欄（統計項目）：使用 bold，其餘正常
                is_first = (ci == 0)
                _pptx_cell(tbl, ri, ci, raw, fg=fg, bg=bg, size=12,
                           align=align, bold=is_first)

        tbl.rows[0].height = Pt(34)
        for ri in range(1, n_rows):
            tbl.rows[ri].height = Pt(26)

        # 分隔粗線：在指定資料列底部加粗線（0-based，不含 header）
        if sep_after_rows:
            # 換算到本頁的區域索引：pg_rows 是整體 rows 的切片
            page_start = pg_idx * max_rows
            for sep_idx in sep_after_rows:
                local_idx = sep_idx - page_start  # 在本頁中的 0-based 位置
                tbl_row_idx = local_idx + 1         # +1 for header row
                if 1 <= tbl_row_idx < n_rows:
                    for ci in range(n_cols):
                        _set_cell_bottom_border(tbl.cell(tbl_row_idx, ci),
                                                width_pt=2.5, color_hex="1B3A5C")


# ═══════════════════════════════════════════════════════
# Annual matrix helpers — 年度計劃表投影片
# ═══════════════════════════════════════════════════════

# 狀態 → (符號, bg RGBColor, fg RGBColor)
_MATRIX_STATUS_STYLE = None   # 延遲初始化（需 pptx 匯入）

def _matrix_status_style():
    """回傳狀態樣式對照表（lazy init）。"""
    global _MATRIX_STATUS_STYLE
    if _MATRIX_STATUS_STYLE is not None:
        return _MATRIX_STATUS_STYLE
    from pptx.dml.color import RGBColor as _R
    _MATRIX_STATUS_STYLE = {
        "completed":    ("✓", _R(0xE8, 0xF5, 0xE9), _R(0x27, 0x7D, 0x27)),  # 綠底綠字
        "in_progress":  ("○", _R(0xFF, 0xFB, 0xE6), _R(0xD4, 0x88, 0x06)),  # 黃底橙字
        "scheduled":    ("○", _R(0xFF, 0xFB, 0xE6), _R(0xD4, 0x88, 0x06)),  # 同上
        "unscheduled":  ("△", _R(0xFF, 0xF4, 0xE6), _R(0xFA, 0x8C, 0x16)),  # 淡橙底橙字
        "overdue":      ("✗", _R(0xFF, 0xF2, 0xF0), _R(0xCF, 0x13, 0x22)),  # 紅底紅字
        "no_data":      ("?", _R(0xFA, 0xFA, 0xFA), _R(0xAA, 0xAA, 0xAA)),  # 灰底灰字
        "non_month":    ("",  None,                  _R(0xDD, 0xDD, 0xDD)),  # 空白
        "no_frequency": ("",  None,                  _R(0xDD, 0xDD, 0xDD)),  # 空白
    }
    return _MATRIX_STATUS_STYLE


def _get_annual_matrix_rows(module: str, year: int, db) -> list:
    """
    依模組取得年度計劃表資料，回傳統一格式的 list。
    每筆: {"source": str, "category": str, "task_name": str,
           "frequency": str, "cells": [{month, status}, ...]}
    dazhi   → hotel periodic-maintenance (PMSchedule)
    luqun   → mall periodic-maintenance (MallPMSchedule)
              + full-building-maintenance (FullBldgPMSchedule)
    """
    results = []

    if module == "dazhi":
        # ── 飯店週期保養 ───────────────────────────────────────────────────
        from app.routers.periodic_maintenance import (
            _get_latest_batch_items, _calc_schedule_status, _should_schedule,
        )
        from app.models.pm_schedule import PMSchedule
        items = _get_latest_batch_items(db)
        year_recs = (
            db.query(PMSchedule)
            .filter(PMSchedule.year_month.like(f"{year}/%"))
            .all()
        )
        smap = {}
        for r in year_recs:
            try:
                m = int(r.year_month.split("/")[1])
                smap[(r.item_ragic_id, m)] = r
            except Exception:
                pass
        for item in items:
            cells = []
            for m in range(1, 13):
                rec = smap.get((item.ragic_id, m))
                if rec:
                    cells.append({"month": m, "status": _calc_schedule_status(rec)})
                elif not (item.frequency or "").strip():
                    cells.append({"month": m, "status": "no_frequency"})
                elif _should_schedule(item, year, m):
                    cells.append({"month": m, "status": "no_data"})
                else:
                    cells.append({"month": m, "status": "non_month"})
            results.append({
                "source":    "飯店週期保養",
                "category":  item.category,
                "task_name": item.task_name,
                "frequency": item.frequency or "",
                "cells":     cells,
            })

    else:  # luqun — 商場週期保養 + 全棟例行維護
        # ── 商場週期保養 ───────────────────────────────────────────────────
        from app.routers.mall_periodic_maintenance import (
            _mall_get_latest_batch_items, _mall_calc_schedule_status,
            _should_schedule as _mall_should_schedule,
        )
        from app.models.mall_pm_schedule import MallPMSchedule
        mall_items = _mall_get_latest_batch_items(db)
        mall_recs = (
            db.query(MallPMSchedule)
            .filter(MallPMSchedule.year_month.like(f"{year}/%"))
            .all()
        )
        mall_smap = {}
        for r in mall_recs:
            try:
                m = int(r.year_month.split("/")[1])
                mall_smap[(r.item_ragic_id, m)] = r
            except Exception:
                pass
        for item in mall_items:
            cells = []
            for m in range(1, 13):
                rec = mall_smap.get((item.ragic_id, m))
                if rec:
                    cells.append({"month": m, "status": _mall_calc_schedule_status(rec)})
                elif not (item.frequency or "").strip():
                    cells.append({"month": m, "status": "no_frequency"})
                elif _mall_should_schedule(item, year, m):
                    cells.append({"month": m, "status": "no_data"})
                else:
                    cells.append({"month": m, "status": "non_month"})
            results.append({
                "source":    "商場週期保養",
                "category":  item.category,
                "task_name": item.task_name,
                "frequency": item.frequency or "",
                "cells":     cells,
            })

        # ── 全棟例行維護 ───────────────────────────────────────────────────
        from app.routers.full_building_maintenance import (
            _fb_get_latest_batch_items, _fb_calc_schedule_status,
            _fb_should_schedule_by_freq,
        )
        from app.models.full_bldg_pm_schedule import FullBldgPMSchedule
        fb_items = _fb_get_latest_batch_items(db)
        fb_recs = (
            db.query(FullBldgPMSchedule)
            .filter(FullBldgPMSchedule.year_month.like(f"{year}/%"))
            .all()
        )
        fb_smap = {}
        for r in fb_recs:
            try:
                m = int(r.year_month.split("/")[1])
                fb_smap[(r.item_ragic_id, m)] = r
            except Exception:
                pass
        for item in fb_items:
            cells = []
            for m in range(1, 13):
                rec = fb_smap.get((item.ragic_id, m))
                if rec:
                    cells.append({"month": m, "status": _fb_calc_schedule_status(rec)})
                elif not (item.frequency or "").strip():
                    cells.append({"month": m, "status": "no_frequency"})
                elif _fb_should_schedule_by_freq(item.frequency or "", m):
                    cells.append({"month": m, "status": "no_data"})
                else:
                    cells.append({"month": m, "status": "non_month"})
            results.append({
                "source":    "全棟例行維護",
                "category":  item.category,
                "task_name": item.task_name,
                "frequency": item.frequency or "",
                "cells":     cells,
            })

    return results


def _add_annual_matrix_slide(prs, template_idx: int,
                              title: str, subtitle: str,
                              matrix_rows: list,
                              now_str: str, SW: float, SH: float,
                              title_fn=None):
    """
    年度計劃表矩陣投影片。
    每格依狀態顯示符號（✓/○/✗/△/?）+ 對應底色。
    超過 ROWS_PER_SLIDE 自動分頁。
    """
    from app.routers.hotel_overview import (
        _clone_template_slide, _set_slide_title, _pptx_cell,
        _pptx_header_row, _pptx_txt,
    )
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    _fn = title_fn or _set_slide_title   # IHG: _set_ihg_slide_title

    C_DARK    = RGBColor(0x1B, 0x3A, 0x5C)
    C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
    C_ROW_ALT = RGBColor(0xF7, 0xF9, 0xFC)
    C_GRAY    = RGBColor(0x88, 0x88, 0x88)
    STYLE     = _matrix_status_style()

    CELL_SIZE  = 12   # 字體 pt（與其他內頁一致）
    TABLE_Y    = 0.95
    TABLE_H    = SH - TABLE_Y - 0.45
    TABLE_W    = SW - 0.8

    # 欄位寬度：類別(0.90) + 項目名稱(2.40) + 頻率(0.70) + 12×月份(auto)
    W_CAT   = 0.90
    W_TASK  = 2.40
    W_FREQ  = 0.70
    W_AVAIL = TABLE_W - W_CAT - W_TASK - W_FREQ
    W_MON   = round(W_AVAIL / 12, 3)   # 每月欄寬（最後一欄 auto）

    columns = [
        {"key": "category",  "label": "類別",   "width": W_CAT,  "align": "left"},
        {"key": "task_name", "label": "項目名稱","width": W_TASK, "align": "left"},
        {"key": "frequency", "label": "頻率",    "width": W_FREQ, "align": "center"},
    ]
    for m in range(1, 13):
        columns.append({"key": f"m{m}", "label": f"{m}月", "width": W_MON, "align": "center"})

    pages = [
        matrix_rows[i:i + ROWS_PER_SLIDE]
        for i in range(0, max(len(matrix_rows), 1), ROWS_PER_SLIDE)
    ]
    total_pages = len(pages)

    for pg_idx, pg_rows in enumerate(pages):
        slide = _clone_template_slide(prs, template_idx)
        pg_sub = (subtitle if total_pages == 1
                  else f"{subtitle}　第 {pg_idx + 1} 頁／共 {total_pages} 頁")
        _fn(slide, title, "", now_str, SW, SH)
        if pg_sub:
            _sub_y2 = 0.592 if title_fn is not None else 0.52
            _pptx_txt(slide, pg_sub, 0.35, _sub_y2, SW - 4.5, 0.22,
                      size=10, bold=True, color=RGBColor(0x00, 0x00, 0x00))

        if not pg_rows:
            _pptx_txt(slide, "（本期暫無保養排程資料）",
                      2.0, 3.5, 9.0, 1.0, size=14, color=C_GRAY, italic=True)
            continue

        n_cols = len(columns)
        n_rows = len(pg_rows) + 1
        tbl = slide.shapes.add_table(
            n_rows, n_cols,
            Inches(0.4), Inches(TABLE_Y), Inches(TABLE_W), Inches(TABLE_H)
        ).table

        used = 0.0
        for ci, col in enumerate(columns):
            w = col.get("width", TABLE_W / n_cols)
            if ci == n_cols - 1:
                w = max(TABLE_W - used, 0.3)
            tbl.columns[ci].width = Inches(w)
            used += w

        # Header row
        for ci, col in enumerate(columns):
            _pptx_cell(tbl, 0, ci, col["label"], bold=True)
        _pptx_header_row(tbl, n_cols, size=CELL_SIZE)

        # Data rows
        _ALIGN_MAP = {"right": PP_ALIGN.RIGHT, "center": PP_ALIGN.CENTER}
        for ri, row in enumerate(pg_rows, 1):
            row_bg = C_ROW_ALT if ri % 2 == 0 else None
            for ci, col in enumerate(columns):
                key = col["key"]
                align = _ALIGN_MAP.get(col.get("align", "left"))

                if key.startswith("m") and key[1:].isdigit():
                    m = int(key[1:])
                    cell_info = next(
                        (c for c in row.get("cells", []) if c["month"] == m), None
                    )
                    status = cell_info["status"] if cell_info else "non_month"
                    sym, s_bg, s_fg = STYLE.get(status, ("", None, C_GRAY))
                    bg = s_bg if s_bg is not None else row_bg
                    _pptx_cell(tbl, ri, ci, sym,
                               fg=s_fg, bg=bg, size=CELL_SIZE,
                               align=PP_ALIGN.CENTER, bold=(status == "completed"))
                else:
                    text = _sanitize(row.get(key, ""))
                    _pptx_cell(tbl, ri, ci, text,
                               fg=C_DARK, bg=row_bg, size=CELL_SIZE,
                               align=align, bold=(ci == 0))

        tbl.rows[0].height = Pt(34)
        for ri in range(1, n_rows):
            tbl.rows[ri].height = Pt(26)


# ═══════════════════════════════════════════════════════
# 商場週期保養統計 helpers（luqun 專用）
# ═══════════════════════════════════════════════════════

def _make_mall_pm_stats_table(
    db, year: int, freq_type: str
) -> tuple[list[dict], list[dict]]:
    """
    商場週期保養年度統計 — 橫向格式（統計項目為列，月份為欄）。

    - 固定顯示全年 12 個月欄位
    - 無資料的月份（all zeros）欄位填 "—"
    - 回傳 (cols, rows)：直接傳給 _add_table_slides()

    freq_type: "monthly" | "quarterly" | "yearly"
    """
    from app.routers.mall_periodic_maintenance import _calc_year_matrix

    matrix = _calc_year_matrix(db, year, freq_type)
    all_months = matrix.months  # 固定 12 個月，不過濾

    # 若全年完全無資料，回傳空（整張投影片跳過）
    if not any(m.period_total > 0 or m.prev_carry_over > 0 for m in all_months):
        return [], []

    # 判斷每個月是否有意義的資料（有資料 = 顯示實際值；無資料 = 顯示 "—"）
    def _has_data(m) -> bool:
        return (m.period_total > 0
                or m.prev_carry_over > 0
                or m.prev_resolved_in_period > 0)

    # ── 欄寬：固定 12 月份，統計項目欄 2.20"，月份欄平均，合計欄自動填滿 ────
    TABLE_W = 12.50
    ITEM_W  = 2.20
    month_w = max(0.65, min(0.90, (TABLE_W - ITEM_W - 1.00) / 12))

    cols: list[dict] = [{"key": "item", "label": "統計項目", "width": ITEM_W, "align": "left"}]
    for m in all_months:
        cols.append({"key": f"m{m.month}", "label": f"{m.month}月",
                     "width": month_w, "align": "center"})
    cols.append({"key": "total", "label": "合計", "align": "center"})

    # ── 合計：僅加總有資料的月份 ───────────────────────────────────────────
    active = [m for m in all_months if _has_data(m)]
    sum_resolve = sum(m.prev_resolved_in_period for m in active)
    sum_total   = sum(m.period_total            for m in active)
    sum_done    = sum(m.period_completed        for m in active)
    full_rate   = (
        f"{round(sum_done / sum_total * 100, 1):.1f}%"
        if sum_total else "—"
    )

    def _r(label: str, fn, total_val) -> dict:
        """fn(m) 被呼叫時，無資料月份傳回 '—'。"""
        row: dict = {"item": label}
        for m in all_months:
            row[f"m{m.month}"] = fn(m) if _has_data(m) else "—"
        row["total"] = total_val
        return row

    rows: list[dict] = [
        _r("截至上月底累計未結案數",
           lambda m: str(m.prev_carry_over),
           "—"),
        _r("其中本月已結案數",
           lambda m: str(m.prev_resolved_in_period),
           str(sum_resolve)),
        _r("累計項目完成率",
           lambda m: (f"{m.carry_over_rate:.1f}%" if m.carry_over_rate is not None else "—"),
           "—"),
        _r("本月週期保養項目數",
           lambda m: str(m.period_total),
           str(sum_total)),
        _r("本月週期保養完成數",
           lambda m: str(m.period_completed),
           str(sum_done)),
        _r("本月週期保養完成率",
           lambda m: (f"{m.period_rate:.1f}%" if m.period_rate is not None else "—"),
           full_rate),
    ]
    return cols, rows


_MALL_PM_FREQ_LABELS = {
    "monthly":   "每月維護",
    "quarterly": "每季維護",
    "yearly":    "每年維護",
}



def _make_fb_pm_stats_table(
    db, year: int, freq_type: str
) -> tuple[list[dict], list[dict]]:
    """
    # Full-building-maintenance version of _make_mall_pm_stats_table.
    商場週期保養年度統計 — 橫向格式（統計項目為列，月份為欄）。

    - 固定顯示全年 12 個月欄位
    - 無資料的月份（all zeros）欄位填 "—"
    - 回傳 (cols, rows)：直接傳給 _add_table_slides()

    freq_type: "monthly" | "quarterly" | "yearly"
    """
    from app.routers.full_building_maintenance import _calc_year_matrix

    matrix = _calc_year_matrix(db, year, freq_type)
    all_months = matrix.months  # 固定 12 個月，不過濾

    # 若全年完全無資料，回傳空（整張投影片跳過）
    if not any(m.period_total > 0 or m.prev_carry_over > 0 for m in all_months):
        return [], []

    # 判斷每個月是否有意義的資料（有資料 = 顯示實際值；無資料 = 顯示 "—"）
    def _has_data(m) -> bool:
        return (m.period_total > 0
                or m.prev_carry_over > 0
                or m.prev_resolved_in_period > 0)

    # ── 欄寬：固定 12 月份，統計項目欄 2.20"，月份欄平均，合計欄自動填滿 ────
    TABLE_W = 12.50
    ITEM_W  = 2.20
    month_w = max(0.65, min(0.90, (TABLE_W - ITEM_W - 1.00) / 12))

    cols: list[dict] = [{"key": "item", "label": "統計項目", "width": ITEM_W, "align": "left"}]
    for m in all_months:
        cols.append({"key": f"m{m.month}", "label": f"{m.month}月",
                     "width": month_w, "align": "center"})
    cols.append({"key": "total", "label": "合計", "align": "center"})

    # ── 合計：僅加總有資料的月份 ───────────────────────────────────────────
    active = [m for m in all_months if _has_data(m)]
    sum_resolve = sum(m.prev_resolved_in_period for m in active)
    sum_total   = sum(m.period_total            for m in active)
    sum_done    = sum(m.period_completed        for m in active)
    full_rate   = (
        f"{round(sum_done / sum_total * 100, 1):.1f}%"
        if sum_total else "—"
    )

    def _r(label: str, fn, total_val) -> dict:
        """fn(m) 被呼叫時，無資料月份傳回 '—'。"""
        row: dict = {"item": label}
        for m in all_months:
            row[f"m{m.month}"] = fn(m) if _has_data(m) else "—"
        row["total"] = total_val
        return row

    rows: list[dict] = [
        _r("截至上月底累計未結案數",
           lambda m: str(m.prev_carry_over),
           "—"),
        _r("其中本月已結案數",
           lambda m: str(m.prev_resolved_in_period),
           str(sum_resolve)),
        _r("累計項目完成率",
           lambda m: (f"{m.carry_over_rate:.1f}%" if m.carry_over_rate is not None else "—"),
           "—"),
        _r("本月週期保養項目數",
           lambda m: str(m.period_total),
           str(sum_total)),
        _r("本月週期保養完成數",
           lambda m: str(m.period_completed),
           str(sum_done)),
        _r("本月週期保養完成率",
           lambda m: (f"{m.period_rate:.1f}%" if m.period_rate is not None else "—"),
           full_rate),
    ]
    return cols, rows



def _add_slide_hyperlink(slide, run, target_slide):
    """Add internal PPT hyperlink to run that jumps to target_slide."""
    from lxml import etree
    from pptx.oxml.ns import qn
    rId = slide.part.relate_to(
        target_slide.part,
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide",
    )
    r_elem = run._r
    rPr = r_elem.find(qn("a:rPr"))
    if rPr is None:
        rPr = etree.Element(qn("a:rPr"))
        r_elem.insert(0, rPr)
    for old in rPr.findall(qn("a:hlinkClick")):
        rPr.remove(old)
    hlinkClick = etree.SubElement(rPr, qn("a:hlinkClick"))
    hlinkClick.set(qn("r:id"), rId)


def _make_index_slide(prs, template_idx, link_labels, now_str, SW, SH):
    """Create Z1 index slide. Returns (slide, list_of_runs)."""
    from app.routers.hotel_overview import (
        _clone_template_slide, _set_slide_title,
    )
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    C_LINK = RGBColor(0x4B, 0xA8, 0xE8)
    slide = _clone_template_slide(prs, template_idx)
    _set_slide_title(slide, "\u5404\u9805\u6e05\u55ae\u9023\u7d50", "", now_str, SW, SH)
    tb = slide.shapes.add_textbox(
        Inches(1.5), Inches(1.6), Inches(SW - 3.0), Inches(SH - 2.5)
    )
    tf = tb.text_frame
    tf.word_wrap = False
    runs = []
    for i, label in enumerate(link_labels):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.space_before = Pt(20)
        run = para.add_run()
        run.text = "\u25b6  " + label
        run.font.size      = Pt(22)
        run.font.bold      = True
        run.font.color.rgb = C_LINK
        run.font.underline = True
        runs.append(run)
    return slide, runs


def _make_end_slide(prs, template_idx, now_str, SW, SH):
    """Create a simple ending slide."""
    from app.routers.hotel_overview import (
        _clone_template_slide, _set_slide_title, _pptx_txt,
    )
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    slide = _clone_template_slide(prs, template_idx)
    _set_slide_title(slide, "\u7c21\u5831\u7d50\u675f", "", now_str, SW, SH)
    _pptx_txt(
        slide,
        "\u4ee5\u4e0a\u8cc7\u6599\u7531\u7cfb\u7d71\u81ea\u52d5\u751f\u6210\uff0c\u4ee5\u7cfb\u7d71\u6578\u64da\u70ba\u6e96",
        2.0, SH / 2 - 0.3, SW - 4.0, 0.6,
        size=14, italic=True,
        color=RGBColor(0x88, 0x88, 0x88),
        align=PP_ALIGN.CENTER,
    )
    return slide

# ═══════════════════════════════════════════════════════


def _make_hotel_pm_stats_table(db, year: int, freq_type: str):
    """Hotel periodic maintenance annual stats table (months as columns)."""
    from app.routers.periodic_maintenance import _calc_year_matrix

    matrix = _calc_year_matrix(db, year, freq_type)
    months = matrix.months

    if not any(m.period_total > 0 or m.prev_carry_over > 0 for m in months):
        return [], []

    def _has(m):
        return m.period_total > 0 or m.prev_carry_over > 0 or m.prev_resolved_in_period > 0

    TABLE_W = 12.50
    ITEM_W  = 2.20
    mw = max(0.65, min(0.90, (TABLE_W - ITEM_W - 1.00) / 12))

    cols = [{"key": "item", "label": "\u7d71\u8a08\u9805\u76ee", "width": ITEM_W, "align": "left"}]
    for m in months:
        cols.append({"key": f"m{m.month}", "label": f"{m.month}\u6708", "width": mw, "align": "center"})
    cols.append({"key": "total", "label": "\u5408\u8a08", "align": "center"})

    active = [m for m in months if _has(m)]
    sum_res   = sum(m.prev_resolved_in_period for m in active)
    sum_total = sum(m.period_total            for m in active)
    sum_done  = sum(m.period_completed        for m in active)
    fr = f"{round(sum_done/sum_total*100,1):.1f}%" if sum_total else "\u2014"

    def _r(label, fn, tv):
        row = {"item": label}
        for m in months:
            row[f"m{m.month}"] = fn(m) if _has(m) else "\u2014"
        row["total"] = tv
        return row

    # Labels match frontend TAB display
    rows = [
        _r("\u622a\u81f3\u4e0a\u6708\u5e95\u7d2f\u8a08\u672a\u7d50\u6848\u6578",
           lambda m: str(m.prev_carry_over), "\u2014"),
        _r("\u5176\u4e2d\u672c\u6708\u5df2\u7d50\u6848\u6578",
           lambda m: str(m.prev_resolved_in_period), str(sum_res)),
        _r("\u7d2f\u8a08\u9805\u76ee\u5b8c\u6210\u7387",
           lambda m: (f"{m.carry_over_rate:.1f}%" if m.carry_over_rate is not None else "\u2014"),
           "\u2014"),
        _r("\u672c\u6708\u9031\u671f\u4fdd\u990a\u9805\u76ee\u6578",
           lambda m: str(m.period_total), str(sum_total)),
        _r("\u672c\u6708\u9031\u671f\u4fdd\u990a\u5b8c\u6210\u6578",
           lambda m: str(m.period_completed), str(sum_done)),
        _r("\u672c\u6708\u9031\u671f\u4fdd\u990a\u5b8c\u6210\u7387",
           lambda m: (f"{m.period_rate:.1f}%" if m.period_rate is not None else "\u2014"),
           fr),
    ]
    return cols, rows


def _get_ihg_section_matrix(db, year: int, month: int) -> dict:
    """Query IHG section matrix directly from DB (mirrors /section-matrix API)."""
    from app.models.ihg_room_maintenance import (
        IHGRoomMaintenanceMaster, IHGRoomMaintenanceSection,
    )
    from app.routers.ihg_room_maintenance import (
        CANONICAL_ROOMS, CANONICAL_ROOM_SET, CANONICAL_CATEGORIES, _derive_floor,
    )
    month_zf = str(month).zfill(2)
    year_str = str(year)
    master_q = db.query(
        IHGRoomMaintenanceMaster.ragic_id,
        IHGRoomMaintenanceMaster.room_no,
        IHGRoomMaintenanceMaster.floor,
        IHGRoomMaintenanceMaster.maint_date,
    ).filter(
        IHGRoomMaintenanceMaster.maint_year  == year_str,
        IHGRoomMaintenanceMaster.maint_month == month_zf,
        IHGRoomMaintenanceMaster.room_no.in_(CANONICAL_ROOM_SET),
    ).all()
    masters_by_room = {r.room_no: r for r in master_q}
    master_ids      = [r.ragic_id for r in master_q]
    secs = (db.query(IHGRoomMaintenanceSection)
            .filter(IHGRoomMaintenanceSection.master_ragic_id.in_(master_ids))
            .all()) if master_ids else []
    sec_map: dict = {}
    for s in secs:
        sec_map.setdefault(s.master_ragic_id, {})[s.category] = s.value
    rooms_out = []
    for rno in CANONICAL_ROOMS:
        if rno in masters_by_room:
            mrow = masters_by_room[rno]
            rooms_out.append({"room_no": rno, "floor": mrow.floor,
                               "maint_date": mrow.maint_date,
                               "sections": sec_map.get(mrow.ragic_id, {}),
                               "has_data": True})
        else:
            rooms_out.append({"room_no": rno, "floor": _derive_floor(rno),
                               "maint_date": "", "sections": {}, "has_data": False})
    active_cats = [c for c in CANONICAL_CATEGORIES
                   if any(r["sections"].get(c) for r in rooms_out if r["has_data"])]
    return {"rooms": rooms_out, "categories": active_cats or list(CANONICAL_CATEGORIES),
            "year": year_str, "month": month_zf}


def _build_ihg_section_matrix_slides(
    prs, template_idx: int,
    title: str, subtitle: str,
    matrix_data: dict,
    now_str: str, SW: float, SH: float,
    title_fn=None,
    max_rows: int = 14,
) -> None:
    """Build IHG section matrix slides with per-cell V/tri/X coloring (web parity)."""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from app.routers.hotel_overview import (
        _clone_template_slide, _set_slide_title,
        _pptx_cell, _pptx_header_row, _pptx_txt,
    )
    _fn = title_fn or _set_slide_title

    # Colours matching frontend SECTION_VALUE_CFG
    C_V_BG   = RGBColor(0xE8, 0xF5, 0xE9);  C_V_FG   = RGBColor(0x2E, 0x7D, 0x32)
    C_T_BG   = RGBColor(0xFF, 0xFD, 0xE7);  C_T_FG   = RGBColor(0xF5, 0x7F, 0x17)
    C_X_BG   = RGBColor(0xFC, 0xE4, 0xEC);  C_X_FG   = RGBColor(0xC6, 0x28, 0x28)
    C_NO_BG  = RGBColor(0xF5, 0xF5, 0xF5)
    C_SKIP   = RGBColor(0xFA, 0xFA, 0xFA)
    C_DARK   = RGBColor(0x1B, 0x3A, 0x5C)
    C_GRAY   = RGBColor(0x88, 0x88, 0x88)
    C_LIGHT  = RGBColor(0x4B, 0xA8, 0xE8)
    C_ALT    = RGBColor(0xEE, 0xF5, 0xFB)

    rooms      = matrix_data.get("rooms", [])
    categories = matrix_data.get("categories", [])

    if not rooms or not categories:
        sl = _clone_template_slide(prs, template_idx)
        _fn(sl, title, subtitle, now_str, SW, SH)
        _pptx_txt(sl, "\uff08\u672c\u671f\u66ab\u7121\u8cc7\u6599\uff09",
                  2.0, 3.5, 9.0, 1.0, size=14, color=C_GRAY, italic=True)
        return

    n_cats  = len(categories)
    pages   = [rooms[i:i+max_rows] for i in range(0, max(len(rooms), 1), max_rows)]
    total_p = len(pages)

    TABLE_W = SW - 0.8
    ROOM_W  = 0.65
    FLR_W   = 0.45
    cat_w   = max(0.52, (TABLE_W - ROOM_W - FLR_W) / n_cats)

    for pg_idx, pg_rooms in enumerate(pages):
        sl = _clone_template_slide(prs, template_idx)
        pg_sub = (subtitle if total_p == 1
                  else f"{subtitle}\u3000\u7b2c {pg_idx+1} \u9801\uff0f\u5171 {total_p} \u9801"
                       f"\u3000\uff08\u5171 {len(rooms)} \u9593\uff09")
        _fn(sl, title, "", now_str, SW, SH)
        if pg_sub:
            _pptx_txt(sl, pg_sub, 0.45, 0.592, SW - 5.0, 0.22, size=10, color=C_LIGHT)

        n_cols = 2 + n_cats
        n_rows = len(pg_rooms) + 1
        TABLE_Y = 0.85
        TABLE_H = SH - TABLE_Y - 0.65

        tbl = sl.shapes.add_table(
            n_rows, n_cols,
            Inches(0.4), Inches(TABLE_Y), Inches(TABLE_W), Inches(TABLE_H)
        ).table

        tbl.columns[0].width = Inches(ROOM_W)
        tbl.columns[1].width = Inches(FLR_W)
        used = ROOM_W + FLR_W
        for ci in range(2, n_cols):
            w = max(TABLE_W - used, 0.3) if ci == n_cols - 1 else cat_w
            tbl.columns[ci].width = Inches(w)
            used += w

        # Header
        _pptx_cell(tbl, 0, 0, "\u623f\u865f", bold=True)
        _pptx_cell(tbl, 0, 1, "\u6a13\u5c64", bold=True)
        for ci, cat in enumerate(categories):
            short = cat.replace("\u5ba2\u623f", "").strip() or cat
            _pptx_cell(tbl, 0, 2+ci, short, bold=True)
        _pptx_header_row(tbl, n_cols, size=9)
        tbl.rows[0].height = Pt(28)

        # Data rows
        for ri, room in enumerate(pg_rooms, 1):
            row_bg    = C_ALT if ri % 2 == 0 else None
            has_data  = room.get("has_data", False)
            sections  = room.get("sections", {})

            _pptx_cell(tbl, ri, 0, room.get("room_no", ""),
                       fg=C_DARK, bg=row_bg, size=9, bold=True, align=PP_ALIGN.CENTER)
            _pptx_cell(tbl, ri, 1, room.get("floor", ""),
                       fg=C_GRAY, bg=row_bg, size=8, align=PP_ALIGN.CENTER)

            for ci, cat in enumerate(categories):
                if not has_data:
                    _pptx_cell(tbl, ri, 2+ci, "\u672a\u57f7",
                               fg=RGBColor(0xBF, 0xBF, 0xBF), bg=C_SKIP,
                               size=8, align=PP_ALIGN.CENTER)
                else:
                    val = sections.get(cat, "")
                    if val == "V":
                        _pptx_cell(tbl, ri, 2+ci, "V",
                                   fg=C_V_FG, bg=C_V_BG, size=9, bold=True,
                                   align=PP_ALIGN.CENTER)
                    elif val == "\u25b2":          # ▲
                        _pptx_cell(tbl, ri, 2+ci, "\u25b2",
                                   fg=C_T_FG, bg=C_T_BG, size=9,
                                   align=PP_ALIGN.CENTER)
                    elif val == "X":
                        _pptx_cell(tbl, ri, 2+ci, "X",
                                   fg=C_X_FG, bg=C_X_BG, size=9, bold=True,
                                   align=PP_ALIGN.CENTER)
                    else:
                        _pptx_cell(tbl, ri, 2+ci, "\u2014",
                                   fg=C_GRAY, bg=C_NO_BG, size=8,
                                   align=PP_ALIGN.CENTER)

            tbl.rows[ri].height = Pt(22)

# Core PPT builder
# ═══════════════════════════════════════════════════════

def _build_repair_pptx(module: str, year: int, month: int, db: Session) -> BytesIO:
    from pptx import Presentation
    from app.routers.hotel_overview import (
        _update_cover_date, _clone_template_slide, _delete_slide,
    )

    # ── 載入服務 ──────────────────────────────────────────────────────────────
    if module == "dazhi":
        from app.models.dazhi_repair import DazhiRepairCase as CaseModel
        import app.services.dazhi_repair_service as svc
        module_label = "大直工務部"
    else:
        from app.models.luqun_repair import LuqunRepairCase as CaseModel
        import app.services.luqun_repair_service as svc
        module_label = "盧群商場工務報修"

    from app.services import repair_report_service as rr_svc

    all_cases   = db.query(CaseModel).all()
    now_str     = datetime.now().strftime("%Y-%m-%d %H:%M")
    period_str  = f"{year}年{month:02d}月"

    # ── IHG 模板偵測（dazhi 使用 Hotel Indigo 新版型）──────────────────────────
    is_ihg = (module == "dazhi")
    if is_ihg:
        from app.routers.hotel_ppt_export import (
            _update_ihg_cover, _set_ihg_slide_title, _move_slide_to_end,
        )
        _title_fn = _set_ihg_slide_title
    else:
        _title_fn = None

    # ── 載入 Presentation ─────────────────────────────────────────────────────
    _tpl_path = IHG_TEMPLATE_PATH if is_ihg else TEMPLATE_PATH
    prs = Presentation(_tpl_path)
    SW  = prs.slide_width.inches
    SH  = prs.slide_height.inches

    # ── Slide 0: Cover ───────────────────────────────────────────────────────
    if is_ihg:
        _update_ihg_cover(prs.slides[0], year, month)
    else:
        _update_cover_date(prs.slides[0], year, month)
        _delete_slide(prs, 1)  # IHG 模板無 TOC，舊版才需刪
    TMPL = 1

    # ══════════════════════════════════════════════════════════════════════════
    # Slide A — 3.1 報修統計（橫向：欄=月份1→12，列=統計項目）
    # ══════════════════════════════════════════════════════════════════════════
    repair_stats = svc.compute_repair_stats(all_cases, year)
    months_data  = repair_stats.get("months", {})
    # 每列為一個統計項目，每欄為一個月份（共 9 項，與 portal 完全一致）
    # kind: "count"=直接取值, "rate"=百分比格式, "sum2"=兩欄相加
    from pptx.dml.color import RGBColor as _RGB
    _C_RED_BG  = _RGB(0xFF, 0xF5, 0xF5)   # 淡紅底（③ 列）
    _C_RED_FG  = _RGB(0xC0, 0x39, 0x2B)   # 深紅字（③⑦ 列）
    _C_RATE_FG = _RGB(0xFA, 0x8C, 0x16)   # 橙色字（④⑧ 完成率列）
    _C_SUM_BG  = _RGB(0xE8, 0xF0, 0xFE)   # 淡藍底（⑨ 合計列）
    _C_DARK    = _RGB(0x1B, 0x3A, 0x5C)   # 品牌深藍（預設）

    # (label, kind, field, _row_bg, _row_fg)
    _stat31_defs = [
        ("① 上月累計未完成項目數",          "count", "prev_uncompleted",          None,      None),
        ("② 上月累計未完成，於本月結案",     "count", "closed_from_prev",           None,      None),
        ("③ 上月累計未完成，於本月仍未完成", "count", "prev_remaining",             _C_RED_BG, _C_RED_FG),
        ("④ 累計項目完成率",                 "rate",  "cum_completion_rate",        None,      _C_RATE_FG),
        ("⑤ 本月報修項目數",                 "count", "this_month_total",           None,      None),
        ("⑥ 本月報修項目完成數",             "count", "this_month_completed",       None,      None),
        ("⑦ 本月報修項目未完成",             "count", "this_month_uncompleted",     None,      _C_RED_FG),
        ("⑧ 本月報修項目完成率",             "rate",  "this_month_completion_rate", None,      _C_RATE_FG),
        ("⑨ 項目完成件數（②＋⑥）",          "sum2",  ("closed_from_prev", "this_month_completed"), _C_SUM_BG, _C_DARK),
    ]
    stats_rows = []
    for label, kind, field, row_bg, row_fg in _stat31_defs:
        row: dict = {"item": label}
        if row_bg is not None:
            row["_row_bg"] = row_bg
        if row_fg is not None:
            row["_row_fg"] = row_fg
        for m in range(1, 13):
            if m > month:          # 未來月份統一顯示 —
                row[f"m{m}"] = "—"
            else:
                md = months_data.get(m, {})
                if kind == "rate":
                    row[f"m{m}"] = _fmt_rate(md.get(field))
                elif kind == "sum2":
                    a, b = field
                    row[f"m{m}"] = _sanitize((md.get(a) or 0) + (md.get(b) or 0))
                else:
                    row[f"m{m}"] = _sanitize(md.get(field, 0))
        stats_rows.append(row)
    stats_cols = [{"key": "item", "label": "統計項目", "width": 2.90, "align": "left"}]
    for m in range(1, 13):
        stats_cols.append({"key": f"m{m}", "label": f"{m}月", "width": 0.80, "align": "center"})
    _add_table_slides(
        prs, TMPL,
        title    = "3.1 報修統計",
        subtitle = f"{period_str}",
        columns  = stats_cols,
        rows     = stats_rows,
        now_str  = now_str,
        SW=SW, SH=SH,
        sep_after_rows = [3, 7],   # ④後（累計完成率↔本月報修）、⑧後（本月完成率↔合計件數）
        title_fn       = _title_fn,
    )

    # ══════════════════════════════════════════════════════════════════════════
    # Slide B — 3.2 結案時間（橫向：欄=月份1→12，列=統計項目）
    # ══════════════════════════════════════════════════════════════════════════
    closing = svc.compute_closing_time(all_cases, year)
    monthly_closing = closing.get("monthly", {})
    # 6 列：與 portal 完全一致（小型 3 列 + 中大型 3 列）
    # kind: "count"=整數, "days"=小數1位, "avg"=小數2位（None→"—"）
    from pptx.dml.color import RGBColor as _RGB32
    _C_SMALL_BG = _RGB32(0xEE, 0xF5, 0xFB)   # 淡藍（小型組）
    _C_LARGE_BG = _RGB32(0xF0, 0xEB, 0xFF)   # 淡紫（中大型組）
    _C_DARK32   = _RGB32(0x1B, 0x3A, 0x5C)
    _C_PURPLE   = _RGB32(0x5B, 0x2D, 0x8E)   # 中大型深紫字

    _stat32_defs = [
        # (label,            size,    field,         kind,    row_bg,       row_fg)
        ("小型 結案件數",    "small", "closed_count", "count", _C_SMALL_BG,  _C_DARK32),
        ("小型 天數合計",    "small", "total_days",   "days",  _C_SMALL_BG,  _C_DARK32),
        ("小型 平均天數",    "small", "avg_days",     "avg",   _C_SMALL_BG,  _C_DARK32),
        ("中大型 結案件數",  "large", "closed_count", "count", _C_LARGE_BG,  _C_PURPLE),
        ("中大型 天數合計",  "large", "total_days",   "days",  _C_LARGE_BG,  _C_PURPLE),
        ("中大型 平均天數",  "large", "avg_days",     "avg",   _C_LARGE_BG,  _C_PURPLE),
    ]
    closing_rows = []
    for label, size_key, field, kind, row_bg, row_fg in _stat32_defs:
        row: dict = {"item": label, "_row_bg": row_bg, "_row_fg": row_fg}
        for m in range(1, 13):
            if m > month:
                row[f"m{m}"] = "—"
            else:
                mc  = monthly_closing.get(m, {})
                obj = mc.get(size_key, {})
                v   = obj.get(field)
                if kind == "count":
                    row[f"m{m}"] = _sanitize(v or 0)
                elif kind == "days":
                    row[f"m{m}"] = _sanitize(round(v, 1)) if v else "—"
                else:   # avg
                    row[f"m{m}"] = _sanitize(round(v, 2)) if v else "—"
        closing_rows.append(row)
    closing_cols = [{"key": "item", "label": "類別 / 月份", "width": 2.10, "align": "left"}]
    for m in range(1, 13):
        closing_cols.append({"key": f"m{m}", "label": f"{m}月", "width": 0.87, "align": "center"})
    _add_table_slides(
        prs, TMPL,
        title    = "3.2 結案時間統計",
        subtitle = f"{period_str}",
        columns  = closing_cols,
        rows     = closing_rows,
        now_str  = now_str,
        SW=SW, SH=SH,
        title_fn = _title_fn,
    )

    # ══════════════════════════════════════════════════════════════════════════
    # Slide C — 3.3 柏拉圖分析（圖）
    # ══════════════════════════════════════════════════════════════════════════
    type_stats = svc.compute_type_stats(all_cases, year)
    type_rows  = type_stats.get("rows", [])
    pareto_buf = _make_pareto_chart(type_rows)
    _add_chart_slide(
        prs, TMPL,
        title     = "3.3 報修類型 — 柏拉圖分析",
        subtitle  = f"{period_str}　深藍＝累計≤80%重點類別",
        chart_buf = pareto_buf,
        now_str   = now_str, SW=SW, SH=SH,
        title_fn  = _title_fn,
    )

    # ══════════════════════════════════════════════════════════════════════════
    # Slide D — 3.3 報修類型各類別（表格）
    # ══════════════════════════════════════════════════════════════════════════
    # 年度佔比：直接取 service 計算好的 cum_pct（= row_total / len(year_cases)），
    # 與前端 3.3 TAB「年度佔比」欄位口徑完全一致。
    # 不再自行用 year_total 重算（year_total 只計 REPAIR_TYPE_ORDER 內且 month 非 None 的案件，
    # 分母不同會造成數字差異）。
    sorted_type_rows = sorted(type_rows, key=lambda r: -r.get("row_total", 0))
    type_table_rows = []
    for r in sorted_type_rows:
        if r.get("row_total", 0) == 0:
            continue
        row_d = {
            "type":    r["type"],
            "total":   str(r.get("row_total", 0)),
            "pct":     f"{r.get('cum_pct', 0)}%",
        }
        for m in range(1, 13):
            v = r.get("monthly", {}).get(m, 0)
            row_d[f"m{m}"] = str(v) if v else "—"
        type_table_rows.append(row_d)

    type_cols = [
        {"key": "type",  "label": "類別", "width": 1.30, "align": "left"},
    ]
    for m in range(1, 13):
        type_cols.append({"key": f"m{m}", "label": f"{m}月", "width": 0.72, "align": "center"})
    type_cols += [
        {"key": "total", "label": "合計", "width": 0.80, "align": "center"},
        {"key": "pct",   "label": "佔比%", "width": 0.80, "align": "center"},
    ]

    _add_table_slides(
        prs, TMPL,
        title    = "3.3 報修類型各類別",
        subtitle = f"{period_str}　依件數降序",
        columns  = type_cols,
        rows     = type_table_rows,
        now_str  = now_str,
        SW=SW, SH=SH,
        max_rows = 9999,   # 不分頁，全部類別同一張投影片
        title_fn = _title_fn,
    )

    # ══════════════════════════════════════════════════════════════════════════
    # Slide E — Dashboard 報修類型分布（圓餅）
    # ══════════════════════════════════════════════════════════════════════════
    pie_buf = _make_pie_chart(type_rows)
    _add_chart_slide(
        prs, TMPL,
        title     = "Dashboard — 報修類型分布",
        subtitle  = f"{period_str}",
        chart_buf = pie_buf,
        now_str   = now_str, SW=SW, SH=SH,
        title_fn  = _title_fn,
    )

    # ══════════════════════════════════════════════════════════════════════════
    # Slide F — 3.4 本月客房報修表（房號×報修類別矩陣，將飯店模組）
    # ══════════════════════════════════════════════════════════════════════════
    if module == "dazhi":
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from app.routers.hotel_overview import (
            _clone_template_slide, _pptx_cell, _pptx_header_row,
            _set_slide_title, _pptx_txt,
        )
        from app.services.dazhi_repair_service import REPAIR_TYPE_ORDER

        _room_data = svc.compute_room_repair_table(all_cases, year, month)
        _room_rows = _room_data.get("rows", [])

        # 收集本月出現的類別
        _cats_used: set = set()
        for _rr in _room_rows:
            for _cat, _cases in _rr.get("categories", {}).items():
                if _cases:
                    _cats_used.add(_cat)
        _ordered_cats = [c for c in REPAIR_TYPE_ORDER if c in _cats_used]
        _ordered_cats += sorted(_cats_used - set(REPAIR_TYPE_ORDER))

        if _room_rows and _ordered_cats:
            _rm_pages   = [_room_rows[i:i+14] for i in range(0, max(len(_room_rows),1), 14)]
            _rm_total_p = len(_rm_pages)
            _n_cats     = len(_ordered_cats)
            _TABLE_W    = SW - 0.8
            _ROOM_W     = 0.65; _FLR_W = 0.45
            _cat_w      = max(0.65, (_TABLE_W - _ROOM_W - _FLR_W) / _n_cats)
            _C_DARK     = RGBColor(0x1B, 0x3A, 0x5C)
            _C_GRAY     = RGBColor(0x88, 0x88, 0x88)
            _C_ALT      = RGBColor(0xEE, 0xF5, 0xFB)
            _C_LIGHT    = RGBColor(0x4B, 0xA8, 0xE8)
            _C_HIT_BG   = RGBColor(0xFF, 0xF7, 0xE6)
            _C_HIT_FG   = RGBColor(0xD4, 0x60, 0x06)

            for _pg_idx, _pg_rooms in enumerate(_rm_pages):
                _slide = _clone_template_slide(prs, TMPL)
                _pg_sub = (f"{period_str}" if _rm_total_p == 1
                           else f"{period_str}　第 {_pg_idx+1} 頁/共 {_rm_total_p} 頁")
                _tfn = _title_fn or _set_slide_title
                _tfn(_slide, "3.4 本月客房報修表", "", now_str, SW, SH)
                _pptx_txt(_slide, _pg_sub, 0.45, 0.592, SW-5.0, 0.22, size=10, color=_C_LIGHT)

                _nc = 2 + _n_cats
                _nr = len(_pg_rooms) + 1
                _tbl = _slide.shapes.add_table(
                    _nr, _nc,
                    Inches(0.4), Inches(0.85), Inches(_TABLE_W), Inches(SH - 0.85 - 0.65)
                ).table
                _tbl.columns[0].width = Inches(_ROOM_W)
                _tbl.columns[1].width = Inches(_FLR_W)
                _used = _ROOM_W + _FLR_W
                for _ci in range(2, _nc):
                    _w = max(_TABLE_W-_used, 0.3) if _ci==_nc-1 else _cat_w
                    _tbl.columns[_ci].width = Inches(_w); _used += _w

                _pptx_cell(_tbl, 0, 0, "房號", bold=True)
                _pptx_cell(_tbl, 0, 1, "樓層", bold=True)
                for _ci, _cat in enumerate(_ordered_cats):
                    _pptx_cell(_tbl, 0, 2+_ci, _cat, bold=True)
                _pptx_header_row(_tbl, _nc, size=9)
                _tbl.rows[0].height = Pt(28)

                for _ri, _room in enumerate(_pg_rooms, 1):
                    _rb = _C_ALT if _ri%2==0 else None
                    _pptx_cell(_tbl, _ri, 0, _room.get("room_no",""),
                               fg=_C_DARK, bg=_rb, size=9, bold=True, align=PP_ALIGN.CENTER)
                    _pptx_cell(_tbl, _ri, 1, _room.get("floor",""),
                               fg=_C_GRAY, bg=_rb, size=8, align=PP_ALIGN.CENTER)
                    _cats_d = _room.get("categories", {})
                    for _ci, _cat in enumerate(_ordered_cats):
                        _cs = _cats_d.get(_cat, [])
                        if _cs:
                            _cd = _cs[0]
                            _desc = _sanitize(_cd.get("title","") if isinstance(_cd, dict) else str(_cd))
                            _txt  = (_desc[:18]+"...") if len(_desc)>18 else _desc
                            if len(_cs) > 1:
                                _txt = f"{len(_cs)}. {_txt}"
                            _pptx_cell(_tbl, _ri, 2+_ci, _txt,
                                       fg=_C_HIT_FG, bg=_C_HIT_BG, size=8, align=PP_ALIGN.LEFT)
                        else:
                            _pptx_cell(_tbl, _ri, 2+_ci, "",
                                       fg=_C_GRAY, bg=_rb, size=8, align=PP_ALIGN.CENTER)
                    _tbl.rows[_ri].height = Pt(22)
        else:
            _add_table_slides(
                prs, TMPL,
                title    = "3.4 本月客房報修表",
                subtitle = f"{period_str}",
                columns  = [{"key":"room_no","label":"房號","width":1.0,"align":"center"},
                            {"key":"floor","label":"樓層","width":0.8,"align":"center"},
                            {"key":"total","label":"報修件數","width":1.0,"align":"center"},
                            {"key":"cats","label":"涉及類別","width":9.35,"align":"left"}],
                rows=[], now_str=now_str, SW=SW, SH=SH, max_rows=14, title_fn=_title_fn,
            )

    # ══════════════════════════════════════════════════════════════════════════
    # Slide G — 報修金額統計（橫向：欄=月份1→12+全年合計，列=費用項目）
    # ══════════════════════════════════════════════════════════════════════════
    fee_stats  = svc.compute_fee_stats(all_cases, year)
    monthly_t  = fee_stats.get("monthly_totals", {})
    month_t    = fee_stats.get("month_totals", {})
    ft         = fee_stats.get("fee_totals", {})
    # 動態讀取各費用欄位（dazhi=3項；luqun=4項含扣款專櫃）
    fee_labels = fee_stats.get("fee_labels", {
        "outsource_fee":   "委外費用",
        "maintenance_fee": "維修費用",
        "deduction_fee":   "扣款費用",
    })

    def _fmt_fee_val(fk: str, v) -> str:
        """扣款專櫃為計數（家數），其餘為金額。"""
        if fk == "deduction_counter":
            return str(int(v)) if v else "—"
        return _fmt_money(v)

    fee_rows = []
    for fk, label in fee_labels.items():
        row: dict = {"item": label}
        for m in range(1, 13):
            row[f"m{m}"] = "—" if m > month else _fmt_fee_val(fk, monthly_t.get(m, {}).get(fk, 0))
        row["total"] = _fmt_fee_val(fk, ft.get(fk, 0))
        fee_rows.append(row)
    # 月小計列
    subtotal_row: dict = {"item": "月小計"}
    for m in range(1, 13):
        subtotal_row[f"m{m}"] = "—" if m > month else _fmt_money(month_t.get(m, 0))
    subtotal_row["total"] = _fmt_money(fee_stats.get("grand_total", 0))
    fee_rows.append(subtotal_row)
    # 欄定義：費用項目 + 1~12月 + 全年合計（最後欄自動填滿）
    fee_cols = [{"key": "item", "label": "費用項目", "width": 1.40, "align": "left"}]
    for m in range(1, 13):
        fee_cols.append({"key": f"m{m}", "label": f"{m}月", "width": 0.82, "align": "right"})
    fee_cols.append({"key": "total", "label": "全年合計", "align": "right"})
    _add_table_slides(
        prs, TMPL,
        title    = "報修金額統計",
        subtitle = f"{period_str}",
        columns  = fee_cols,
        rows     = fee_rows,
        now_str  = now_str,
        SW=SW, SH=SH,
        title_fn = _title_fn,
    )

    # ══════════════════════════════════════════════════════════════════════════
    # Slide HPM1~HPM4 — hotel PM stats + annual matrix (dazhi/IHG only)
    if is_ihg:
        _H_FREQ = {"monthly": "\u6bcf\u6708\u7dad\u8b77",
                   "quarterly": "\u6bcf\u5b63\u7dad\u8b77",
                   "yearly": "\u6bcf\u5e74\u7dad\u8b77"}
        for _hfreq in ("monthly", "quarterly", "yearly"):
            _hfl = _H_FREQ[_hfreq]
            try:
                _hcols, _hrows = _make_hotel_pm_stats_table(db, year, _hfreq)
                if _hrows:
                    _add_table_slides(
                        prs, TMPL,
                        title    = f"\u98ef\u5e97\u9031\u671f\u4fdd\u990a \u2014 {_hfl}\u5e74\u5ea6\u7d71\u8a08",
                        subtitle = f"{year}\u5e74",
                        columns  = _hcols,
                        rows     = _hrows,
                        now_str  = now_str,
                        SW=SW, SH=SH,
                        title_fn = _title_fn,
                    )
            except Exception as _he:
                logger.warning("Hotel PM stats (%s) failed: %s", _hfl, _he, exc_info=True)


    # Slide PM1~PM3 — 商場週期保養年度統計（每月／每季／每年，luqun 專用）
    # 置於「商場週期保養年度計劃表」之前
    # ══════════════════════════════════════════════════════════════════════════
    if module == "luqun":
        for freq_type in ("monthly", "quarterly", "yearly"):
            freq_label = _MALL_PM_FREQ_LABELS[freq_type]
            try:
                pm_cols, pm_rows = _make_mall_pm_stats_table(db, year, freq_type)
                if not pm_rows:
                    continue
                _add_table_slides(
                    prs, TMPL,
                    title    = f"商場週期保養 — {freq_label}年度統計",
                    subtitle = f"{year}年",
                    columns  = pm_cols,
                    rows     = pm_rows,
                    now_str  = now_str,
                    SW=SW, SH=SH,
                )
            except Exception as _pm_e:
                logger.warning("PM stats slide (%s) failed (skipped): %s",
                               freq_label, _pm_e, exc_info=True)

        for freq_type in ("monthly", "quarterly", "yearly"):
            freq_label = _MALL_PM_FREQ_LABELS[freq_type]
            if not is_ihg:
                try:
                    fb_cols, fb_rows = _make_fb_pm_stats_table(db, year, freq_type)
                    if not fb_rows:
                        continue
                    _add_table_slides(
                        prs, TMPL,
                        title    = f"全棟例行維護 — {freq_label}年度統計",
                        subtitle = f"{year}年",
                        columns  = fb_cols,
                        rows     = fb_rows,
                        now_str  = now_str,
                        SW=SW, SH=SH,
                    )
                except Exception as _fb_e:
                    logger.warning("FB stats slide (%s) failed (skipped): %s",
                                   freq_label, _fb_e, exc_info=True)


        logger.info("[CHECKPOINT] Reached before Slide I, module=%s", module)
        # ══════════════════════════════════════════════════════════════════════════
        # Slide I — Annual maintenance plan matrix
        #   dazhi -> hotel periodic maintenance (1 slide)
        #   luqun -> mall periodic maintenance + full building maintenance
        # ══════════════════════════════════════════════════════════════════════════
        # Cross-section variables for hyperlink application in Slide H
        z1_slide       = None
        z1_runs        = []
        ia_first_slide = None
        ib_first_slide = None

        try:
            all_matrix_rows = _get_annual_matrix_rows(module, year, db)

            if module == "dazhi":
                src_rows = [r for r in all_matrix_rows if r["source"] == "飯店週期保養"]
                _add_annual_matrix_slide(
                    prs, TMPL,
                    title       = "飯店週期保養年度計劃表",
                    subtitle    = f"{{year}}年　狀態：✓已完成 ○已排定 ✗逾期 △未排定 ?待排",
                    matrix_rows = src_rows,
                    now_str     = now_str, SW=SW, SH=SH,
                )

            else:  # luqun
                import traceback as _luqun_tb
                logger.info("[Slide I] luqun: all_matrix_rows=%d", len(all_matrix_rows))
                src_mall = [r for r in all_matrix_rows if r["source"] == "商場週期保養"]
                src_fb   = [r for r in all_matrix_rows if r["source"] == "全棟例行維護"]
                logger.info("[Slide I] luqun: src_mall=%d, src_fb=%d", len(src_mall), len(src_fb))

                # Z1 各項清單連結（索引頁）— 置於兩張年度計劃表之前
                # （超連結在未完成附表建立後再套用，見本區段末端）
                _z1_labels = [
                    "商場週期保養年度計劃表",
                    "全棟例行維護年度計劃表",
                    "未完成附表",
                ]
                z1_slide, z1_runs = _make_index_slide(prs, TMPL, _z1_labels, now_str, SW, SH)

                # I-a: 商場週期保養年度計劃表
                ia_first_idx = len(prs.slides)
                try:
                    _add_annual_matrix_slide(
                        prs, TMPL,
                        title       = "商場週期保養年度計劃表",
                        subtitle    = f"{{year}}年　狀態：✓已完成 ○已排定 ✗逾期 △未排定 ?待排",
                        matrix_rows = src_mall,
                        now_str     = now_str, SW=SW, SH=SH,
                    )
                    ia_first_slide = prs.slides[ia_first_idx]
                    logger.info("[Slide I] 商場週期保養年度計劃表 created OK")
                except Exception as _ia_err:
                    from app.routers.hotel_overview import _delete_slide as _ds_ia
                    while len(prs.slides) > ia_first_idx:
                        _ds_ia(prs, len(prs.slides) - 1)
                    logger.error("[Slide I] 商場週期保養年度計劃表 FAILED: %s\n%s", _ia_err, _luqun_tb.format_exc())

                # I-b: 全棟例行維護年度計劃表
                ib_first_idx = len(prs.slides)
                try:
                    _add_annual_matrix_slide(
                        prs, TMPL,
                        title       = "全棟例行維護年度計劃表",
                        subtitle    = f"{{year}}年　狀態：✓已完成 ○已排定 ✗逾期 △未排定 ?待排",
                        matrix_rows = src_fb,
                        now_str     = now_str, SW=SW, SH=SH,
                    )
                    ib_first_slide = prs.slides[ib_first_idx]
                    logger.info("[Slide I] 全棟例行維護年度計劃表 created OK")
                except Exception as _ib_err:
                    from app.routers.hotel_overview import _delete_slide as _ds_ib
                    while len(prs.slides) > ib_first_idx:
                        _ds_ib(prs, len(prs.slides) - 1)
                    logger.error("[Slide I] 全棟例行維護年度計劃表 FAILED: %s\n%s", _ib_err, _luqun_tb.format_exc())

        except Exception as _e:
            import traceback as _tb
            logger.error("Annual matrix slide FAILED: %s\n%s", _e, _tb.format_exc())

        # ══════════════════════════════════════════════════════════════════════════
        # Slide H — Unfinished cases appendix
        # ══════════════════════════════════════════════════════════════════════════
        _is_hotel = (module == "dazhi")
        unfinished = rr_svc.get_all_unfinished_cases(
            db=db, year=year, month=month,
            include_hotel=_is_hotel,
            include_mall=not _is_hotel,
        )
        uf_title = "\u672a\u5b8c\u6210\u9644\u8868\uff08\u98ef\u5e97\uff09" if _is_hotel else "\u672a\u5b8c\u6210\u9644\u8868\uff08\u5546\u5834\uff09"
        uf_rows = []
        for c in unfinished:
            uf_rows.append({
                "case_no":  _sanitize(c.get("case_no", "")),
                "occurred": (_sanitize(c.get("occurred_at", "")) or "")[:10],
                "floor":    _sanitize(c.get("floor", "")),
                "rtype":    _sanitize(c.get("repair_type", "")),
                "title":    _sanitize(c.get("title", "")),
                "status":   _sanitize(c.get("status", "")),
                "days":     _sanitize(c.get("pending_days", "")),
                "unit":     _sanitize(c.get("responsible_unit", "")),
            })
        # Z1 各項清單連結索引頁已於兩張年度計劃表之前建立；超連結於下方 UF 投影片建立後再套用
        uf_first_idx = len(prs.slides)
        _add_table_slides(
            prs, TMPL,
            title    = uf_title,
            subtitle = f"{period_str}",
            columns  = [
                {"key": "case_no",  "label": "\u6848\u4ef6\u7de8\u865f",   "width": 1.40, "align": "center"},
                {"key": "occurred", "label": "\u5831\u4fee\u65e5\u671f",   "width": 1.05, "align": "center"},
                {"key": "floor",    "label": "\u767c\u751f\u6a13\u5c64",   "width": 1.30, "align": "center"},
                {"key": "rtype",    "label": "\u5de5\u9805\u985e\u5225",   "width": 1.35, "align": "center"},
                {"key": "title",    "label": "\u5831\u4fee\u5167\u5bb9",   "width": 3.20, "align": "left"},
                {"key": "status",   "label": "\u72c0\u614b",                "width": 1.05, "align": "center"},
                {"key": "days",     "label": "\u7b49\u5f85\u5929\u6578",  "width": 1.00, "align": "center"},
                {"key": "unit",     "label": "\u5de5\u52d9\u8655\u7406\u4eba\u54e1", "width": 1.80, "align": "left"},
            ],
            rows     = uf_rows,
            now_str  = now_str,
            SW=SW, SH=SH,
            max_rows = ROWS_PER_SLIDE,
        )

        # Remove the content template placeholder slide
        # Apply hyperlinks to Z1 index slide (luqun only)
        if z1_slide is not None and z1_runs and ia_first_slide and ib_first_slide:
            try:
                uf_first_slide = (prs.slides[uf_first_idx]
                                  if len(prs.slides) > uf_first_idx else None)
                _add_slide_hyperlink(z1_slide, z1_runs[0], ia_first_slide)
                _add_slide_hyperlink(z1_slide, z1_runs[1], ib_first_slide)
                if uf_first_slide is not None:
                    _add_slide_hyperlink(z1_slide, z1_runs[2], uf_first_slide)
            except Exception as _hl_e:
                logger.warning("Index slide hyperlinks failed: %s", _hl_e, exc_info=True)



    # IHG 客房保養表明細（section matrix）
    if is_ihg:
        try:
            _sm_data = _get_ihg_section_matrix(db, year, month)
            _sm_year = _sm_data.get("year", str(year))
            _sm_month = _sm_data.get("month", str(month).zfill(2))
            _build_ihg_section_matrix_slides(
                prs, TMPL,
                title    = "\u5ba2\u623f\u4fdd\u990a\u8868\u660e\u7d30",
                subtitle = f"{_sm_year}\u5e74{_sm_month}\u6708",
                matrix_data = _sm_data,
                now_str  = now_str, SW=SW, SH=SH,
                title_fn = _title_fn,
                max_rows = 14,
            )
        except Exception as _sm_e:
            logger.warning("IHG section matrix slides failed: %s", _sm_e, exc_info=True)

    # ── 飯店 Z1 + 未完成附表（飯店），dazhi/IHG 專用 ─────────────────────────
    if is_ihg:
        # 記錄飯店週期保養年度計劃表的投影片（供 Z1 超連結使用）
        _hpm_matrix_slide = None
        try:
            # 飯店週期保養年度計劃表已在 if is_ihg: block 中建立，取最後建立的矩陣投影片
            _hmat_idx = c2.rfind("_add_annual_matrix_slide")  # 只用來找，不重複建立
        except Exception:
            pass

        # 未完成附表（飯店）
        _h_unfinished = rr_svc.get_all_unfinished_cases(
            db=db, year=year, month=month,
            include_hotel=True, include_mall=False,
        )
        _huf_rows = []
        for _hc in _h_unfinished:
            _huf_rows.append({
                "case_no":  _sanitize(_hc.get("case_no", "")),
                "occurred": (_sanitize(_hc.get("occurred_at", "")) or "")[:10],
                "floor":    _sanitize(_hc.get("floor", "")),
                "rtype":    _sanitize(_hc.get("repair_type", "")),
                "title":    _sanitize(_hc.get("title", "")),
                "status":   _sanitize(_hc.get("status", "")),
                "days":     _sanitize(_hc.get("pending_days", "")),
                "unit":     _sanitize(_hc.get("responsible_unit", "")),
            })

        # Z1 各項清單連結（倒數第2頁）
        _h_z1_labels = [
            "\u98ef\u5e97\u9031\u671f\u4fdd\u990a\u5e74\u5ea6\u8a08\u5283\u8868",
            "\u672a\u5b8c\u6210\u9644\u8868\uff08\u98ef\u5e97\uff09",
        ]
        _h_z1_slide, _h_z1_runs = _make_index_slide(
            prs, TMPL, _h_z1_labels, now_str, SW, SH
        )

        # 飯店週期保養年度計劃表（在 Z1 之後、UF 之前）
        _hmat_first_slide = None
        try:
            _hmat_rows2 = _get_annual_matrix_rows("dazhi", year, db)
            _hdazhi2 = [r for r in _hmat_rows2 if r["source"] == "飯店週期保養"]
            if _hdazhi2:
                _n_before2 = len(prs.slides)
                try:
                    _add_annual_matrix_slide(
                        prs, TMPL,
                        title       = "飯店週期保養年度計劃表",
                        subtitle    = f"{year}年　狀態：✓已完成 ○已排定 ✗逾期 △未排定 ?待排",
                        matrix_rows = _hdazhi2,
                        now_str     = now_str, SW=SW, SH=SH,
                        title_fn    = _title_fn,
                    )
                    _hmat_first_slide = prs.slides[_n_before2]
                except Exception as _inner2:
                    from app.routers.hotel_overview import _delete_slide as _ds2
                    while len(prs.slides) > _n_before2:
                        _ds2(prs, len(prs.slides) - 1)
                    logger.warning("Hotel annual matrix error: %s", _inner2, exc_info=True)
        except Exception as _hme2:
            logger.warning("Hotel annual matrix failed: %s", _hme2, exc_info=True)


        # 未完成附表（飯店）投影片
        _h_uf_first_idx = len(prs.slides)
        _add_table_slides(
            prs, TMPL,
            title    = "\u672a\u5b8c\u6210\u9644\u8868\uff08\u98ef\u5e97\uff09",
            subtitle = f"{period_str}",
            columns  = [
                {"key": "case_no",  "label": "\u6848\u4ef6\u7de8\u865f",   "width": 1.40, "align": "center"},
                {"key": "occurred", "label": "\u5831\u4fee\u65e5\u671f",   "width": 1.05, "align": "center"},
                {"key": "floor",    "label": "\u767c\u751f\u6a13\u5c64",   "width": 1.30, "align": "center"},
                {"key": "rtype",    "label": "\u5de5\u9805\u985e\u5225",   "width": 1.35, "align": "center"},
                {"key": "title",    "label": "\u5831\u4fee\u5167\u5bb9",   "width": 3.20, "align": "left"},
                {"key": "status",   "label": "\u72c0\u614b",               "width": 1.05, "align": "center"},
                {"key": "days",     "label": "\u7b49\u5f85\u5929\u6578",  "width": 1.00, "align": "center"},
                {"key": "unit",     "label": "\u5de5\u52d9\u8655\u7406\u4eba\u54e1","width": 1.80, "align": "left"},
            ],
            rows     = _huf_rows,
            now_str  = now_str,
            SW=SW, SH=SH,
            max_rows = ROWS_PER_SLIDE,
            title_fn = _title_fn,
        )

        # 套用 Z1 超連結（runs[0]=年度計劃表, runs[1]=未完成附表）
        try:
            if _hmat_first_slide is not None:
                _add_slide_hyperlink(_h_z1_slide, _h_z1_runs[0], _hmat_first_slide)
            if len(prs.slides) > _h_uf_first_idx:
                _h_uf_slide = prs.slides[_h_uf_first_idx]
                _add_slide_hyperlink(_h_z1_slide, _h_z1_runs[1], _h_uf_slide)
        except Exception as _hlh:
            logger.warning("Hotel Z1 hyperlinks failed: %s", _hlh, exc_info=True)

    # ── IHG 模板：刪除 ContentTemplate，封底移到最末頁 ────────────────────────
    if is_ihg:
        _delete_slide(prs, TMPL)
        _move_slide_to_end(prs, 1)
    else:
        _delete_slide(prs, TMPL)

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ═══════════════════════════════════════════════════════
# Router endpoint
# ═══════════════════════════════════════════════════════

@router.post("/export", summary="trigger repair PPTX export (dazhi / luqun)")
def export_repair_pptx(
    body:         RepairPptBody = Body(...),
    db:           Session       = Depends(get_db),
    current_user: User          = Depends(get_current_user),
):
    if body.module not in ("dazhi", "luqun"):
        raise HTTPException(status_code=400, detail="module must be dazhi or luqun")

    module_prefix = {
        "dazhi": "\u98ef\u5e97\u5de5\u52d9\u5831\u4fee",
        "luqun": "\u5546\u5834\u5de5\u52d9\u5831\u4fee",
    }
    try:
        pptx_buf = _build_repair_pptx(body.module, body.year, body.month, db)
    except Exception as e:
        logger.error("repair pptx build error: %s", e, exc_info=True)
        raise HTTPException(status_code=500, detail=f"PPTX build failed: {e}")

    filename = f"{module_prefix[body.module]}{body.month}\u6708\u5831\u544a.pptx"
    encoded  = quote(filename)
    return StreamingResponse(
        pptx_buf,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={
            "Content-Disposition": f"attachment; filename*=UTF-8\'\'{encoded}",
        },
    )


# ═══════════════════════════════════════════════════════
# Diagnostic endpoint
# ═══════════════════════════════════════════════════════

@router.get("/diag/chart", summary="test matplotlib chart generation (no auth required)")
def diag_chart():
    import sys, traceback as _tb
    result: dict = {}
    try:
        import matplotlib
        result["matplotlib_version"] = matplotlib.__version__
        result["matplotlib_path"]    = matplotlib.__file__
    except Exception as e:
        result["matplotlib_import_error"] = str(e)
        return result
    try:
        font = _cjk_font()
        result["cjk_font"] = font or "not found"
    except Exception as e:
        result["cjk_font_error"] = str(e)
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        result["backend"] = matplotlib.get_backend()
    except Exception as e:
        result["backend_error"] = _tb.format_exc()
        return result
    try:
        from io import BytesIO
        fig, ax = plt.subplots(figsize=(4, 2), dpi=72)
        ax.bar(["A", "B", "C"], [3, 7, 5])
        ax.set_title("test")
        buf = BytesIO()
        fig.savefig(buf, format="png", dpi=72, bbox_inches="tight")
        buf.seek(0)
        size = len(buf.read())
        plt.close(fig)
        result["simple_chart_ok"]    = True
        result["simple_chart_bytes"] = size
    except Exception:
        result["simple_chart_ok"]    = False
        result["simple_chart_error"] = _tb.format_exc()
        return result
    fake_rows = [
        {"type": "A", "row_total": 40, "cum_pct": 40.0},
        {"type": "B", "row_total": 25, "cum_pct": 25.0},
        {"type": "C", "row_total": 15, "cum_pct": 15.0},
    ]
    try:
        buf = _make_pareto_chart(fake_rows)
        result["pareto_chart_ok"]    = buf is not None
        result["pareto_chart_bytes"] = len(buf.read()) if buf else 0
    except Exception:
        result["pareto_chart_ok"]    = False
        result["pareto_chart_error"] = _tb.format_exc()
    try:
        buf = _make_pie_chart(fake_rows)
        result["pie_chart_ok"]    = buf is not None
        result["pie_chart_bytes"] = len(buf.read()) if buf else 0
    except Exception:
        result["pie_chart_ok"]    = False
        result["pie_chart_error"] = _tb.format_exc()
    result["python_version"] = sys.version
    return result
