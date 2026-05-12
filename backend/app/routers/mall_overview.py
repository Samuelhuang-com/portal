"""
商場管理 Dashboard — 跨模組彙整 API
GET /api/v1/mall/daily-hours   每日工時彙總（五項工項）

來源（均查本地 DB，不打 Ragic）：
  ① 現場報修 — luqun_repair_cases（_stat_dt 口徑：已結案→completed_at，其餘→occurred_at，排除取消）
  ② 上級交辦 — 固定 0（模組未開發）
  ③ 緊急事件 — 固定 0（模組未開發）
  ④ 例行維護 — mall_pm_batch_item + full_bldg_pm_batch_item（start_time / end_time 實際保養時間）
  ⑤ 每日巡檢 — mall_facility_inspection_batch + rf_inspection_batch（start/end_time）

回傳格式與 work-category-analysis _build_daily 一致，供前端直接套用相同表格元件。
"""
import calendar
import re
from collections import defaultdict
from datetime import date, datetime
from io import BytesIO
from typing import List, Optional
from urllib.parse import quote

from fastapi import APIRouter, Body, Depends, Query
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
from sqlalchemy.orm import Session

from app.core.database import get_db
from app.models.full_building_maintenance import FullBldgPMBatch, FullBldgPMItem
from app.models.luqun_repair import LuqunRepairCase
from app.services.time_utils import parse_minutes as _parse_minutes
from app.models.mall_facility_inspection import MallFIBatch
from app.models.mall_periodic_maintenance import (
    MallPeriodicMaintenanceBatch,
    MallPeriodicMaintenanceItem,
)
from app.models.rf_inspection import RFInspectionBatch

router = APIRouter(prefix="/mall", tags=["商場管理 Dashboard"])

# 固定五項工項（順序即表格列順序）
MALL_CATEGORIES = ["現場報修", "上級交辦", "緊急事件", "例行維護", "每日巡檢"]



@router.get("/daily-hours", summary="商場管理 — 每日工時彙總（五項工項）")
def get_mall_daily_hours(
    year:  int = Query(..., ge=2020, le=2030, description="年份"),
    month: int = Query(..., ge=1,    le=12,   description="月份（1–12）"),
    db: Session = Depends(get_db),
):
    """
    彙整五項工作類別的每日工時（HR），供商場管理 Dashboard「B. 每日累計」Tab 使用。

    回傳格式：
    ```json
    {
      "year": 2026, "month": 4,
      "days": [1, 2, ..., 30],
      "weekdays": ["二", "三", ...],
      "rows": [
        {"category": "現場報修", "hours": [1.5, 0.0, ...], "total": 10.5, "pct": 35.0},
        ...
        {"category": "TOTAL",   "hours": [...],            "total": 30.0, "pct": 100.0}
      ]
    }
    ```
    """
    _, days_in_month = calendar.monthrange(year, month)
    days = list(range(1, days_in_month + 1))
    zh = ["一", "二", "三", "四", "五", "六", "日"]
    weekdays = [zh[date(year, month, d).weekday()] for d in days]

    bucket: dict[str, dict[int, float]] = {c: defaultdict(float) for c in MALL_CATEGORIES}
    case_bucket: dict[str, dict[int, int]] = {c: defaultdict(int) for c in MALL_CATEGORIES}

    # ── ① 現場報修：_stat_dt 口徑（已結案→completed_at，其餘→occurred_at，排除取消）─
    for c in db.query(LuqunRepairCase).all():
        if c.is_excluded_flag:
            continue
        stat_dt = c.completed_at if (c.is_completed_flag and c.completed_at) else c.occurred_at
        if not stat_dt:
            continue
        if stat_dt.year == year and stat_dt.month == month:
            case_bucket["現場報修"][stat_dt.day] += 1
            if (c.work_hours or 0) > 0:
                bucket["現場報修"][stat_dt.day] += c.work_hours

    # ── ④ 例行維護：mall_pm_batch_item ───────────────────────────────────────
    # period_month 可能儲存為 "2026/04" 或 "2026/4"，用 LIKE + Python 過濾
    # 確保不遺漏任何零填充格式差異
    def _pm_day(sched: str) -> int:
        """
        scheduled_date 格式 "MM/DD" → day；
        無法解析（空白或格式異常）→ 1（落回月初，確保實際保養時間不漏算）
        """
        s = (sched or "").strip()
        if "/" in s:
            try:
                d = int(s.split("/")[1])
                if 1 <= d <= days_in_month:
                    return d
            except (ValueError, IndexError):
                pass
        return 1   # 未排定 / 格式不符 → 歸入第 1 天

    mall_pm_batches = (
        db.query(MallPeriodicMaintenanceBatch)
        .filter(MallPeriodicMaintenanceBatch.period_month.like(f"{year}/%"))
        .all()
    )
    mall_pm_batch_ids = [
        b.ragic_id for b in mall_pm_batches
        if b.period_month and int(b.period_month.split("/")[1]) == month
    ]
    if mall_pm_batch_ids:
        for item in (
            db.query(MallPeriodicMaintenanceItem)
            .filter(MallPeriodicMaintenanceItem.batch_ragic_id.in_(mall_pm_batch_ids))
            .all()
        ):
            mins = _parse_minutes(item.start_time or "", item.end_time or "")
            case_bucket["例行維護"][_pm_day(item.scheduled_date)] += 1
            if mins > 0:
                bucket["例行維護"][_pm_day(item.scheduled_date)] += mins / 60

    # ── ④ 例行維護：full_bldg_pm_batch_item ──────────────────────────────────
    fb_pm_batches = (
        db.query(FullBldgPMBatch)
        .filter(FullBldgPMBatch.period_month.like(f"{year}/%"))
        .all()
    )
    fb_batch_ids = [
        b.ragic_id for b in fb_pm_batches
        if b.period_month and int(b.period_month.split("/")[1]) == month
    ]
    if fb_batch_ids:
        for item in (
            db.query(FullBldgPMItem)
            .filter(FullBldgPMItem.batch_ragic_id.in_(fb_batch_ids))
            .all()
        ):
            mins = _parse_minutes(item.start_time or "", item.end_time or "")
            case_bucket["例行維護"][_pm_day(item.scheduled_date)] += 1
            if mins > 0:
                bucket["例行維護"][_pm_day(item.scheduled_date)] += mins / 60

    # ── ⑤ 每日巡檢：mall_facility_inspection_batch（實際+缺漏補算）─────────────
    # 5 張固定巡檢表（4F/3F/1F-3F/1F/B1F-B4F），每天每表應巡一次
    MALL_FI_SHEET_COUNT = 5
    _today = date.today()
    counting_end_day = _today.day if (year == _today.year and month == _today.month) else days_in_month
    date_prefix = f"{year}/{month:02d}/"

    fi_sheets_by_day: dict[int, set] = defaultdict(set)
    for b in (
        db.query(MallFIBatch)
        .filter(MallFIBatch.inspection_date.like(f"{date_prefix}%"))
        .all()
    ):
        try:
            day = int(b.inspection_date.split("/")[2])
            if 1 <= day <= days_in_month:
                fi_sheets_by_day[day].add(b.sheet_key)
                mins = _parse_minutes(b.start_time or "", b.end_time or "")
                bucket["每日巡檢"][day] += mins / 60
        except (ValueError, IndexError):
            pass

    # 實際場次 + 缺漏場次（≤ counting_end_day 才補缺漏）
    for d in days:
        actual = len(fi_sheets_by_day.get(d, set()))
        if d <= counting_end_day:
            case_bucket["每日巡檢"][d] += actual + max(0, MALL_FI_SHEET_COUNT - actual)
        else:
            case_bucket["每日巡檢"][d] += actual

    # ── ⑤ 每日巡檢：rf_inspection_batch（整棟巡檢，實際批次數）─────────────────
    for b in (
        db.query(RFInspectionBatch)
        .filter(RFInspectionBatch.inspection_date.like(f"{date_prefix}%"))
        .all()
    ):
        try:
            day = int(b.inspection_date.split("/")[2])
            if 1 <= day <= days_in_month:
                mins = _parse_minutes(b.start_time or "", b.end_time or "")
                case_bucket["每日巡檢"][day] += 1
                bucket["每日巡檢"][day] += mins / 60
        except (ValueError, IndexError):
            pass

    # ── 組裝結果（與 WCA _build_daily 格式完全一致）──────────────────────────
    result_rows: list[dict] = []
    grand_total = 0.0
    grand_day = [0.0] * len(days)
    grand_cases_total = 0
    grand_cases_day = [0] * len(days)

    for cat in MALL_CATEGORIES:
        day_h = [round(bucket[cat][d], 1) for d in days]
        total = round(sum(day_h), 1)
        grand_total += total
        for i, h in enumerate(day_h):
            grand_day[i] += h
        day_c = [case_bucket[cat][d] for d in days]
        cases_total = sum(day_c)
        grand_cases_total += cases_total
        for i, c in enumerate(day_c):
            grand_cases_day[i] += c
        result_rows.append({"category": cat, "hours": day_h, "total": total, "pct": 0.0,
                            "cases": day_c, "cases_total": cases_total, "cases_pct": 0.0})

    # 計算各列 %
    for row in result_rows:
        row["pct"] = round(row["total"] / grand_total * 100, 1) if grand_total else 0.0
        row["cases_pct"] = round(row["cases_total"] / grand_cases_total * 100, 1) if grand_cases_total else 0.0

    # TOTAL 合計列
    result_rows.append({
        "category":    "TOTAL",
        "hours":       [round(h, 1) for h in grand_day],
        "total":       round(grand_total, 1),
        "pct":         100.0,
        "cases":       grand_cases_day,
        "cases_total": grand_cases_total,
        "cases_pct":   100.0,
    })

    return {
        "year":     year,
        "month":    month,
        "days":     days,
        "weekdays": weekdays,
        "rows":     result_rows,
    }


@router.get("/monthly-hours", summary="商場管理 — 每月工時彙總（五項工項）")
def get_mall_monthly_hours(
    year: int = Query(..., ge=2020, le=2030, description="年份"),
    db: Session = Depends(get_db),
):
    """
    彙整五項工作類別的每月工時（HR），供商場管理 Dashboard「C. 每月累計」Tab 使用。

    回傳格式：
    ```json
    {
      "year": 2026,
      "months": [1, 2, ..., 12],
      "rows": [
        {"category": "現場報修", "hours": [1.5, 3.0, ...], "total": 45.0, "pct": 35.0},
        ...
        {"category": "TOTAL",   "hours": [...],            "total": 128.0, "pct": 100.0}
      ]
    }
    ```
    """
    bucket: dict[str, dict[int, float]] = {c: defaultdict(float) for c in MALL_CATEGORIES}
    case_bucket: dict[str, dict[int, int]] = {c: defaultdict(int) for c in MALL_CATEGORIES}
    year_prefix = f"{year}/"

    # ── ① 現場報修：_stat_dt 口徑（已結案→completed_at，其餘→occurred_at，排除取消）─
    for c in db.query(LuqunRepairCase).all():
        if c.is_excluded_flag:
            continue
        stat_dt = c.completed_at if (c.is_completed_flag and c.completed_at) else c.occurred_at
        if not stat_dt:
            continue
        if stat_dt.year == year and 1 <= stat_dt.month <= 12:
            case_bucket["現場報修"][stat_dt.month] += 1
            if (c.work_hours or 0) > 0:
                bucket["現場報修"][stat_dt.month] += c.work_hours

    # ── ④ 例行維護：mall_pm ──────────────────────────────────────────────────
    # 先撈 batch（1 次），再用 IN 一次撈全部 items（1 次），避免 N+1
    _mall_pm_batches = (
        db.query(MallPeriodicMaintenanceBatch)
        .filter(MallPeriodicMaintenanceBatch.period_month.like(f"{year_prefix}%"))
        .all()
    )
    _mall_pm_batch_month: dict[str, int] = {}
    for _b in _mall_pm_batches:
        try:
            _m = int(_b.period_month.split("/")[1])
            if 1 <= _m <= 12:
                _mall_pm_batch_month[_b.ragic_id] = _m
        except (ValueError, IndexError, AttributeError):
            pass

    if _mall_pm_batch_month:
        for item in db.query(MallPeriodicMaintenanceItem).filter(
            MallPeriodicMaintenanceItem.batch_ragic_id.in_(_mall_pm_batch_month.keys())
        ).all():
            m = _mall_pm_batch_month.get(item.batch_ragic_id)
            if m is None:
                continue
            mins = _parse_minutes(item.start_time or "", item.end_time or "")
            case_bucket["例行維護"][m] += 1
            bucket["例行維護"][m] += mins / 60

    # ── ④ 例行維護：full_bldg_pm ─────────────────────────────────────────────
    # 先撈 batch（1 次），再用 IN 一次撈全部 items（1 次），避免 N+1
    _fbldg_pm_batches = (
        db.query(FullBldgPMBatch)
        .filter(FullBldgPMBatch.period_month.like(f"{year_prefix}%"))
        .all()
    )
    _fbldg_pm_batch_month: dict[str, int] = {}
    for _b in _fbldg_pm_batches:
        try:
            _m = int(_b.period_month.split("/")[1])
            if 1 <= _m <= 12:
                _fbldg_pm_batch_month[_b.ragic_id] = _m
        except (ValueError, IndexError, AttributeError):
            pass

    if _fbldg_pm_batch_month:
        for item in db.query(FullBldgPMItem).filter(
            FullBldgPMItem.batch_ragic_id.in_(_fbldg_pm_batch_month.keys())
        ).all():
            m = _fbldg_pm_batch_month.get(item.batch_ragic_id)
            if m is None:
                continue
            mins = _parse_minutes(item.start_time or "", item.end_time or "")
            case_bucket["例行維護"][m] += 1
            bucket["例行維護"][m] += mins / 60

    # ── ⑤ 每日巡檢：mall_facility ────────────────────────────────────────────
    for b in db.query(MallFIBatch).filter(MallFIBatch.inspection_date.like(f"{year_prefix}%")).all():
        try:
            m = int(b.inspection_date.split("/")[1])
        except (ValueError, IndexError, AttributeError):
            continue
        if 1 <= m <= 12:
            case_bucket["每日巡檢"][m] += 1
            bucket["每日巡檢"][m] += _parse_minutes(b.start_time or "", b.end_time or "") / 60

    # ── ⑤ 每日巡檢：rf_inspection ────────────────────────────────────────────
    for b in db.query(RFInspectionBatch).filter(RFInspectionBatch.inspection_date.like(f"{year_prefix}%")).all():
        try:
            m = int(b.inspection_date.split("/")[1])
        except (ValueError, IndexError, AttributeError):
            continue
        if 1 <= m <= 12:
            case_bucket["每日巡檢"][m] += 1
            bucket["每日巡檢"][m] += _parse_minutes(b.start_time or "", b.end_time or "") / 60

    # ── 組裝結果 ─────────────────────────────────────────────────────────────
    result_rows: list[dict] = []
    grand_total = 0.0
    grand_m = [0.0] * 12
    grand_cases_total = 0
    grand_cases_m = [0] * 12

    for cat in MALL_CATEGORIES:
        mh = [round(bucket[cat][m], 1) for m in range(1, 13)]
        total = round(sum(mh), 1)
        grand_total += total
        for i, h in enumerate(mh):
            grand_m[i] += h
        mc = [case_bucket[cat][m] for m in range(1, 13)]
        cases_total = sum(mc)
        grand_cases_total += cases_total
        for i, c in enumerate(mc):
            grand_cases_m[i] += c
        result_rows.append({"category": cat, "hours": mh, "total": total, "pct": 0.0,
                            "cases": mc, "cases_total": cases_total, "cases_pct": 0.0})

    for row in result_rows:
        row["pct"] = round(row["total"] / grand_total * 100, 1) if grand_total else 0.0
        row["cases_pct"] = round(row["cases_total"] / grand_cases_total * 100, 1) if grand_cases_total else 0.0

    result_rows.append({
        "category":    "TOTAL",
        "hours":       [round(h, 1) for h in grand_m],
        "total":       round(grand_total, 1),
        "pct":         100.0,
        "cases":       grand_cases_m,
        "cases_total": grand_cases_total,
        "cases_pct":   100.0,
    })

    return {"year": year, "months": list(range(1, 13)), "rows": result_rows}


@router.get("/person-hours", summary="商場管理 — 人員工時佔比（五項工項）")
def get_mall_person_hours(
    year: int = Query(..., ge=2020, le=2030, description="年份"),
    db: Session = Depends(get_db),
):
    """
    彙整五項工作類別各人員工時佔比，供商場管理 Dashboard「D. 人員工時%」Tab 使用。
    格式與 WCA _build_person_table 完全一致。

    人員識別規則：
      ① 現場報修 — LuqunRepairCase.acceptor（結案人）
      ④ 例行維護 — MallPeriodicMaintenanceItem / FullBldgPMItem.executor_name（可空格分隔多人）
      ⑤ 每日巡檢 — MallFIBatch / RFInspectionBatch.inspector_name
      ②③ 上級交辦 / 緊急事件 — 模組未開發，無人員資料

    回傳格式：
    ```json
    {
      "year": 2026,
      "persons": ["王小明", "李大華", ...],
      "rows": [
        {"category": "現場報修", "pct_by_person": [45.2, 30.1, ...]},
        ...
      ]
    }
    ```
    """
    year_prefix = f"{year}/"

    # person → category → hours 彙整
    ph: dict[str, dict[str, float]] = defaultdict(lambda: defaultdict(float))

    # ── ① 現場報修：acceptor ─────────────────────────────────────────────────
    for c in db.query(LuqunRepairCase).filter(LuqunRepairCase.occ_year == year).all():
        person = (c.acceptor or "").strip()
        if person and person != "未指定" and (c.work_hours or 0) > 0:
            ph[person]["現場報修"] += c.work_hours

    # ── ④ 例行維護：executor_name（可多人空格分隔）──────────────────────────────
    # 先撈 batch（1 次），再用 IN 一次撈全部 items（1 次），避免 N+1
    _mall_pm_ids = [
        b.ragic_id for b in
        db.query(MallPeriodicMaintenanceBatch)
        .filter(MallPeriodicMaintenanceBatch.period_month.like(f"{year_prefix}%"))
        .all()
    ]
    if _mall_pm_ids:
        for item in db.query(MallPeriodicMaintenanceItem).filter(
            MallPeriodicMaintenanceItem.batch_ragic_id.in_(_mall_pm_ids)
        ).all():
            names = [n.strip() for n in (item.executor_name or "").split() if n.strip() and n.strip() != "未指定"]
            mins = _parse_minutes(item.start_time or "", item.end_time or "")
            if names and mins > 0:
                share = (mins / 60) / len(names)
                for n in names:
                    ph[n]["例行維護"] += share

    _fbldg_pm_ids = [
        b.ragic_id for b in
        db.query(FullBldgPMBatch)
        .filter(FullBldgPMBatch.period_month.like(f"{year_prefix}%"))
        .all()
    ]
    if _fbldg_pm_ids:
        for item in db.query(FullBldgPMItem).filter(
            FullBldgPMItem.batch_ragic_id.in_(_fbldg_pm_ids)
        ).all():
            names = [n.strip() for n in (item.executor_name or "").split() if n.strip() and n.strip() != "未指定"]
            mins = _parse_minutes(item.start_time or "", item.end_time or "")
            if names and mins > 0:
                share = (mins / 60) / len(names)
                for n in names:
                    ph[n]["例行維護"] += share

    # ── ⑤ 每日巡檢：inspector_name ───────────────────────────────────────────
    for b in db.query(MallFIBatch).filter(MallFIBatch.inspection_date.like(f"{year_prefix}%")).all():
        person = (b.inspector_name or "").strip()
        if person and person != "未指定":
            mins = _parse_minutes(b.start_time or "", b.end_time or "")
            if mins > 0:
                ph[person]["每日巡檢"] += mins / 60

    for b in db.query(RFInspectionBatch).filter(RFInspectionBatch.inspection_date.like(f"{year_prefix}%")).all():
        person = (b.inspector_name or "").strip()
        if person and person != "未指定":
            mins = _parse_minutes(b.start_time or "", b.end_time or "")
            if mins > 0:
                ph[person]["每日巡檢"] += mins / 60

    # ── 找出 Top-15 人員（依全類別合計工時降冪）─────────────────────────────────
    person_totals: dict[str, float] = {
        p: sum(cats.values()) for p, cats in ph.items()
    }
    persons = sorted(person_totals, key=lambda p: -person_totals[p])[:15]

    if not persons:
        return {"year": year, "persons": [], "rows": []}

    # ── 組裝結果（格式與 WCA _build_person_table 完全一致）───────────────────────
    result_rows = []
    for cat in MALL_CATEGORIES:
        cat_total = sum(ph[p][cat] for p in persons)
        result_rows.append({
            "category":      cat,
            "pct_by_person": [
                round(ph[p][cat] / cat_total * 100, 1) if cat_total else 0.0
                for p in persons
            ],
        })

    return {
        "year":          year,
        "persons":       persons,
        "person_totals": [round(person_totals[p], 1) for p in persons],
        "rows":          result_rows,
    }


# =============================================================================
# PPTX 匯出 — Pydantic 資料模型
# =============================================================================

class KpiSummaryIn(BaseModel):
    total_cases:      int
    completed_cases:  int
    total_work_hours: float
    abnormal_count:   int
    overdue_count:    int

class SourceCardIn(BaseModel):
    source_name:      str
    source_key:       str
    case_count:       int
    completed_count:  int
    completion_rate:  float
    abnormal_count:   int
    overdue_count:    int
    work_hours:       float
    actual_hours:     Optional[float] = None

class RepairCostsIn(BaseModel):
    outsource_fee:    float
    maintenance_fee:  float
    deduction_fee:    float
    month_total_fee:  float
    period_label:     str

class MallPptxPayload(BaseModel):
    kpi_summary:  KpiSummaryIn
    source_cards: List[SourceCardIn]
    repair_costs: RepairCostsIn


# =============================================================================
# PPTX 匯出 — 共用工具函式
# =============================================================================

def _mall_pptx_txt(slide, text: str, x: float, y: float, w: float, h: float,
                   size: int = 12, bold: bool = False, color=None,
                   align=None, italic: bool = False, wrap: bool = False):
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    if color is None:
        from pptx.dml.color import RGBColor
        color = RGBColor(0x1B, 0x3A, 0x5C)
    if align is None:
        align = PP_ALIGN.LEFT
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def _mall_pptx_rect(slide, x: float, y: float, w: float, h: float,
                    fill_rgb, border_rgb=None):
    from pptx.util import Inches
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.color.rgb = border_rgb if border_rgb else fill_rgb
    return shape


def _mall_pptx_header(slide, title: str, subtitle: str,
                      period_str: str, now_str: str,
                      sw: float = 13.33, sh: float = 7.5):
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    C_DARK   = RGBColor(0x1B, 0x3A, 0x5C)
    C_ACCENT = RGBColor(0x4B, 0xA8, 0xE8)
    C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
    C_GRAY   = RGBColor(0x88, 0x88, 0x88)
    _mall_pptx_rect(slide, 0, 0, sw, 0.85, C_DARK)
    _mall_pptx_txt(slide, title, 0.4, 0.07, 9.5, 0.5,
                   size=22, bold=True, color=C_WHITE)
    if subtitle:
        _mall_pptx_txt(slide, subtitle, 0.4, 0.52, 7.0, 0.28,
                       size=10, color=C_ACCENT)
    _mall_pptx_txt(slide, period_str, 9.5, 0.07, 3.5, 0.35,
                   size=11, color=C_WHITE, align=PP_ALIGN.RIGHT)
    _mall_pptx_txt(slide, f"匯出時間：{now_str}", 8.8, sh - 0.35, 4.2, 0.3,
                   size=8, color=C_GRAY, align=PP_ALIGN.RIGHT)
    _mall_pptx_txt(slide,
                   "商場管理系統  ·  自動生成，資料以系統為準",
                   0.4, sh - 0.35, 8.0, 0.3, size=8, color=C_GRAY)


def _mall_pptx_cell(tbl, row: int, col: int, text: str,
                    bold: bool = False, align=None, fg=None, bg=None,
                    size: int = 9):
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    if align is None:
        align = PP_ALIGN.LEFT
    cell = tbl.cell(row, col)
    if bg:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
    cell.text = str(text)
    cell.margin_left   = Inches(0.06)
    cell.margin_right  = Inches(0.04)
    cell.margin_top    = Inches(0.02)
    cell.margin_bottom = Inches(0.02)
    p = cell.text_frame.paragraphs[0]
    p.alignment = align
    r = p.runs[0] if p.runs else p.add_run()
    r.font.size = Pt(size)
    r.font.bold = bold
    if fg:
        r.font.color.rgb = fg


def _mall_pptx_header_row(tbl, n_cols: int, size: int = 10):
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    C_DARK  = RGBColor(0x1B, 0x3A, 0x5C)
    C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    for c in range(n_cols):
        _mall_pptx_cell(tbl, 0, c, tbl.cell(0, c).text,
                        bold=True, align=PP_ALIGN.CENTER,
                        fg=C_WHITE, bg=C_DARK, size=size)


# =============================================================================
# PPTX Slide 2 — 三層 KPI 總覽（商場管理版）
# =============================================================================

# 商場主管交辦 / 商場緊急事件：placeholder source_keys
_MALL_PH_KEYS = {"mall_supervisor", "mall_emergency", "__ph__"}

def _build_mall_slide2_kpi(slide, payload: "MallPptxPayload",
                            period_str: str,
                            SW: float = 13.33, SH: float = 7.5):
    """Render slide 2 with 3 layers: 主管摘要 / 各來源狀態 / 報修費用。"""
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    MARGIN = 0.3
    W      = SW - 2 * MARGIN   # 12.73"

    C_DARK   = RGBColor(0x1B, 0x3A, 0x5C)
    C_ACCENT = RGBColor(0x4B, 0xA8, 0xE8)
    C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
    C_GRAY   = RGBColor(0x88, 0x88, 0x88)
    C_GREEN  = RGBColor(0x52, 0xC4, 0x1A)
    C_RED    = RGBColor(0xFF, 0x4D, 0x4F)
    C_ORANGE = RGBColor(0xFA, 0x8C, 0x16)

    SRC_COLORS = {
        "mall_pm":        RGBColor(0x1B, 0x3A, 0x5C),
        "full_bldg_pm":   RGBColor(0x4B, 0xA8, 0xE8),
        "mall_facility":  RGBColor(0x72, 0x2E, 0xD1),
        "full_bldg_insp": RGBColor(0x52, 0xC4, 0x1A),
        "luqun_repair":   RGBColor(0xFA, 0x8C, 0x16),
    }

    # ── Layer 1: 主管摘要（5 KPI boxes）─────────────────────────────────────
    _mall_pptx_txt(slide, "▌ 主管摘要", MARGIN, 0.88, 4.0, 0.18,
                   size=9, bold=True, color=C_GRAY)

    s = payload.kpi_summary
    L1_Y, L1_H = 1.07, 0.88
    box_w = W / 5
    kpi_items = [
        ("本期總工項",   f"{s.total_cases:,}",        "筆",  C_DARK),
        ("已完成工項",   f"{s.completed_cases:,}",    "筆",  C_GREEN),
        ("本期工時合計", f"{s.total_work_hours:.1f}",  "HR",  C_ORANGE),
        ("異常/未完成",  f"{s.abnormal_count:,}",      "件",  C_RED),
        ("逾期未完成",   f"{s.overdue_count:,}",       "項",
         RGBColor(0xC0, 0x39, 0x2B)),
    ]
    for i, (lbl, val, unit, color) in enumerate(kpi_items):
        bx = MARGIN + i * box_w
        bg = RGBColor(0xF0, 0xF8, 0xFF) if i % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF)
        _mall_pptx_rect(slide, bx, L1_Y, box_w - 0.06, L1_H, bg, C_ACCENT)
        _mall_pptx_txt(slide, lbl,
                       bx + 0.08, L1_Y + 0.05, box_w - 0.18, 0.20,
                       size=8, color=C_GRAY)
        _mall_pptx_txt(slide, f"{val} {unit}",
                       bx + 0.06, L1_Y + 0.27, box_w - 0.14, 0.48,
                       size=19, bold=True, color=color,
                       align=PP_ALIGN.CENTER)

    # ── Layer 2: 各來源本期狀態（2 × 4 卡）──────────────────────────────────
    L2_LABEL_Y = 2.04
    _mall_pptx_txt(slide, "▌ 各來源本期狀態", MARGIN, L2_LABEL_Y, 5.0, 0.18,
                   size=9, bold=True, color=C_GRAY)

    cards = list(payload.source_cards)
    # 補齊到 8 張卡（2 列 × 4 欄）
    ph_names = ["商場主管交辦", "商場緊急事件"]
    ph_idx   = 0
    while len(cards) < 8:
        name = ph_names[ph_idx] if ph_idx < len(ph_names) else "—"
        ph_idx += 1
        cards.append(SourceCardIn(
            source_name=name, source_key="__ph__",
            case_count=-1, completed_count=0, completion_rate=-1,
            abnormal_count=0, overdue_count=0, work_hours=-1,
        ))

    L2_Y   = 2.24
    card_h = 1.52
    card_w = W / 4
    row_gap = 0.08

    for idx, card in enumerate(cards[:8]):
        row = idx // 4
        col = idx % 4
        cx  = MARGIN + col * card_w
        cy  = L2_Y + row * (card_h + row_gap)
        is_ph = card.source_key in _MALL_PH_KEYS
        src_color = SRC_COLORS.get(card.source_key, C_GRAY)

        # Card background + top color bar
        _mall_pptx_rect(slide, cx + 0.04, cy, card_w - 0.08, card_h,
                        RGBColor(0xFF, 0xFF, 0xFF), src_color)
        _mall_pptx_rect(slide, cx + 0.04, cy, card_w - 0.08, 0.10, src_color)

        if is_ph:
            _mall_pptx_txt(slide, card.source_name,
                           cx + 0.12, cy + 0.12, card_w - 0.24, 0.24,
                           size=10, bold=True, color=C_GRAY)
            _mall_pptx_txt(slide, "數據準備中",
                           cx + 0.12, cy + 0.65, card_w - 0.24, 0.25,
                           size=9, color=C_GRAY, align=PP_ALIGN.CENTER)
            continue

        # Source name
        _mall_pptx_txt(slide, card.source_name,
                       cx + 0.12, cy + 0.12, card_w - 0.24, 0.24,
                       size=10, bold=True, color=src_color)

        # Case / completed counts
        half     = (card_w - 0.24) / 2
        case_cnt = max(0, card.case_count)
        _mall_pptx_txt(slide, f"工項  {case_cnt:,} 筆",
                       cx + 0.12, cy + 0.38, half, 0.20, size=9, color=C_DARK)
        _mall_pptx_txt(slide, f"完成  {card.completed_count:,} 筆",
                       cx + 0.12 + half, cy + 0.38, half, 0.20,
                       size=9, color=C_GREEN)

        # Completion rate bar
        if card.completion_rate >= 0:
            bar_x = cx + 0.12
            bar_y = cy + 0.61
            bw    = card_w - 0.24
            _mall_pptx_rect(slide, bar_x, bar_y, bw, 0.09,
                            RGBColor(0xDD, 0xDD, 0xDD))
            rate_clamped = min(card.completion_rate / 100, 1.0)
            fill_w = bw * rate_clamped
            if fill_w > 0.01:
                fc = (C_GREEN  if card.completion_rate >= 80
                      else C_ORANGE if card.completion_rate >= 50
                      else C_RED)
                _mall_pptx_rect(slide, bar_x, bar_y, fill_w, 0.09, fc)
            _mall_pptx_txt(slide, f"完成率 {card.completion_rate:.0f}%",
                           bar_x, bar_y + 0.10, bw, 0.17, size=8, color=C_GRAY)

        # Info row: abnormal / overdue / work hours / actual hours
        parts = []
        if card.abnormal_count > 0:
            lbl = "未完成" if card.source_key == "luqun_repair" else "異常"
            parts.append(f"{lbl}：{card.abnormal_count:,}")
        if card.overdue_count > 0:
            parts.append(f"逾期：{card.overdue_count:,}")
        hrs = max(0.0, card.work_hours)
        if hrs > 0:
            parts.append(f"預估：{hrs:.1f} HR")
        if card.actual_hours is not None and card.actual_hours > 0:
            parts.append(f"保養：{card.actual_hours:.1f} HR")
        if parts:
            _mall_pptx_txt(slide, "  ".join(parts),
                           cx + 0.12, cy + 1.27, card_w - 0.24, 0.20,
                           size=8, color=C_DARK)

    # ── Layer 3: 報修費用摘要（3 費用 box）───────────────────────────────────
    L3_Y = L2_Y + 2 * (card_h + row_gap) + 0.10
    _mall_pptx_txt(slide, "▌ 報修費用摘要", MARGIN, L3_Y - 0.20, 5.0, 0.18,
                   size=9, bold=True, color=C_GRAY)

    c     = payload.repair_costs
    fee_w = W / 3
    FEE_H = 0.80
    fee_items = [
        (f"委外+維修費用（{c.period_label}）",
         c.outsource_fee + c.maintenance_fee,
         f"委外 ${c.outsource_fee:,.0f}  /  維修 ${c.maintenance_fee:,.0f}"),
        (f"扣款費用（{c.period_label}）",
         c.deduction_fee, None),
        ("本月費用合計",
         c.month_total_fee, None),
    ]
    for i, (lbl, val, sub) in enumerate(fee_items):
        fx        = MARGIN + i * fee_w
        val_color = C_RED if val > 0 else C_DARK
        bg        = (RGBColor(0xFF, 0xF5, 0xF5) if val > 0
                     else RGBColor(0xF6, 0xFF, 0xED))
        _mall_pptx_rect(slide, fx + 0.05, L3_Y, fee_w - 0.10, FEE_H,
                        bg, C_ACCENT)
        _mall_pptx_txt(slide, lbl,
                       fx + 0.14, L3_Y + 0.04, fee_w - 0.28, 0.20,
                       size=8, color=C_GRAY)
        _mall_pptx_txt(slide, f"${val:,.0f}",
                       fx + 0.10, L3_Y + 0.24, fee_w - 0.20, 0.38,
                       size=20, bold=True, color=val_color,
                       align=PP_ALIGN.CENTER)
        if sub:
            _mall_pptx_txt(slide, sub,
                           fx + 0.10, L3_Y + 0.63, fee_w - 0.20, 0.16,
                           size=8, color=C_GRAY, align=PP_ALIGN.CENTER)


# =============================================================================
# PPTX 主建構函式
# =============================================================================

def _build_mall_pptx(year: int, month: int,
                     daily: dict, monthly: dict, persons: dict,
                     kpi_payload: Optional["MallPptxPayload"] = None) -> BytesIO:
    """
    產生 5 張投影片的商場管理 Dashboard 報告（.pptx）。
    Slide 1: 封面
    Slide 2: 本期績效總覽（KPI payload，由前端傳入）
    Slide 3: 每日工時累計（daily-hours API 資料）
    Slide 4: 每月工時累計（monthly-hours API 資料）
    Slide 5: 人員工時排名（person-hours API 資料）
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    SW, SH = 13.33, 7.5
    C_DARK    = RGBColor(0x1B, 0x3A, 0x5C)
    C_DARK_BG = RGBColor(0x0F, 0x26, 0x40)
    C_ACCENT  = RGBColor(0x4B, 0xA8, 0xE8)
    C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
    C_ROW_ALT = RGBColor(0xEE, 0xF5, 0xFB)
    C_TOTAL   = RGBColor(0xE8, 0xF4, 0xFF)
    C_GRAY    = RGBColor(0x88, 0x88, 0x88)

    prs = Presentation()
    prs.slide_width  = Inches(SW)
    prs.slide_height = Inches(SH)
    blank = prs.slide_layouts[6]

    now_str    = datetime.now().strftime("%Y-%m-%d %H:%M")
    period_str = f"{year}年{month:02d}月"

    # ── Slide 1: Cover ────────────────────────────────────────────────────────
    s1 = prs.slides.add_slide(blank)
    s1.background.fill.solid()
    s1.background.fill.fore_color.rgb = C_DARK_BG
    _mall_pptx_rect(s1, 0, 0, SW, 0.12, C_ACCENT)
    _mall_pptx_txt(s1, "商場管理 Dashboard", SW/2 - 5, 2.2, 10, 1.0,
                   size=38, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    _mall_pptx_txt(s1, "匯出報告", SW/2 - 5, 3.1, 10, 0.7,
                   size=38, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    _mall_pptx_rect(s1, SW/2 - 2.0, 4.05, 4.0, 0.55, C_ACCENT)
    _mall_pptx_txt(s1, f"統計月份：{period_str}", SW/2 - 2.0, 4.1, 4.0, 0.5,
                   size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    _mall_pptx_txt(s1,
                   f"匯出時間：{now_str}　　商場管理系統  ·  自動生成",
                   0.6, SH - 0.55, SW - 1.2, 0.4,
                   size=10, color=RGBColor(0xAA, 0xBB, 0xCC),
                   align=PP_ALIGN.CENTER)

    # ── Slide 2: 本期績效總覽 ─────────────────────────────────────────────────
    s2 = prs.slides.add_slide(blank)
    _mall_pptx_header(s2, "本期績效總覽",
                      f"{period_str}  主管摘要 · 各來源狀態 · 報修費用",
                      period_str, now_str)
    if kpi_payload is not None:
        _build_mall_slide2_kpi(s2, kpi_payload, period_str, SW=SW, SH=SH)

    # ── Slide 3: 每日工時累計 ─────────────────────────────────────────────────
    s3 = prs.slides.add_slide(blank)
    _mall_pptx_header(s3, "每日工時累計",
                      f"{period_str}  各工項每日工時分布摘要",
                      period_str, now_str)
    n_rows3 = len(daily["rows"]) + 1
    t3 = s3.shapes.add_table(
        n_rows3, 5,
        Inches(0.4), Inches(0.95), Inches(12.53), Inches(5.8)
    ).table
    t3.columns[0].width = Inches(2.8)
    t3.columns[1].width = Inches(2.5)
    t3.columns[2].width = Inches(1.8)
    t3.columns[3].width = Inches(2.3)
    t3.columns[4].width = Inches(3.13)
    for c_i, h in enumerate(["工項", "月合計(HR)", "占比 %", "最高工時日", "最高日工時(HR)"]):
        _mall_pptx_cell(t3, 0, c_i, h, bold=True)
    _mall_pptx_header_row(t3, 5, size=10)
    for ri, row_d in enumerate(daily["rows"], 1):
        cat      = row_d["category"]
        total    = round(row_d["total"], 1)
        pct      = round(row_d["pct"], 1)
        hrs_list = row_d["hours"]
        if hrs_list and max(hrs_list) > 0:
            peak_idx = max(range(len(hrs_list)), key=lambda i: hrs_list[i])
            peak_day = daily["days"][peak_idx] if peak_idx < len(daily["days"]) else peak_idx + 1
            peak_val = round(hrs_list[peak_idx], 1)
            peak_str = f"第 {peak_day} 日"
        else:
            peak_str, peak_val = "—", 0.0
        is_total = cat == "TOTAL"
        bg  = C_TOTAL if is_total else (C_ROW_ALT if ri % 2 == 0 else None)
        num = PP_ALIGN.RIGHT
        _mall_pptx_cell(t3, ri, 0, cat,                                     bold=is_total, fg=C_DARK, bg=bg, size=10)
        _mall_pptx_cell(t3, ri, 1, f"{total:,.1f}",                        bold=is_total, align=num,  fg=C_DARK, bg=bg, size=10)
        _mall_pptx_cell(t3, ri, 2, f"{pct:.1f}",                           bold=is_total, align=num,  fg=C_DARK, bg=bg, size=10)
        _mall_pptx_cell(t3, ri, 3, peak_str,                               bold=is_total, align=PP_ALIGN.CENTER, fg=C_DARK, bg=bg, size=10)
        _mall_pptx_cell(t3, ri, 4, f"{peak_val:.1f}" if peak_val else "—", bold=is_total, align=num,  fg=C_DARK, bg=bg, size=10)
    t3.rows[0].height = Pt(30)
    for ri in range(1, n_rows3):
        t3.rows[ri].height = Pt(28)
    _mall_pptx_txt(s3,
                   f"※ 統計範圍：{year}年{month:02d}月，共 {len(daily['days'])} 天",
                   0.4, SH - 0.5, 10, 0.3, size=8, color=C_GRAY, italic=True)

    # ── Slide 4: 每月工時累計 ─────────────────────────────────────────────────
    s4 = prs.slides.add_slide(blank)
    _mall_pptx_header(s4, "每月工時累計",
                      f"{year} 全年各工項月度工時彙總",
                      f"{year}年全年", now_str)
    n_rows4 = len(monthly["rows"]) + 1
    t4 = s4.shapes.add_table(
        n_rows4, 14,
        Inches(0.4), Inches(0.95), Inches(12.53), Inches(5.8)
    ).table
    t4.columns[0].width = Inches(2.0)
    for i in range(1, 13):
        t4.columns[i].width = Inches(0.80)
    t4.columns[13].width = Inches(0.93)
    mth_labels = ["1月","2月","3月","4月","5月","6月",
                  "7月","8月","9月","10月","11月","12月"]
    for c_i, h in enumerate(["工項"] + mth_labels + ["全年合計"]):
        _mall_pptx_cell(t4, 0, c_i, h, bold=True)
    _mall_pptx_header_row(t4, 14, size=9)
    for ri, row_d in enumerate(monthly["rows"], 1):
        cat     = row_d["category"]
        mth_hrs = row_d["hours"]
        total   = round(row_d["total"], 1)
        is_total = cat == "TOTAL"
        bg       = C_TOTAL if is_total else (C_ROW_ALT if ri % 2 == 0 else None)
        _mall_pptx_cell(t4, ri, 0, cat, bold=is_total, fg=C_DARK, bg=bg, size=9)
        for ci, h in enumerate(mth_hrs, 1):
            val = round(h, 1)
            _mall_pptx_cell(t4, ri, ci,
                            f"{val:.1f}" if val > 0 else "—",
                            bold=is_total, align=PP_ALIGN.RIGHT,
                            fg=C_DARK, bg=bg, size=8)
        _mall_pptx_cell(t4, ri, 13, f"{total:,.1f}",
                        bold=is_total, align=PP_ALIGN.RIGHT,
                        fg=C_DARK, bg=bg, size=9)
    for ri in range(n_rows4):
        t4.rows[ri].height = Pt(26)
    _mall_pptx_txt(s4,
                   f"※ 統計年份：{year}，工時單位：小時（HR）",
                   0.4, SH - 0.5, 10, 0.3, size=8, color=C_GRAY, italic=True)

    # ── Slide 5: 人員工時排名 ─────────────────────────────────────────────────
    s5 = prs.slides.add_slide(blank)
    person_names  = persons.get("persons", [])
    person_totals = persons.get("person_totals", [])
    _mall_pptx_header(s5, "人員工時排名",
                      f"{year} 全年  Top {len(person_names)} 人",
                      f"{year}年全年", now_str)
    grand_person = sum(person_totals)
    if person_names:
        n_rows5 = len(person_names) + 2
        t5 = s5.shapes.add_table(
            n_rows5, 4,
            Inches(0.4), Inches(0.95), Inches(8.2), Inches(5.8)
        ).table
        t5.columns[0].width = Inches(0.8)
        t5.columns[1].width = Inches(3.1)
        t5.columns[2].width = Inches(2.3)
        t5.columns[3].width = Inches(2.0)
        for c_i, h in enumerate(["排名", "姓名", f"{year}年工時(HR)", "全年占比 %"]):
            _mall_pptx_cell(t5, 0, c_i, h, bold=True)
        _mall_pptx_header_row(t5, 4, size=10)
        for ri, (name, hrs) in enumerate(zip(person_names, person_totals), 1):
            pct = round(hrs / grand_person * 100, 1) if grand_person else 0.0
            bg  = C_ROW_ALT if ri % 2 == 0 else None
            _mall_pptx_cell(t5, ri, 0, str(ri),        bold=False, align=PP_ALIGN.CENTER, fg=C_DARK, bg=bg, size=10)
            _mall_pptx_cell(t5, ri, 1, name,           bold=False, fg=C_DARK, bg=bg, size=10)
            _mall_pptx_cell(t5, ri, 2, f"{hrs:.1f}",   bold=False, align=PP_ALIGN.RIGHT, fg=C_DARK, bg=bg, size=10)
            _mall_pptx_cell(t5, ri, 3, f"{pct:.1f} %", bold=False, align=PP_ALIGN.RIGHT, fg=C_DARK, bg=bg, size=10)
        tr = n_rows5 - 1
        _mall_pptx_cell(t5, tr, 0, "",                  bg=C_TOTAL)
        _mall_pptx_cell(t5, tr, 1, "合計",               bold=True, fg=C_DARK, bg=C_TOTAL, size=10)
        _mall_pptx_cell(t5, tr, 2, f"{grand_person:.1f}", bold=True, align=PP_ALIGN.RIGHT, fg=C_DARK, bg=C_TOTAL, size=10)
        _mall_pptx_cell(t5, tr, 3, "100.0 %",            bold=True, align=PP_ALIGN.RIGHT, fg=C_DARK, bg=C_TOTAL, size=10)
        for ri in range(n_rows5):
            t5.rows[ri].height = Pt(24)
        _mall_pptx_txt(s5,
                       "工時來源：\n現場報修 / 例行維護\n每日巡檢",
                       8.9, 1.0, 4.0, 2.0,
                       size=9, color=C_GRAY, italic=True, wrap=True)
    else:
        _mall_pptx_txt(s5, f"（{year} 年暫無人員工時記錄）",
                       2.0, 3.0, 9.0, 1.0,
                       size=18, color=C_GRAY, align=PP_ALIGN.CENTER)
    _mall_pptx_txt(s5,
                   f"※ 統計來源：五項工作類別全年工時，人員依總工時降冪排列（Top 15）",
                   0.4, SH - 0.5, 12, 0.3, size=8, color=C_GRAY, italic=True)

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# =============================================================================
# Export endpoint — POST /mall/overview/export/pptx
# =============================================================================

@router.post("/overview/export/pptx", summary="商場管理 Dashboard — 匯出 PowerPoint 報告")
def export_mall_overview_pptx(
    year:  int              = Query(..., ge=2020, le=2030, description="年份"),
    month: int              = Query(..., ge=1,    le=12,   description="月份（1–12）"),
    body:  MallPptxPayload  = Body(...),
    db:    Session          = Depends(get_db),
):
    """
    匯出商場管理 Dashboard 完整 PowerPoint 報告（.pptx）。
    body 由前端帶入已計算好的 KPI 資料（主管摘要 / 各來源狀態 / 費用摘要），
    Slide 3–5 工時資料由後端自行查 DB。
    """
    daily   = get_mall_daily_hours(year, month, db)
    monthly = get_mall_monthly_hours(year, db)
    persons = get_mall_person_hours(year, db)
    buf     = _build_mall_pptx(year, month, daily, monthly, persons,
                                kpi_payload=body)
    filename = f"商場管理報告_{year}年{month:02d}月.pptx"
    content_disposition = f"attachment; filename*=UTF-8''{quote(filename)}"
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": content_disposition},
    )
