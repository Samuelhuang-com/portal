"""
飯店 Dashboard PPT 匯出模組 — 獨立 Router
==========================================

路由前綴：/hotel/ppt-export
  GET  /sections  → 所有已註冊 section 的 metadata
  GET  /config    → 使用者設定（merge registry metadata）
  POST /config    → 儲存使用者設定
  POST /export    → 觸發匯出（Registry-driven _build_hotel_pptx_v2）

Section Registry 於本模組 import 時自動執行（共 17 個 sections）：
  - hotel_overview 原有 12 個（Dashboard 7 + 每日/每月/每年/人員工時%/人員排名）
  - 報修管理 2 個（報修未完成報表 / 本月結案工單）
  - 大直工務部 3 個（年度報修統計 / 未完成工單 / 本月結案工單）
"""

import json
import logging
import traceback as _tb
from datetime import datetime
from io import BytesIO
from typing import Any, Optional, List
from urllib.parse import quote

from fastapi import APIRouter, Body, Depends, HTTPException, Query
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
from sqlalchemy.orm import Session

from app.core.database import get_db
from app.dependencies import get_current_user
from app.models.ppt_export_config import PptExportConfig
from app.models.ppt_export_history import PptExportHistory
from app.models.user import User
from app.services.ppt_section_registry import (
    PptSectionDef,
    build_default_config,
    fetch_section_data,
    get_section_def,
    get_sections,
    register_section,
)

logger = logging.getLogger("hotel_ppt_export")

router = APIRouter(prefix="/hotel/ppt-export", tags=["飯店 PPT 匯出"])

MODULE_KEY = "hotel_overview"


# ═════════════════════════════════════════════════════════════════════════════
# Pydantic schemas
# ═════════════════════════════════════════════════════════════════════════════

class ExportConfigItem(BaseModel):
    export_key:            str
    enabled:               bool
    include_detail:        bool = False
    sort_order:            int
    second_title_override: Optional[str] = None   # B-1：使用者自訂投影片標題


class SaveConfigBody(BaseModel):
    config:      List[ExportConfigItem]
    template_id: str = "default"


class FrontendData(BaseModel):
    """前端計算好的 frontend_payload sections 資料。"""
    kpi_summary:     Optional[dict]        = None
    source_cards:    Optional[List[dict]]  = None
    # repair_costs 接受兩種格式：
    #   ① dict   → {"outsource_fee": x, "maintenance_fee": y, "deduction_fee": z}
    #   ② list   → [{"category": "外包費", "amount": x}, ...]（前端 Dashboard 格式）
    repair_costs:    Optional[Any]         = None
    bar_chart_data:  Optional[List[dict]]  = None
    rate_chart_data: Optional[List[dict]]  = None
    dazhi_trend_data: Optional[List[dict]] = None
    hours_pie_data:  Optional[List[dict]]  = None


class ExportBody(BaseModel):
    year:            int
    month:           int
    inspection_date: str = ""
    frontend_data:   FrontendData = FrontendData()


# ═════════════════════════════════════════════════════════════════════════════
# Helper — merge registry + DB config
# ═════════════════════════════════════════════════════════════════════════════

# Dashboard 群組定義（合為一張投影片的 section 群組）
DASHBOARD_KPI_GROUP   = ["dashboard_kpi_summary", "dashboard_source_status", "dashboard_repair_costs"]
DASHBOARD_CHART_GROUP = ["dashboard_bar_chart", "dashboard_rate_chart", "dashboard_dazhi_trend", "dashboard_hours_pie"]
DASHBOARD_GROUPS      = [DASHBOARD_KPI_GROUP, DASHBOARD_CHART_GROUP]

# export_key → slide_group_id 映射（前端用於顯示分組標示）
SLIDE_GROUP_MAP: dict[str, str] = {}
for _g in DASHBOARD_GROUPS:
    for _k in _g:
        SLIDE_GROUP_MAP[_k] = _g[0]  # 以群組第一個 key 作為群組代表


def _sync_dashboard_groups(config_list: list[dict]) -> list[dict]:
    """
    Dashboard 投影片群組同步：
    KPI 群組（sections 1-3）和圖表群組（sections 4-7）各自共用一張投影片，
    因此同群組內任一 section 的 enabled 狀態必須一致。
    規則：以群組中第一個出現（按 sort_order）的 section enabled 值為準，同步至全群組。
    """
    key_map = {c["export_key"]: c for c in config_list}
    for group in DASHBOARD_GROUPS:
        # 找出群組中排序最前面（sort_order 最小）的 section 作為代表
        members = [key_map[k] for k in group if k in key_map]
        if not members:
            continue
        representative = min(members, key=lambda c: c.get("sort_order", 99))
        enabled = representative["enabled"]
        for m in members:
            m["enabled"] = enabled
    return config_list


def _load_merged_config(db: Session, user_id: Optional[str] = None) -> tuple[list[dict], str]:
    """
    從 DB 讀取使用者設定，merge Registry metadata，回傳完整 config list。
    若 DB 無記錄，回傳 Registry 預設值。

    Parameters
    ----------
    user_id : 指定用戶 ID，None 表示全局設定（向後相容）

    Returns (config_list, template_id)
    """
    query = db.query(PptExportConfig).filter(PptExportConfig.module_key == MODULE_KEY)
    if user_id is not None:
        row = query.filter(PptExportConfig.user_id == user_id).first()
        # fallback：若無個人設定，取全局設定
        if row is None:
            row = query.filter(PptExportConfig.user_id.is_(None)).first()
    else:
        row = query.filter(PptExportConfig.user_id.is_(None)).first()

    # Registry 中所有 section 的 metadata（按 sort_order 排序）
    registry_sections = get_sections(MODULE_KEY)
    registry_map = {s.export_key: s for s in registry_sections}

    if row is None:
        # 無 DB 記錄 → 回傳 Registry 預設（全部 enabled，sort_order 照 registry）
        config_list = [
            _merge_item(s, {"enabled": True, "include_detail": False, "sort_order": s.sort_order})
            for s in registry_sections
        ]
        return _sync_dashboard_groups(config_list), "default"

    # 有 DB 記錄 → merge
    saved: list[dict] = json.loads(row.config_json)

    config_list = []
    # ① DB 中已有的 section（按 DB sort_order 排序）
    for item in sorted(saved, key=lambda x: x.get("sort_order", 99)):
        s = registry_map.get(item["export_key"])
        if s is None:
            continue  # registry 已移除的 section，跳過
        config_list.append(_merge_item(s, item))

    # ② Registry 中新增（DB 尚未記錄）的 section，附在末尾，預設 enabled=True
    existing_keys = {item["export_key"] for item in saved}
    for s in registry_sections:
        if s.export_key not in existing_keys:
            config_list.append(_merge_item(s, {
                "enabled": True,
                "include_detail": False,
                "sort_order": s.sort_order,
            }))

    return _sync_dashboard_groups(config_list), row.template_id or "default"


def _merge_item(section: PptSectionDef, user_pref: dict) -> dict:
    """將 Registry metadata 與使用者偏好合併為一個 config dict。"""
    # second_title 優先使用使用者自訂覆寫值（B-1）
    override_title = user_pref.get("second_title_override") or ""
    return {
        # Registry metadata（唯讀）
        "export_key":             section.export_key,
        "tab_name":               section.tab_name,
        "second_title":           override_title or section.second_title,
        "second_title_default":   section.second_title,   # 始終回傳 registry 預設值供前端 placeholder 用
        "description":            section.description,
        "export_type":            section.export_type,
        "slide_layout":           section.slide_layout,
        "supports_detail":        section.supports_detail,
        "detail_description":     section.detail_description,
        "data_source":            section.data_source,
        # 滑動群組資訊（前端分組顯示用）
        "slide_group_id":         SLIDE_GROUP_MAP.get(section.export_key),
        # 使用者偏好（可修改）
        "enabled":                user_pref.get("enabled", True),
        "include_detail":         user_pref.get("include_detail", False),
        "sort_order":             user_pref.get("sort_order", section.sort_order),
        "second_title_override":  override_title or None,
    }


# ═════════════════════════════════════════════════════════════════════════════
# GET /sections
# ═════════════════════════════════════════════════════════════════════════════

@router.get("/sections", summary="取得所有可匯出 Section 的 metadata（Registry）")
def get_ppt_sections(current_user: User = Depends(get_current_user)):
    """回傳 Registry 中所有已註冊 sections 的 metadata，供設定頁 UI 使用。"""
    sections = get_sections(MODULE_KEY)
    return {
        "module_key": MODULE_KEY,
        "sections": [
            {
                "export_key":        s.export_key,
                "tab_name":          s.tab_name,
                "second_title":      s.second_title,
                "description":       s.description,
                "export_type":       s.export_type,
                "slide_layout":      s.slide_layout,
                "supports_detail":   s.supports_detail,
                "detail_description": s.detail_description,
                "data_source":       s.data_source,
                "sort_order":        s.sort_order,
            }
            for s in sections
        ],
    }


# ═════════════════════════════════════════════════════════════════════════════
# GET /config
# ═════════════════════════════════════════════════════════════════════════════

@router.get("/config", summary="取得使用者 PPT 匯出設定（merge Registry metadata）")
def get_ppt_config(
    db:           Session = Depends(get_db),
    current_user: User    = Depends(get_current_user),
):
    config_list, template_id = _load_merged_config(db, user_id=str(current_user.id))
    # 取個人設定 row（用於顯示 updated_by / updated_at）
    row = (
        db.query(PptExportConfig)
        .filter(
            PptExportConfig.module_key == MODULE_KEY,
            PptExportConfig.user_id == str(current_user.id),
        )
        .first()
    )
    return {
        "module_key":  MODULE_KEY,
        "template_id": template_id,
        "updated_by":  row.updated_by  if row else None,
        "updated_at":  row.updated_at.isoformat() if (row and row.updated_at) else None,
        "config":      config_list,
    }


# ═════════════════════════════════════════════════════════════════════════════
# POST /config
# ═════════════════════════════════════════════════════════════════════════════

@router.post("/config", summary="儲存使用者 PPT 匯出設定")
def save_ppt_config(
    body:         SaveConfigBody = Body(...),
    db:           Session        = Depends(get_db),
    current_user: User           = Depends(get_current_user),
):
    # 只儲存使用者偏好部分（不含 registry metadata）
    to_save = [
        {
            "export_key":            item.export_key,
            "enabled":               item.enabled,
            "include_detail":        item.include_detail,
            "sort_order":            item.sort_order,
            "second_title_override": item.second_title_override or None,
        }
        for item in body.config
    ]
    uid = str(current_user.id)
    row = (
        db.query(PptExportConfig)
        .filter(
            PptExportConfig.module_key == MODULE_KEY,
            PptExportConfig.user_id == uid,
        )
        .first()
    )
    now = datetime.now()
    if row is None:
        row = PptExportConfig(
            module_key  = MODULE_KEY,
            user_id     = uid,
            config_json = json.dumps(to_save, ensure_ascii=False),
            template_id = body.template_id,
            updated_by  = current_user.email,
            updated_at  = now,
        )
        db.add(row)
    else:
        row.config_json = json.dumps(to_save, ensure_ascii=False)
        row.template_id = body.template_id
        row.updated_by  = current_user.email
        row.updated_at  = now
    db.commit()
    return {"ok": True, "updated_at": now.isoformat()}


# ═════════════════════════════════════════════════════════════════════════════
# PPTX Builder v2 — Generic, Registry-driven
# ═════════════════════════════════════════════════════════════════════════════

def _sanitize_cell_text(val: str) -> str:
    """移除 XML 不允許的控制字元（CR / LF / Tab → 空格），防止 PowerPoint 開啟修復警告。"""
    return val.replace('\r\n', ' ').replace('\r', ' ').replace('\n', ' ').replace('\t', ' ')


def _build_generic_table_slide(
    slide, title: str, subtitle: str,
    columns: list[dict], rows: list[dict],
    now_str: str, SW: float, SH: float,
) -> None:
    """
    通用表格 Slide builder。
    columns: [{"key": str, "label": str, "width": float, "align": "left"|"right"|"center"}, ...]
    rows:    [{"欄位名稱": "值", ...}, ...]
    呼叫端需先分頁（每頁 ≤ ROWS_PER_SLIDE），此函式僅顯示傳入的 rows。
    """
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from app.routers.hotel_overview import (
        _set_slide_title, _pptx_cell, _pptx_header_row, _pptx_txt,
    )

    C_DARK    = RGBColor(0x1B, 0x3A, 0x5C)
    C_ROW_ALT = RGBColor(0xEE, 0xF5, 0xFB)
    C_GRAY    = RGBColor(0x88, 0x88, 0x88)

    TABLE_Y = 0.95
    TABLE_H = SH - TABLE_Y - 0.45
    TABLE_W = SW - 0.8

    _set_slide_title(slide, title, subtitle, now_str, SW, SH)

    if not rows or not columns:
        _pptx_txt(slide, "（本期暫無資料）",
                  2.0, 3.5, 9.0, 1.0, size=14, color=C_GRAY, italic=True)
        return

    n_cols = len(columns)
    n_rows = len(rows) + 1  # +1 header

    tbl = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(0.4), Inches(TABLE_Y), Inches(TABLE_W), Inches(TABLE_H)
    ).table

    # 欄寬設定（最後一欄取剩餘，避免超出版面）
    used = 0.0
    for ci, col in enumerate(columns):
        w = col.get("width", TABLE_W / n_cols)
        if ci == n_cols - 1:
            w = max(TABLE_W - used, 0.3)
        tbl.columns[ci].width = Inches(w)
        used += w

    # Header
    for ci, col in enumerate(columns):
        _pptx_cell(tbl, 0, ci, col.get("label", col["key"]), bold=True)
    _pptx_header_row(tbl, n_cols, size=9)

    # Data rows — 9pt 字體，並清除 CR/LF 避免 XML 警告
    _ALIGN_MAP = {"right": PP_ALIGN.RIGHT, "center": PP_ALIGN.CENTER}
    for ri, row in enumerate(rows, 1):
        bg = C_ROW_ALT if ri % 2 == 0 else None
        for ci, col in enumerate(columns):
            raw = str(row.get(col["key"], "") or "—")
            val = _sanitize_cell_text(raw)
            align = _ALIGN_MAP.get(col.get("align", "left"))
            _pptx_cell(tbl, ri, ci, val, fg=C_DARK, bg=bg,
                       size=9, align=align)

    # Row heights — header 稍高，資料列一致
    tbl.rows[0].height = Pt(26)
    for ri in range(1, n_rows):
        tbl.rows[ri].height = Pt(20)


def _make_chart_image(chart_type: str, data: list[dict], title: str) -> Optional[BytesIO]:
    """
    C-3：用 matplotlib 生成圖表 PNG，回傳 BytesIO。
    chart_type: "bar" | "hbar" | "line" | "pie"
    失敗時回傳 None（由呼叫端 fallback 至表格模式）。
    """
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        import matplotlib.font_manager as fm

        # 嘗試使用系統中文字型（Windows / macOS / Linux 均適用）
        # 優先順序：繁中 Windows → 簡中 Windows → macOS → Linux Noto
        _CJK_PRIORITY = [
            "Microsoft JhengHei",   # 微軟正黑體（Windows 繁中）
            "Microsoft JhengHei UI",
            "PMingLiU",             # 新細明體（Windows 繁中）
            "MingLiU",              # 細明體（Windows 繁中）
            "DFKai-SB",             # 標楷體（Windows 繁中）
            "Microsoft YaHei",      # 微軟雅黑（Windows 簡中）
            "Microsoft YaHei UI",
            "SimHei",               # 黑體（Windows 簡中）
            "SimSun",               # 宋體（Windows 簡中）
            "PingFang TC",          # macOS 繁中
            "PingFang SC",          # macOS 簡中
            "Heiti TC",             # macOS 繁中
            "Noto Sans CJK TC",     # Linux Noto 繁中
            "Noto Sans CJK SC",     # Linux Noto 簡中
            "WenQuanYi Micro Hei",  # Linux 文泉驛
        ]
        _available = {f.name for f in fm.fontManager.ttflist}
        _chosen_font = next((fn for fn in _CJK_PRIORITY if fn in _available), None)

        # 若優先清單找不到，退而求其次：關鍵字模糊比對
        if _chosen_font is None:
            _cjk_kw = ("jhenghei", "mingliu", "dfkai", "yahei", "simhei",
                       "pingfang", "heiti", "noto", "cjk", "wenquanyi",
                       "chinese", "gothic", "mincho", "taipei", "arial unicode")
            _chosen_font = next(
                (f.name for f in fm.fontManager.ttflist
                 if any(kw in f.name.lower() for kw in _cjk_kw)),
                None,
            )

        if _chosen_font:
            plt.rcParams["font.family"] = _chosen_font
            logger.info("[PPT] matplotlib 使用字型：%s", _chosen_font)
        else:
            logger.warning("[PPT] 找不到 CJK 字型，中文可能顯示為方塊")
        plt.rcParams["axes.unicode_minus"] = False

        C_BRAND  = "#1B3A5C"
        C_ACCENT = "#4BA8E8"
        C_ACCENT2 = "#F5A623"
        COLORS   = [C_BRAND, C_ACCENT, C_ACCENT2, "#2ECC71", "#E74C3C", "#9B59B6"]

        fig, ax = plt.subplots(figsize=(5.6, 3.2), dpi=150)
        fig.patch.set_facecolor("white")

        if chart_type == "hbar":
            # 橫條圖：各來源工項/完成數
            labels = [str(r.get("date", ""))[:8] for r in data]
            vals1  = [float(r.get("工項數", 0)) for r in data]
            vals2  = [float(r.get("完成數", 0)) for r in data]
            y = range(len(labels))
            ax.barh([i + 0.2 for i in y], vals1, 0.4, color=C_BRAND,   label="工項數")
            ax.barh([i - 0.2 for i in y], vals2, 0.4, color=C_ACCENT,  label="已完成")
            ax.set_yticks(list(y))
            ax.set_yticklabels(labels, fontsize=8)
            ax.legend(fontsize=7)
            ax.set_xlabel("件數", fontsize=8)

        elif chart_type == "bar":
            # 直條圖：完成率%
            labels = [str(r.get("date", ""))[:8] for r in data]
            vals   = [float(r.get("rate", 0)) for r in data]
            bars   = ax.bar(range(len(labels)), vals, color=C_BRAND, width=0.5)
            ax.set_xticks(range(len(labels)))
            ax.set_xticklabels(labels, fontsize=8, rotation=20, ha="right")
            ax.set_ylim(0, 110)
            ax.set_ylabel("%", fontsize=8)
            for bar, v in zip(bars, vals):
                ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 1,
                        f"{v:.1f}%", ha="center", va="bottom", fontsize=7)

        elif chart_type == "line":
            # 折線圖：趨勢
            labels = [str(r.get("date", ""))[-5:] for r in data]  # 取後5碼 "MM/DD" 或 "YY/MM"
            total  = [float(r.get("total",     0)) for r in data]
            done   = [float(r.get("completed", 0)) for r in data]
            x = range(len(labels))
            ax.plot(x, total, marker="o", color=C_BRAND,  linewidth=1.5, markersize=4, label="總件數")
            ax.plot(x, done,  marker="s", color=C_ACCENT, linewidth=1.5, markersize=4, label="完成件數")
            ax.set_xticks(list(x))
            ax.set_xticklabels(labels, fontsize=7, rotation=30, ha="right")
            ax.legend(fontsize=7)
            ax.set_ylabel("件數", fontsize=8)

        elif chart_type == "pie":
            # 圓餅圖：工時占比
            labels = [str(r.get("category", "")) for r in data]
            sizes  = [max(float(r.get("hours", 0)), 0) for r in data]
            if sum(sizes) == 0:
                sizes = [1] * len(sizes)
            wedge_colors = COLORS[:len(labels)]
            ax.pie(sizes, labels=labels, colors=wedge_colors,
                   autopct="%1.1f%%", pctdistance=0.75,
                   textprops={"fontsize": 7},
                   startangle=90)
            ax.axis("equal")

        ax.set_title(title, fontsize=9, fontweight="bold", color=C_BRAND, pad=8)
        ax.tick_params(labelsize=8)
        for spine in ax.spines.values():
            spine.set_linewidth(0.5)
            spine.set_color("#dddddd")

        fig.tight_layout(pad=0.8)
        buf = BytesIO()
        fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
        buf.seek(0)
        plt.close(fig)
        return buf
    except Exception as e:
        logger.warning("matplotlib chart generation failed (%s): %s", chart_type, e)
        return None


def _build_decision_charts_slide(
    slide, frontend_data: FrontendData,
    now_str: str, SW: float, SH: float,
) -> None:
    """
    C-3：Dashboard 決策分析圖表 → 4 張 matplotlib 圖形（2×2 版面）。
    若 matplotlib 圖形生成失敗，自動 fallback 至原有數據表格。
    """
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from app.routers.hotel_overview import (
        _set_slide_title, _pptx_txt, _pptx_cell, _pptx_header_row,
    )

    C_DARK    = RGBColor(0x1B, 0x3A, 0x5C)
    C_ROW_ALT = RGBColor(0xEE, 0xF5, 0xFB)
    C_GRAY    = RGBColor(0x88, 0x88, 0x88)

    _set_slide_title(slide, "決策分析 — 數據彙整", "各來源工項/完成率/趨勢/工時占比",
                     now_str, SW, SH)

    # 2×2 版面尺寸
    MARGIN_L = 0.35
    MARGIN_T = 1.0
    GAP_X    = 0.15
    GAP_Y    = 0.15
    IMG_W    = (SW - MARGIN_L * 2 - GAP_X) / 2
    IMG_H    = (SH - MARGIN_T - 0.35 - GAP_Y) / 2

    # 4 個圖表的定義：(資料, chart_type, title, grid_position)
    bar   = frontend_data.bar_chart_data   or []
    rate  = frontend_data.rate_chart_data  or []
    trend = frontend_data.dazhi_trend_data or []
    pie   = frontend_data.hours_pie_data   or []

    charts = [
        (bar,   "hbar", "① 各來源工項/案件數比較", 0, 0),
        (rate,  "bar",  "② 各來源完成率（%）",     1, 0),
        (trend, "line", "③ 工務12個月報修趨勢",   0, 1),
        (pie,   "pie",  "④ 工時來源占比",          1, 1),
    ]

    def _mini_table_fallback(title_text, headers, data_rows, ox, oy):
        """matplotlib 失敗時的表格 fallback"""
        LABEL_H = 0.22
        _pptx_txt(slide, title_text, ox, oy, IMG_W - 0.1, LABEL_H,
                  size=9, bold=True, color=C_DARK)
        n_cols = len(headers)
        n_rows = len(data_rows) + 1
        tbl = slide.shapes.add_table(
            n_rows, n_cols,
            Inches(ox), Inches(oy + LABEL_H + 0.04),
            Inches(IMG_W - 0.15), Inches(IMG_H - LABEL_H - 0.12),
        ).table
        col_w = (IMG_W - 0.15) / n_cols
        for ci in range(n_cols):
            tbl.columns[ci].width = Inches(col_w)
        for ci, h in enumerate(headers):
            _pptx_cell(tbl, 0, ci, h, bold=True, size=8)
        _pptx_header_row(tbl, n_cols, size=8)
        for ri, row in enumerate(data_rows[:10], 1):
            bg = C_ROW_ALT if ri % 2 == 0 else None
            for ci, val in enumerate(row):
                _pptx_cell(tbl, ri, ci, _sanitize_cell_text(str(val)),
                           fg=C_DARK, bg=bg, size=7)
        for ri in range(n_rows):
            tbl.rows[ri].height = Pt(20)

    for data, ctype, title, col, row_idx in charts:
        ox = MARGIN_L + col * (IMG_W + GAP_X)
        oy = MARGIN_T + row_idx * (IMG_H + GAP_Y)

        # C-3：嘗試 matplotlib
        img_buf = _make_chart_image(ctype, data, title) if data else None
        if img_buf:
            slide.shapes.add_picture(
                img_buf, Inches(ox), Inches(oy),
                width=Inches(IMG_W), height=Inches(IMG_H),
            )
        else:
            # fallback 表格
            if ctype == "hbar":
                _mini_table_fallback(title, ["來源","工項數","已完成"],
                    [[r.get("date",""), r.get("工項數",0), r.get("完成數",0)] for r in data],
                    ox, oy)
            elif ctype == "bar":
                _mini_table_fallback(title, ["來源","完成率%"],
                    [[r.get("date",""), r.get("rate",0)] for r in data],
                    ox, oy)
            elif ctype == "line":
                _mini_table_fallback(title, ["月份","總件數","完成件數"],
                    [[r.get("date",""), r.get("total",0), r.get("completed",0)] for r in data],
                    ox, oy)
            elif ctype == "pie":
                total_hrs = sum(r.get("hours", 0) for r in data) or 1
                _mini_table_fallback(title, ["來源","工時(HR)","占比%"],
                    [[r.get("category",""), r.get("hours",0),
                      f"{r.get('hours',0)/total_hrs*100:.1f}%"] for r in data],
                    ox, oy)


def _build_hotel_pptx_v2(
    params: dict,
    frontend_data: FrontendData,
    enabled_config: list[dict],
    db: Session,
    template_id: str = "default",
) -> BytesIO:
    """
    Registry-driven PPTX builder v2。

    Parameters
    ----------
    params:
        {"year": int, "month": int, "inspection_date": str}
    frontend_data:
        前端帶入的 frontend_payload sections 資料
    enabled_config:
        已按 sort_order 排序、僅包含 enabled=True 的 config list
    db:
        SQLAlchemy Session（供 backend_db sections 查詢）
    """
    import os
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from app.routers.hotel_overview import (
        _clone_template_slide, _delete_slide, _update_cover_date,
        _set_slide_title, _build_slide2_kpi,
        get_hotel_daily_hours, get_hotel_monthly_hours, get_hotel_person_hours,
        _pptx_txt, _pptx_cell, _pptx_header_row,
        HotelPptxPayload, KpiSummaryIn, SourceCardIn, RepairCostsIn,
    )

    # C-2：依 template_id 選擇模板檔案
    _tpl_meta = next(
        (t for t in AVAILABLE_TEMPLATES if t["id"] == template_id),
        AVAILABLE_TEMPLATES[0],
    )
    TEMPLATE_PATH = os.path.abspath(
        os.path.join(os.path.dirname(__file__), '..', 'static',
                     'pptx_templates', _tpl_meta["filename"])
    )
    prs  = Presentation(TEMPLATE_PATH)
    SW   = prs.slide_width.inches
    SH   = prs.slide_height.inches
    now_str    = datetime.now().strftime("%Y-%m-%d %H:%M")
    year       = params["year"]
    month      = params["month"]
    period_str = f"{year}年{month:02d}月"

    C_GRAY = RGBColor(0x88, 0x88, 0x88)

    # ── Slide 1: Cover（固定）──────────────────────────────────────────────
    _update_cover_date(prs.slides[0], year, month)

    # ── 預先算出需要 clone 幾次（每個 enabled section = 1 張，決策圖表合 1 張）──
    # 簡單做法：先 clone 足夠多張，之後再刪除多餘的（python-pptx 刪除較難）
    # 這裡改為逐一 clone + build，不預先 clone

    # Dashboard Sections 的特殊合併規則：
    # dashboard_kpi_summary / dashboard_source_status / dashboard_repair_costs → 同一頁
    # dashboard_bar_chart / dashboard_rate_chart / dashboard_dazhi_trend / dashboard_hours_pie → 同一頁

    enabled_keys  = {c["export_key"] for c in enabled_config}
    detail_keys   = {c["export_key"] for c in enabled_config if c.get("include_detail")}
    sorted_config = sorted(enabled_config, key=lambda c: c.get("sort_order", 99))

    # 決定要輸出的 "slide groups"（按排序後順序）
    DASHBOARD_KPI_KEYS    = {"dashboard_kpi_summary", "dashboard_source_status", "dashboard_repair_costs"}
    DASHBOARD_CHART_KEYS  = {"dashboard_bar_chart", "dashboard_rate_chart", "dashboard_dazhi_trend", "dashboard_hours_pie"}

    slide_groups: list[str] = []  # group keys
    kpi_group_added   = False
    chart_group_added = False

    for item in sorted_config:
        k = item["export_key"]
        if k in DASHBOARD_KPI_KEYS:
            if not kpi_group_added:
                slide_groups.append("__kpi_group__")
                kpi_group_added = True
        elif k in DASHBOARD_CHART_KEYS:
            if not chart_group_added:
                slide_groups.append("__chart_group__")
                chart_group_added = True
        else:
            slide_groups.append(k)
            if k in detail_keys:
                slide_groups.append(f"__detail__{k}")

    # ── 刪除 TOC slide（index 1），保留 cover(0) + content_template(2)──
    _delete_slide(prs, 1)
    # 現在: [0]=cover, [1]=content_template

    # ── 逐一 clone + build ───────────────────────────────────────────────
    for group_key in slide_groups:
        slide = _clone_template_slide(prs, 1)  # clone content_template（始終是 index 1）

        # ── Slide: KPI 總覽（主管摘要 + 各來源狀態 + 費用）──────────────
        if group_key == "__kpi_group__":
            _set_slide_title(slide, "本期績效總覽",
                             f"{period_str}  主管摘要 · 各來源狀態 · 報修費用",
                             now_str, SW, SH)
            fd = frontend_data
            if (fd.kpi_summary and fd.source_cards and fd.repair_costs):
                try:
                    # repair_costs 支援兩種格式：
                    #   ① dict  → {"outsource_fee": x, "maintenance_fee": y, "deduction_fee": z}
                    #   ② list  → [{"category": "外包費", "amount": x}, ...]（前端 Dashboard 格式）
                    _CAT_MAP = {"外包費": "outsource_fee", "保養費": "maintenance_fee", "扣款費": "deduction_fee"}
                    if isinstance(fd.repair_costs, list):
                        costs_dict = {
                            _CAT_MAP.get(item.get("category", ""), item.get("category", "")): item.get("amount", 0)
                            for item in fd.repair_costs
                        }
                    else:
                        costs_dict = fd.repair_costs
                    kpi_obj = HotelPptxPayload(
                        kpi_summary  = KpiSummaryIn(**fd.kpi_summary),
                        source_cards = [SourceCardIn(**c) for c in fd.source_cards],
                        repair_costs = RepairCostsIn(**costs_dict),
                    )
                    _build_slide2_kpi(slide, kpi_obj, period_str, SW=SW, SH=SH)
                except Exception as e:
                    logger.warning("KPI slide build error: %s", e)
            continue

        # ── Slide: 決策分析圖表（4 個數據表格）────────────────────────
        if group_key == "__chart_group__":
            _build_decision_charts_slide(slide, frontend_data, now_str, SW, SH)
            continue

        # ── Slide: 明細附錄 ─────────────────────────────────────────────
        if group_key.startswith("__detail__"):
            base_key    = group_key[len("__detail__"):]
            section_def = get_section_def(base_key)
            if section_def is None:
                continue
            result = fetch_section_data(base_key, db, params, include_detail=True)
            detail_rows = result.get("detail", []) if result else []
            _build_generic_table_slide(
                slide,
                title    = f"{section_def.second_title} — 明細附錄",
                subtitle = period_str,
                columns  = section_def.columns,
                rows     = detail_rows,
                now_str  = now_str, SW=SW, SH=SH,
            )
            continue

        # ── Slide: 一般 backend_db section ─────────────────────────────
        section_def = get_section_def(group_key)
        if section_def is None:
            continue

        if section_def.data_source == "backend_db":
            result = fetch_section_data(group_key, db, params)
            rows   = result.get("main", []) if result else []
            # 人員工時佔比：欄位動態（依 DB 人員名單），需呼叫 _person_pct_cols()
            if group_key == "staff_hours_percent":
                columns = _person_pct_cols(db, params)
            else:
                columns = section_def.columns

            # ── 分頁邏輯：每頁最多 ROWS_PER_SLIDE 筆 ─────────────────────────
            ROWS_PER_SLIDE = 20
            if len(rows) > ROWS_PER_SLIDE:
                pages       = [rows[i:i+ROWS_PER_SLIDE]
                                for i in range(0, len(rows), ROWS_PER_SLIDE)]
                total_pages = len(pages)
                for pg_idx, pg_rows in enumerate(pages):
                    if pg_idx > 0:
                        # 第 2 頁起另外 clone 一張
                        slide = _clone_template_slide(prs, 1)
                    pg_subtitle = (
                        f"{period_str}　第 {pg_idx + 1} 頁 / 共 {total_pages} 頁"
                        f"　（共 {len(rows)} 筆）"
                    )
                    _build_generic_table_slide(
                        slide,
                        title    = section_def.second_title,
                        subtitle = pg_subtitle,
                        columns  = columns,
                        rows     = pg_rows,
                        now_str  = now_str, SW=SW, SH=SH,
                    )
            else:
                _build_generic_table_slide(
                    slide,
                    title    = section_def.second_title,
                    subtitle = period_str,
                    columns  = columns,
                    rows     = rows,
                    now_str  = now_str, SW=SW, SH=SH,
                )

        # frontend_payload 個別 sections（日後擴充用；目前均已合入 kpi_group）
        else:
            _set_slide_title(slide, section_def.second_title, period_str, now_str, SW, SH)
            _pptx_txt(slide, "（此 section 由前端資料產生，版型開發中）",
                      2.0, 3.5, 9.0, 1.0, size=14, color=C_GRAY, italic=True)

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ═════════════════════════════════════════════════════════════════════════════
# POST /export
# ═════════════════════════════════════════════════════════════════════════════

@router.post("/export", summary="觸發飯店 Dashboard PPT 匯出（Registry-driven v2）")
def export_hotel_ppt(
    body:         ExportBody = Body(...),
    db:           Session    = Depends(get_db),
    current_user: User       = Depends(get_current_user),
):
    try:
        config_list, template_id = _load_merged_config(db, user_id=str(current_user.id))
        enabled_config  = [c for c in config_list if c.get("enabled")]
        buf = _build_hotel_pptx_v2(
            params         = {"year": body.year, "month": body.month,
                              "inspection_date": body.inspection_date},
            frontend_data  = body.frontend_data,
            enabled_config = enabled_config,
            db             = db,
            template_id    = template_id,
        )
    except Exception as exc:
        logger.error("PPT export v2 failed:\n%s", _tb.format_exc())
        raise HTTPException(status_code=500, detail=f"PPT 建立失敗：{exc}")

    # B-3：寫入匯出歷史紀錄
    try:
        history = PptExportHistory(
            module_key    = MODULE_KEY,
            year          = body.year,
            month         = body.month,
            exported_by   = current_user.email,
            exported_at   = datetime.now(),
            sections_json = json.dumps(
                [c["export_key"] for c in enabled_config], ensure_ascii=False
            ),
            template_id   = template_id,
        )
        db.add(history)
        db.commit()
    except Exception:
        logger.warning("Failed to write export history: %s", _tb.format_exc())

    filename = f"飯店管理報告_{body.year}年{body.month:02d}月.pptx"
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename*=UTF-8''{quote(filename)}"},
    )


# ═════════════════════════════════════════════════════════════════════════════
# GET /history  (B-3)
# ═════════════════════════════════════════════════════════════════════════════

@router.get("/history", summary="取得最近 PPT 匯出歷史（最新 30 筆）")
def get_export_history(
    limit:        int     = Query(default=30, le=100),
    db:           Session = Depends(get_db),
    current_user: User    = Depends(get_current_user),
):
    rows = (
        db.query(PptExportHistory)
        .filter(PptExportHistory.module_key == MODULE_KEY)
        .order_by(PptExportHistory.exported_at.desc())
        .limit(limit)
        .all()
    )
    return [
        {
            "id":           r.id,
            "year":         r.year,
            "month":        r.month,
            "exported_by":  r.exported_by,
            "exported_at":  r.exported_at.isoformat(),
            "template_id":  r.template_id,
            "sections":     json.loads(r.sections_json or "[]"),
            "section_count": len(json.loads(r.sections_json or "[]")),
        }
        for r in rows
    ]


# ═════════════════════════════════════════════════════════════════════════════
# GET /templates  (C-2)
# ═════════════════════════════════════════════════════════════════════════════

import os as _os

_TEMPLATE_DIR = _os.path.abspath(
    _os.path.join(_os.path.dirname(__file__), '..', 'static', 'pptx_templates')
)

# 可用模板清單（key → 顯示名稱 + 檔案名稱）
AVAILABLE_TEMPLATES: list[dict] = [
    {
        "id":          "default",
        "label":       "預設模板（深藍主色）",
        "filename":    "hotel_report_template.pptx",
        "description": "品牌主色 #1B3A5C，標準版面",
    },
    # 未來新增模板範例（檔案放至 static/pptx_templates/ 後取消註解）:
    # {
    #     "id":       "minimal",
    #     "label":    "極簡白底版",
    #     "filename": "hotel_report_minimal.pptx",
    #     "description": "白底黑字，適合對外報告",
    # },
]


@router.get("/templates", summary="取得可用 PPTX 模板清單（C-2）")
def get_ppt_templates(current_user: User = Depends(get_current_user)):
    result = []
    for t in AVAILABLE_TEMPLATES:
        path = _os.path.join(_TEMPLATE_DIR, t["filename"])
        result.append({**t, "available": _os.path.exists(path)})
    return result


# ═════════════════════════════════════════════════════════════════════════════
# Section 註冊 — hotel_overview 原有 12 個
# ═════════════════════════════════════════════════════════════════════════════
#
# 命名規則：
#   data_source="frontend_payload" → 資料由前端 ExportBody.frontend_data 帶入
#   data_source="backend_db"       → 後端在匯出時呼叫 data_provider 查 DB
#
# ─────────────────────────────────────────────────────────────────────────────
# Dashboard 分組（frontend_payload）
# ─────────────────────────────────────────────────────────────────────────────

register_section(PptSectionDef(
    export_key="dashboard_kpi_summary", module_key=MODULE_KEY,
    tab_name="Dashboard", second_title="主管摘要",
    description="本期 5 個 KPI 指標：總工項、已完成、工時合計、異常、逾期",
    export_type="kpi_cards", slide_layout="kpi_summary",
    data_source="frontend_payload", sort_order=1,
))

register_section(PptSectionDef(
    export_key="dashboard_source_status", module_key=MODULE_KEY,
    tab_name="Dashboard", second_title="各來源本期狀態",
    description="4 個來源的工項數、完成率、工時卡片（含 2 個佔位卡）",
    export_type="kpi_cards", slide_layout="chart_two_col",
    data_source="frontend_payload", sort_order=2,
))

register_section(PptSectionDef(
    export_key="dashboard_repair_costs", module_key=MODULE_KEY,
    tab_name="Dashboard", second_title="報修費用摘要",
    description="委外費用、維修費用、扣款費用、本月合計（累計至篩選月份）",
    export_type="summary_text", slide_layout="summary_page",
    data_source="frontend_payload", sort_order=3,
))

register_section(PptSectionDef(
    export_key="dashboard_bar_chart", module_key=MODULE_KEY,
    tab_name="Dashboard", second_title="各來源工項/案件數比較",
    description="各來源工項總數與已完成數對比（數據表格，非圖形）",
    export_type="table", slide_layout="chart_full",
    data_source="frontend_payload", sort_order=4,
))

register_section(PptSectionDef(
    export_key="dashboard_rate_chart", module_key=MODULE_KEY,
    tab_name="Dashboard", second_title="各來源完成率比較",
    description="各來源完成率%（數據表格）",
    export_type="table", slide_layout="chart_full",
    data_source="frontend_payload", sort_order=5,
))

register_section(PptSectionDef(
    export_key="dashboard_dazhi_trend", module_key=MODULE_KEY,
    tab_name="Dashboard", second_title="工務12個月報修趨勢",
    description="工務部最近12個月報修總件數與完成件數趨勢（數據表格）",
    export_type="table", slide_layout="chart_full",
    data_source="frontend_payload", sort_order=6,
))

register_section(PptSectionDef(
    export_key="dashboard_hours_pie", module_key=MODULE_KEY,
    tab_name="Dashboard", second_title="工時來源占比",
    description="各來源工時(HR)與占比%（數據表格）",
    export_type="table", slide_layout="summary_page",
    data_source="frontend_payload", sort_order=7,
))

# ─────────────────────────────────────────────────────────────────────────────
# B. 每日累計（backend_db）
# ─────────────────────────────────────────────────────────────────────────────

def _provide_daily(db: Session, params: dict) -> list[dict]:
    from app.routers.hotel_overview import get_hotel_daily_hours
    data = get_hotel_daily_hours(params["year"], params["month"], db)
    rows = []
    for row in data["rows"]:
        hrs = row["hours"]
        max_val = max(hrs) if hrs else 0
        peak_day = data["days"][hrs.index(max_val)] if max_val > 0 else "—"
        rows.append({
            "來源":       row["category"],
            "月合計(HR)": f"{row['total']:.1f}",
            "占比%":      f"{row['pct']:.1f}",
            "最高工時日": f"第 {peak_day} 日" if peak_day != "—" else "—",
            "最高日(HR)": f"{max_val:.1f}" if max_val > 0 else "—",
        })
    return rows


def _provide_daily_detail(db: Session, params: dict) -> list[dict]:
    """每日累計明細：完整 N日 × 各來源 交叉表（以每日工時為主）"""
    from app.routers.hotel_overview import get_hotel_daily_hours
    data = get_hotel_daily_hours(params["year"], params["month"], db)
    rows = []
    for row in data["rows"]:
        entry = {"來源": row["category"]}
        for i, day in enumerate(data["days"]):
            val = row["hours"][i] if i < len(row["hours"]) else 0
            entry[f"{day}日"] = f"{val:.1f}" if val > 0 else "—"
        entry["合計"] = f"{row['total']:.1f}"
        rows.append(entry)
    return rows


_daily_cols = [
    {"key": "來源",       "label": "來源",       "width": 2.4},
    {"key": "月合計(HR)", "label": "月合計(HR)", "width": 1.8, "align": "right"},
    {"key": "占比%",      "label": "占比%",      "width": 1.2, "align": "right"},
    {"key": "最高工時日", "label": "最高工時日", "width": 2.0, "align": "center"},
    {"key": "最高日(HR)", "label": "最高日(HR)", "width": 1.8, "align": "right"},
]

register_section(
    PptSectionDef(
        export_key="daily_accumulate_table", module_key=MODULE_KEY,
        tab_name="B. 每日累計", second_title="每日累計工時摘要",
        description="各來源月合計工時、占比、最高工時日",
        export_type="table", slide_layout="table_full",
        supports_detail=True,
        detail_description="附加完整 N日×來源交叉明細表（以工時 HR 為主）",
        data_source="backend_db", sort_order=8,
        columns=_daily_cols,
    ),
    data_provider=_provide_daily,
    detail_provider=_provide_daily_detail,
)

# ─────────────────────────────────────────────────────────────────────────────
# C. 每月累計（backend_db）
# ─────────────────────────────────────────────────────────────────────────────

def _provide_monthly(db: Session, params: dict) -> list[dict]:
    from app.routers.hotel_overview import get_hotel_monthly_hours
    data = get_hotel_monthly_hours(params["year"], db)
    rows = []
    for row in data["rows"]:
        entry = {"來源": row["category"]}
        for i, m in enumerate(range(1, 13)):
            val = row["hours"][i] if i < len(row["hours"]) else 0
            entry[f"{m}月"] = f"{val:.1f}" if val > 0 else "—"
        entry["全年合計"] = f"{row['total']:.1f}"
        rows.append(entry)
    return rows


def _provide_monthly_detail(db: Session, params: dict) -> list[dict]:
    """每月累計明細：同時呈現工時(HR) + 案件數的對照"""
    from app.routers.hotel_overview import get_hotel_monthly_hours
    data = get_hotel_monthly_hours(params["year"], db)
    rows = []
    for row in data["rows"]:
        # 工時列
        hrs_entry = {"類型": f"{row['category']}（工時HR）"}
        for i, m in enumerate(range(1, 13)):
            val = row["hours"][i] if i < len(row["hours"]) else 0
            hrs_entry[f"{m}月"] = f"{val:.1f}" if val > 0 else "—"
        hrs_entry["合計"] = f"{row['total']:.1f}"
        rows.append(hrs_entry)
        # 案件數列
        case_entry = {"類型": f"{row['category']}（案件數）"}
        cases = row.get("cases", [])
        for i, m in enumerate(range(1, 13)):
            val = cases[i] if i < len(cases) else 0
            case_entry[f"{m}月"] = str(val) if val > 0 else "—"
        case_entry["合計"] = str(row.get("cases_total", 0))
        rows.append(case_entry)
    return rows


_monthly_cols = (
    [{"key": "來源", "label": "來源", "width": 2.2}] +
    [{"key": f"{m}月", "label": f"{m}月", "width": 0.76, "align": "right"} for m in range(1, 13)] +
    [{"key": "全年合計", "label": "全年合計", "width": 1.0, "align": "right"}]
)

register_section(
    PptSectionDef(
        export_key="monthly_accumulate_table", module_key=MODULE_KEY,
        tab_name="C. 每月累計", second_title="每月工時累計",
        description="各來源 × 12 個月工時(HR)彙總",
        export_type="table", slide_layout="table_full",
        supports_detail=True,
        detail_description="附加工時(HR) + 案件數對照表（雙行呈現）",
        data_source="backend_db", sort_order=9,
        columns=_monthly_cols,
    ),
    data_provider=_provide_monthly,
    detail_provider=_provide_monthly_detail,
)

# ─────────────────────────────────────────────────────────────────────────────
# D. 每年累計（backend_db，Running Total，年份同 Dashboard year）
# ─────────────────────────────────────────────────────────────────────────────

def _provide_yearly(db: Session, params: dict) -> list[dict]:
    from app.routers.hotel_overview import get_hotel_monthly_hours
    data = get_hotel_monthly_hours(params["year"], db)
    rows = []
    for row in data["rows"]:
        entry = {"來源": row["category"]}
        cases = row.get("cases", [])
        cumsum = 0
        for i, m in enumerate(range(1, 13)):
            cumsum += (cases[i] if i < len(cases) else 0)
            entry[f"{m}月"] = str(cumsum) if cumsum > 0 else "—"
        entry["全年合計"] = str(cumsum)
        rows.append(entry)
    return rows


def _provide_yearly_detail(db: Session, params: dict) -> list[dict]:
    """每年累計明細：月度實際值 vs 累計值對照"""
    from app.routers.hotel_overview import get_hotel_monthly_hours
    data = get_hotel_monthly_hours(params["year"], db)
    rows = []
    for row in data["rows"]:
        cases = row.get("cases", [])
        # 月度實際值
        actual_entry = {"說明": f"{row['category']}（月度實際）"}
        for i, m in enumerate(range(1, 13)):
            val = cases[i] if i < len(cases) else 0
            actual_entry[f"{m}月"] = str(val) if val > 0 else "—"
        actual_entry["全年"] = str(sum(cases))
        rows.append(actual_entry)
        # 累計值
        cum_entry = {"說明": f"{row['category']}（累計至月）"}
        cumsum = 0
        for i, m in enumerate(range(1, 13)):
            cumsum += (cases[i] if i < len(cases) else 0)
            cum_entry[f"{m}月"] = str(cumsum) if cumsum > 0 else "—"
        cum_entry["全年"] = str(cumsum)
        rows.append(cum_entry)
    return rows


_yearly_cols = (
    [{"key": "來源", "label": "來源", "width": 2.2}] +
    [{"key": f"{m}月", "label": f"{m}月", "width": 0.76, "align": "right"} for m in range(1, 13)] +
    [{"key": "全年合計", "label": "全年合計", "width": 1.0, "align": "right"}]
)

register_section(
    PptSectionDef(
        export_key="yearly_accumulate_table", module_key=MODULE_KEY,
        tab_name="D. 每年累計", second_title="年度累計案件數",
        description="各來源 × 12個月 Running Total（累計至各月）",
        export_type="table", slide_layout="table_full",
        supports_detail=True,
        detail_description="附加月度實際值 vs 累計值對照表",
        data_source="backend_db", sort_order=10,
        columns=_yearly_cols,
    ),
    data_provider=_provide_yearly,
    detail_provider=_provide_yearly_detail,
)

# ─────────────────────────────────────────────────────────────────────────────
# 人員工時%（backend_db）
# ─────────────────────────────────────────────────────────────────────────────

def _provide_person_pct(db: Session, params: dict) -> list[dict]:
    from app.routers.hotel_overview import get_hotel_person_hours
    data = get_hotel_person_hours(params["year"], db)
    persons = data.get("persons", [])
    rows = []
    for row in data.get("rows", []):
        entry = {"來源": row["category"]}
        for i, name in enumerate(persons):
            pct = row["pct_by_person"][i] if i < len(row["pct_by_person"]) else 0
            entry[name] = f"{pct:.1f}%" if pct > 0 else "—"
        rows.append(entry)
    return rows


def _provide_person_pct_detail(db: Session, params: dict) -> list[dict]:
    """人員工時% 明細：實際工時(HR) 對照"""
    from app.routers.hotel_overview import get_hotel_person_hours
    data   = get_hotel_person_hours(params["year"], db)
    persons = data.get("persons", [])
    totals  = data.get("person_totals", [])
    rows = []
    for row in data.get("rows", []):
        entry = {"來源": row["category"]}
        for i, name in enumerate(persons):
            pct = row["pct_by_person"][i] if i < len(row["pct_by_person"]) else 0
            hrs = round(pct / 100 * totals[i], 1) if i < len(totals) else 0
            entry[name] = f"{hrs:.1f}" if hrs > 0 else "—"
        rows.append(entry)
    return rows


def _person_pct_cols(db: Session, params: dict) -> list[dict]:
    from app.routers.hotel_overview import get_hotel_person_hours
    data = get_hotel_person_hours(params["year"], db)
    persons = data.get("persons", [])
    return (
        [{"key": "來源", "label": "來源", "width": 2.0}] +
        [{"key": name, "label": name, "width": 0.8, "align": "center"} for name in persons]
    )

# 人員工時% 的 columns 在 register 時尚未知（需查 DB），用空 list；builder 會動態處理
register_section(
    PptSectionDef(
        export_key="staff_hours_percent", module_key=MODULE_KEY,
        tab_name="人員工時%", second_title="人員工時佔比",
        description="各來源 × Top-15 人員工時佔比%",
        export_type="table", slide_layout="table_full",
        supports_detail=True,
        detail_description="附加各人員實際工時(HR) 對照表",
        data_source="backend_db", sort_order=11,
        columns=[],  # 動態欄位，由 data_provider 提供
    ),
    data_provider=_provide_person_pct,
    detail_provider=_provide_person_pct_detail,
)

# ─────────────────────────────────────────────────────────────────────────────
# 人員排名（backend_db）
# ─────────────────────────────────────────────────────────────────────────────

def _provide_ranking(db: Session, params: dict) -> list[dict]:
    from app.routers.hotel_overview import get_hotel_person_hours
    data    = get_hotel_person_hours(params["year"], db)
    persons = data.get("persons", [])
    totals  = data.get("person_totals", [])
    grand   = sum(totals) or 1
    medals  = ["🥇", "🥈", "🥉"]
    return [
        {
            "排名":       medals[i] if i < 3 else str(i + 1),
            "姓名":       name,
            "全年工時(HR)": f"{totals[i]:.1f}" if i < len(totals) else "0.0",
            "全年占比%":  f"{totals[i]/grand*100:.1f}%" if i < len(totals) else "0.0%",
        }
        for i, name in enumerate(persons)
    ]


def _provide_ranking_detail(db: Session, params: dict) -> list[dict]:
    """人員排名明細：各人員 × 各來源工時"""
    from app.routers.hotel_overview import get_hotel_person_hours
    data    = get_hotel_person_hours(params["year"], db)
    persons = data.get("persons", [])
    totals  = data.get("person_totals", [])
    medals  = ["🥇", "🥈", "🥉"]
    rows    = []
    for pi, name in enumerate(persons):
        entry = {
            "排名": medals[pi] if pi < 3 else str(pi + 1),
            "姓名": name,
            "合計(HR)": f"{totals[pi]:.1f}" if pi < len(totals) else "0.0",
        }
        for row in data.get("rows", []):
            pct  = row["pct_by_person"][pi] if pi < len(row["pct_by_person"]) else 0
            hrs  = round(pct / 100 * (totals[pi] if pi < len(totals) else 0), 1)
            entry[row["category"]] = f"{hrs:.1f}" if hrs > 0 else "—"
        rows.append(entry)
    return rows


_ranking_cols = [
    {"key": "排名",        "label": "排名",        "width": 0.7, "align": "center"},
    {"key": "姓名",        "label": "姓名",        "width": 1.8},
    {"key": "全年工時(HR)", "label": "全年工時(HR)", "width": 1.6, "align": "right"},
    {"key": "全年占比%",   "label": "占比%",       "width": 1.2, "align": "right"},
]

register_section(
    PptSectionDef(
        export_key="staff_ranking", module_key=MODULE_KEY,
        tab_name="人員排名", second_title="人員工時排名",
        description="Top-15 人員全年工時排名（含名次、工時、占比）",
        export_type="table", slide_layout="table_full",
        supports_detail=True,
        detail_description="附加各人員 × 各來源工時明細表",
        data_source="backend_db", sort_order=12,
        columns=_ranking_cols,
    ),
    data_provider=_provide_ranking,
    detail_provider=_provide_ranking_detail,
)

# ─────────────────────────────────────────────────────────────────────────────
# 報修管理 — 新增 Section ①：報修未完成報表
# ─────────────────────────────────────────────────────────────────────────────

def _provide_repair_unfinished(db: Session, params: dict) -> list[dict]:
    """未完成的報修工單（排除取消案件）"""
    from app.models.dazhi_repair import DazhiRepairCase
    # is_completed_flag 是 @property，不能用於 SQLAlchemy filter
    # 正確做法：completed_at IS NULL AND is_completed = False
    cases = (
        db.query(DazhiRepairCase)
        .filter(
            DazhiRepairCase.completed_at == None,   # noqa: E711
            DazhiRepairCase.is_completed == False,  # noqa: E712
        )
        .order_by(DazhiRepairCase.occurred_at)
        .all()
    )
    rows = []
    for c in cases:
        if (c.status or "").strip() == "取消":
            continue
        rows.append({
            "報修編號": c.case_no     or "—",
            "報修日期": c.occurred_at.strftime("%Y/%m/%d") if c.occurred_at else "—",
            "地點":     c.floor       or "—",
            "工作說明": (c.title or "—")[:30],  # 截斷過長文字
            "目前狀態": c.status      or "—",
            "負責人":   c.acceptor    or "—",
        })
    return rows


def _provide_repair_unfinished_detail(db: Session, params: dict) -> list[dict]:
    """未完成報修明細：含費用欄位"""
    from app.models.dazhi_repair import DazhiRepairCase
    cases = (
        db.query(DazhiRepairCase)
        .filter(
            DazhiRepairCase.completed_at == None,   # noqa: E711
            DazhiRepairCase.is_completed == False,  # noqa: E712
        )
        .order_by(DazhiRepairCase.occurred_at)
        .all()
    )
    rows = []
    for c in cases:
        if (c.status or "").strip() == "取消":
            continue
        rows.append({
            "報修編號": c.case_no     or "—",
            "報修日期": c.occurred_at.strftime("%Y/%m/%d") if c.occurred_at else "—",
            "地點":     c.floor       or "—",
            "工作說明": (c.title or "—")[:25],
            "目前狀態": c.status      or "—",
            "負責人":   c.acceptor    or "—",
            "委外費用": f"${c.outsource_fee:,.0f}"    if c.outsource_fee    else "—",
            "維修費用": f"${c.maintenance_fee:,.0f}"  if c.maintenance_fee  else "—",
        })
    return rows


_unfinished_cols = [
    {"key": "報修編號", "label": "報修編號", "width": 1.3},
    {"key": "報修日期", "label": "報修日期", "width": 1.3},
    {"key": "地點",     "label": "地點",     "width": 1.8},
    {"key": "工作說明", "label": "工作說明", "width": 3.5},
    {"key": "目前狀態", "label": "狀態",     "width": 1.2},
    {"key": "負責人",   "label": "負責人",   "width": 1.4},
]

register_section(
    PptSectionDef(
        export_key="repair_unfinished_report", module_key=MODULE_KEY,
        tab_name="報修管理", second_title="報修未完成報表",
        description="目前狀態為未完成的報修工單清單（排除取消案件）",
        export_type="table", slide_layout="table_full",
        supports_detail=True,
        detail_description="附加委外費用 / 維修費用欄位",
        data_source="backend_db", sort_order=13,
        columns=_unfinished_cols,
    ),
    data_provider=_provide_repair_unfinished,
    detail_provider=_provide_repair_unfinished_detail,
)

# ─────────────────────────────────────────────────────────────────────────────
# 報修管理 — 新增 Section ②：本月結案工單
# ─────────────────────────────────────────────────────────────────────────────

def _provide_repair_closed(db: Session, params: dict) -> list[dict]:
    """本月結案工單（以 completed_at 歸月）"""
    from app.models.dazhi_repair import DazhiRepairCase
    year, month = params["year"], params["month"]
    cases = db.query(DazhiRepairCase).all()
    rows  = []
    for c in cases:
        if (c.status or "").strip() == "取消":
            continue
        if not c.is_completed_flag:
            continue
        if not (c.completed_at and c.completed_at.year == year and c.completed_at.month == month):
            continue
        rows.append({
            "報修編號": c.case_no     or "—",
            "報修日期": c.occurred_at.strftime("%Y/%m/%d") if c.occurred_at else "—",
            "結案日期": c.completed_at.strftime("%Y/%m/%d"),
            "地點":     c.floor       or "—",
            "工作說明": (c.title or "—")[:28],
            "維修工時": f"{c.work_hours:.1f} HR" if c.work_hours else "—",
            "負責人":   c.acceptor    or "—",
        })
    return rows


def _provide_repair_closed_detail(db: Session, params: dict) -> list[dict]:
    """本月結案工單明細：附加委外費用 / 維修費用欄位"""
    from app.models.dazhi_repair import DazhiRepairCase
    year, month = params["year"], params["month"]
    cases = db.query(DazhiRepairCase).all()
    rows  = []
    for c in cases:
        if (c.status or "").strip() == "取消":
            continue
        if not c.is_completed_flag:
            continue
        if not (c.completed_at and c.completed_at.year == year and c.completed_at.month == month):
            continue
        rows.append({
            "報修編號": c.case_no     or "—",
            "報修日期": c.occurred_at.strftime("%Y/%m/%d") if c.occurred_at else "—",
            "結案日期": c.completed_at.strftime("%Y/%m/%d"),
            "地點":     c.floor       or "—",
            "工作說明": (c.title or "—")[:25],
            "維修工時": f"{c.work_hours:.1f} HR" if c.work_hours else "—",
            "負責人":   c.acceptor    or "—",
            "委外費用": f"${c.outsource_fee:,.0f}"   if c.outsource_fee   else "—",
            "維修費用": f"${c.maintenance_fee:,.0f}" if c.maintenance_fee else "—",
        })
    return rows


_closed_cols = [
    {"key": "報修編號", "label": "報修編號", "width": 1.3},
    {"key": "報修日期", "label": "報修日期", "width": 1.3},
    {"key": "結案日期", "label": "結案日期", "width": 1.3},
    {"key": "地點",     "label": "地點",     "width": 1.7},
    {"key": "工作說明", "label": "工作說明", "width": 3.2},
    {"key": "維修工時", "label": "工時",     "width": 1.0, "align": "right"},
    {"key": "負責人",   "label": "負責人",   "width": 1.3},
]

register_section(
    PptSectionDef(
        export_key="repair_closed_list", module_key=MODULE_KEY,
        tab_name="報修管理", second_title="本月結案工單",
        description="本月（篩選年月）已完成結案的報修工單清單",
        export_type="table", slide_layout="table_full",
        supports_detail=True,
        detail_description="附加委外費用 / 維修費用欄位",
        data_source="backend_db", sort_order=14,
        columns=_closed_cols,
    ),
    data_provider=_provide_repair_closed,
    detail_provider=_provide_repair_closed_detail,
)


# =============================================================================
# 大直工務部 — Section 15：年度報修統計表（對應前端 3.1 報修 TAB）
# =============================================================================

def _provide_dazhi_repair_monthly_stats(db: Session, params: dict) -> list[dict]:
    """
    大直工務部 3.1 報修統計：全年 12 月份 × 9 欄交叉表。
    直接呼叫 compute_repair_stats() — 確保口徑與前端 TAB 完全一致。
    year 取自 params["year"]；month 不用於此 Section（全年展開）。
    """
    from app.models.dazhi_repair import DazhiRepairCase
    from app.services.dazhi_repair_service import compute_repair_stats

    year  = params.get("year", datetime.now().year)
    cases = db.query(DazhiRepairCase).all()
    stats = compute_repair_stats(cases, year)

    rows: list[dict] = []
    sum_closed_from_prev = 0
    sum_this_total       = 0
    sum_this_completed   = 0
    sum_this_uncompleted = 0

    for m in range(1, 13):
        d  = stats["months"].get(m, {})
        pv = d.get("prev_uncompleted",       0) or 0
        cf = d.get("closed_from_prev",       0) or 0
        pr = d.get("prev_remaining",         0) or 0
        tt = d.get("this_month_total",       0) or 0
        tc = d.get("this_month_completed",   0) or 0
        tu = d.get("this_month_uncompleted", 0) or 0
        mr = d.get("this_month_completion_rate")
        cr = d.get("cum_completion_rate")

        rows.append({
            "月份":           f"{m:02d} 月",
            "上月累計未完成":  str(pv),
            "本月結案上月":    str(cf),
            "上月遙留":        str(pr),
            "本月報修":        str(tt),
            "本月完成":        str(tc),
            "本月未完成":     str(tu),
            "本月完成率":      f"{mr:.1f}%" if mr is not None else "—",
            "累計完成率":      f"{cr:.1f}%" if cr is not None else "—",
        })
        sum_closed_from_prev += cf
        sum_this_total       += tt
        sum_this_completed   += tc
        sum_this_uncompleted += tu

    full_rate = (
        f"{round(sum_this_completed / sum_this_total * 100, 1):.1f}%"
        if sum_this_total else "—"
    )
    rows.append({
        "月份":           "[全年合計]",
        "上月累計未完成":  "—",
        "本月結案上月":    str(sum_closed_from_prev),
        "上月遙留":        "—",
        "本月報修":        str(sum_this_total),
        "本月完成":        str(sum_this_completed),
        "本月未完成":     str(sum_this_uncompleted),
        "本月完成率":      full_rate,
        "累計完成率":      "—",
    })
    return rows


_dazhi_monthly_stats_cols = [
    {"key": "月份",          "label": "月份",           "width": 0.60, "align": "center"},
    {"key": "上月累計未完成", "label": "上月累計未完成", "width": 0.95, "align": "right"},
    {"key": "本月結案上月",   "label": "本月結案(上月)", "width": 0.95, "align": "right"},
    {"key": "上月遙留",       "label": "上月遙留",        "width": 0.85, "align": "right"},
    {"key": "本月報修",       "label": "本月報修",        "width": 0.80, "align": "right"},
    {"key": "本月完成",       "label": "本月完成",        "width": 0.80, "align": "right"},
    {"key": "本月未完成",     "label": "本月未完成",      "width": 0.85, "align": "right"},
    {"key": "本月完成率",     "label": "本月完成率",    "width": 0.90, "align": "right"},
    {"key": "累計完成率",     "label": "累計完成率",    "width": 0.90, "align": "right"},
]

register_section(
    PptSectionDef(
        export_key="dazhi_repair_monthly_stats", module_key=MODULE_KEY,
        tab_name="大直工務部", second_title="年度報修統計",
        description="全年 1~12 月報修指標交叉表（對應前端 3.1 報修 TAB，口徑完全一致）",
        export_type="table", slide_layout="table_full",
        supports_detail=False,
        data_source="backend_db", sort_order=15,
        columns=_dazhi_monthly_stats_cols,
    ),
    data_provider=_provide_dazhi_repair_monthly_stats,
)


# =============================================================================
# 大直工務部 — Section 16：未完成工單（按等待天數降序）
# =============================================================================

def _provide_dazhi_repair_unfinished(db: Session, params: dict) -> list[dict]:
    """所有尚未結案的大直工務部工單，按等待天數降序（最久在最前）。"""
    from app.models.dazhi_repair import DazhiRepairCase
    from app.services.dazhi_repair_service import is_completed, is_excluded

    cases = db.query(DazhiRepairCase).all()
    now   = datetime.now()
    enriched: list = []

    for c in cases:
        if is_excluded(c.status or ""):
            continue
        if is_completed(c.status or ""):
            continue
        if not c.occurred_at:
            continue
        wait = int((now - c.occurred_at).total_seconds() // 86400)
        enriched.append((wait, c))

    enriched.sort(key=lambda x: x[0], reverse=True)

    return [
        {
            "案件編號": c.case_no or "—",
            "報修日期": c.occurred_at.strftime("%Y/%m/%d"),
            "地點":     c.floor    or "—",
            "工作說明": (c.title   or "—")[:25],
            "目前狀態": c.status   or "—",
            "負責人":   c.acceptor or "—",
            "等待天數": f"{wait} 天",
        }
        for wait, c in enriched
    ]


def _provide_dazhi_repair_unfinished_detail(db: Session, params: dict) -> list[dict]:
    """明細附頁：附加委外費用 / 維修費用欄位（排序與主表相同）。"""
    from app.models.dazhi_repair import DazhiRepairCase
    from app.services.dazhi_repair_service import is_completed, is_excluded

    cases = db.query(DazhiRepairCase).all()
    now   = datetime.now()
    enriched: list = []

    for c in cases:
        if is_excluded(c.status or ""):
            continue
        if is_completed(c.status or ""):
            continue
        if not c.occurred_at:
            continue
        wait = int((now - c.occurred_at).total_seconds() // 86400)
        enriched.append((wait, c))

    enriched.sort(key=lambda x: x[0], reverse=True)

    return [
        {
            "案件編號": c.case_no or "—",
            "報修日期": c.occurred_at.strftime("%Y/%m/%d"),
            "地點":     c.floor    or "—",
            "工作說明": (c.title   or "—")[:25],
            "目前狀態": c.status   or "—",
            "負責人":   c.acceptor or "—",
            "等待天數": f"{wait} 天",
            "委外費用": f"${c.outsource_fee:,.0f}"   if c.outsource_fee   else "—",
            "維修費用": f"${c.maintenance_fee:,.0f}" if c.maintenance_fee else "—",
        }
        for wait, c in enriched
    ]


_dazhi_unfinished_cols = [
    {"key": "案件編號", "label": "案件編號", "width": 1.20},
    {"key": "報修日期", "label": "報修日期", "width": 1.10},
    {"key": "地點",     "label": "地點",     "width": 1.50},
    {"key": "工作說明", "label": "工作說明", "width": 2.60},
    {"key": "目前狀態", "label": "狀態",     "width": 1.00},
    {"key": "負責人",   "label": "負責人",   "width": 1.10},
    {"key": "等待天數", "label": "等待天數", "width": 0.90, "align": "right"},
]

register_section(
    PptSectionDef(
        export_key="dazhi_repair_unfinished", module_key=MODULE_KEY,
        tab_name="大直工務部", second_title="大直工務部未完成工單",
        description="所有尚未結案工單，依等待天數由久至近排序（排除取消案件）",
        export_type="table", slide_layout="table_full",
        supports_detail=True,
        detail_description="附加委外費用 / 維修費用欄位",
        data_source="backend_db", sort_order=16,
        columns=_dazhi_unfinished_cols,
    ),
    data_provider=_provide_dazhi_repair_unfinished,
    detail_provider=_provide_dazhi_repair_unfinished_detail,
)


# =============================================================================
# 大直工務部 — Section 17：本月結案工單（含費用）
# =============================================================================

def _provide_dazhi_repair_closed_this_month(db: Session, params: dict) -> list[dict]:
    """本月（簾選年月）已結案工單，主表直接含委外 / 維修費用。"""
    from app.models.dazhi_repair import DazhiRepairCase
    from app.services.dazhi_repair_service import is_completed, is_excluded

    year, month = params.get("year"), params.get("month")
    cases = db.query(DazhiRepairCase).all()
    rows: list[dict] = []

    for c in cases:
        if is_excluded(c.status or ""):
            continue
        if not is_completed(c.status or ""):
            continue
        if not (c.completed_at
                and c.completed_at.year  == year
                and c.completed_at.month == month):
            continue
        rows.append({
            "案件編號": c.case_no or "—",
            "報修日期": c.occurred_at.strftime("%Y/%m/%d") if c.occurred_at else "—",
            "結案日期": c.completed_at.strftime("%Y/%m/%d"),
            "地點":     c.floor    or "—",
            "工作說明": (c.title   or "—")[:22],
            "工時":     f"{c.work_hours:.1f}" if c.work_hours else "—",
            "委外費用": f"${c.outsource_fee:,.0f}"   if c.outsource_fee   else "—",
            "維修費用": f"${c.maintenance_fee:,.0f}" if c.maintenance_fee else "—",
        })
    return rows


def _provide_dazhi_repair_closed_detail(db: Session, params: dict) -> list[dict]:
    """明細附頁：附加扣款費用 / 負責人 / 備註。"""
    from app.models.dazhi_repair import DazhiRepairCase
    from app.services.dazhi_repair_service import is_completed, is_excluded

    year, month = params.get("year"), params.get("month")
    cases = db.query(DazhiRepairCase).all()
    rows: list[dict] = []

    for c in cases:
        if is_excluded(c.status or ""):
            continue
        if not is_completed(c.status or ""):
            continue
        if not (c.completed_at
                and c.completed_at.year  == year
                and c.completed_at.month == month):
            continue
        rows.append({
            "案件編號": c.case_no or "—",
            "報修日期": c.occurred_at.strftime("%Y/%m/%d") if c.occurred_at else "—",
            "結案日期": c.completed_at.strftime("%Y/%m/%d"),
            "地點":     c.floor    or "—",
            "工作說明": (c.title   or "—")[:22],
            "工時":     f"{c.work_hours:.1f}" if c.work_hours else "—",
            "委外費用": f"${c.outsource_fee:,.0f}"   if c.outsource_fee   else "—",
            "維修費用": f"${c.maintenance_fee:,.0f}" if c.maintenance_fee else "—",
            "扣款費用": f"${c.deduction_fee:,.0f}"   if c.deduction_fee   else "—",
            "負責人":   c.acceptor or "—",
            "備註":     (c.finance_note or "")[:20] or "—",
        })
    return rows


_dazhi_closed_cols = [
    {"key": "案件編號", "label": "案件編號", "width": 1.10},
    {"key": "報修日期", "label": "報修日期", "width": 1.00},
    {"key": "結案日期", "label": "結案日期", "width": 1.00},
    {"key": "地點",     "label": "地點",     "width": 1.30},
    {"key": "工作說明", "label": "工作說明", "width": 2.20},
    {"key": "工時",     "label": "工時",     "width": 0.75, "align": "right"},
    {"key": "委外費用", "label": "委外費",   "width": 0.95, "align": "right"},
    {"key": "維修費用", "label": "維修費",   "width": 0.95, "align": "right"},
]

register_section(
    PptSectionDef(
        export_key="dazhi_repair_closed_month", module_key=MODULE_KEY,
        tab_name="大直工務部", second_title="大直工務部本月結案工單",
        description="本月（簾選年月）已結案工單，主表直接顯示委外 / 維修費用",
        export_type="table", slide_layout="table_full",
        supports_detail=True,
        detail_description="附加扣款費用 / 負責人 / 備註欄位",
        data_source="backend_db", sort_order=17,
        columns=_dazhi_closed_cols,
    ),
    data_provider=_provide_dazhi_repair_closed_this_month,
    detail_provider=_provide_dazhi_repair_closed_detail,
)
