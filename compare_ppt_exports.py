# -*- coding: utf-8 -*-
"""
比對腳本：hotel/overview(dazhi) vs mall/overview(luqun) 兩個 PowerPoint 匯出檔
─────────────────────────────────────────────────────────────
用途：純比對，不改任何程式。並排列出兩邊匯出檔的
      ① 投影片順序與標題（表頭順序）
      ② 每張表格的欄位表頭
      ③ 每張表格的資料列數（檢查是否有 0 列＝資料短缺）

執行（本機，能讀到 C:/portal_data/portal.db 的同一個 Python 環境）：
    python compare_ppt_exports.py            # 預設 2026 年 6 月
    python compare_ppt_exports.py 2026 6     # 指定年月

會在 backend/ 產生 cmp_dazhi_<年>.pptx 與 cmp_luqun_<年>.pptx 供開檔檢視；不改原始程式。
"""
import os
import sys
import traceback

try:
    sys.stdout.reconfigure(encoding="utf-8")
except Exception:
    pass

# ── 定位 backend ───────────────────────────────────────────────────────────────
HERE = os.path.dirname(os.path.abspath(__file__))
backend_dir = None
for c in (HERE, os.path.join(HERE, "backend"), os.path.dirname(HERE)):
    if os.path.isdir(os.path.join(c, "app", "routers")):
        backend_dir = c
        break
if backend_dir is None:
    print("[ERROR] 找不到 backend/app，請把本腳本放在 portal 根目錄或 backend/ 下")
    sys.exit(1)
sys.path.insert(0, backend_dir)
os.chdir(backend_dir)

YEAR = int(sys.argv[1]) if len(sys.argv) > 1 else 2026
MONTH = int(sys.argv[2]) if len(sys.argv) > 2 else 6

from app.core.database import SessionLocal          # noqa: E402
from app.routers.repair_ppt_export import _build_repair_pptx  # noqa: E402
from pptx import Presentation                        # noqa: E402


def _slide_title(slide):
    """取投影片標題：選字級最大的非空文字框；同級取最上方。"""
    best, best_sz, best_top = "", -1.0, 1e9
    for sh in slide.shapes:
        if not getattr(sh, "has_text_frame", False):
            continue
        txt = sh.text_frame.text.strip()
        if not txt:
            continue
        sz = 0.0
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                if r.font.size:
                    sz = max(sz, r.font.size.pt)
        top = sh.top or 0
        first = txt.replace("\n", " ").strip()
        if sz > best_sz or (sz == best_sz and top < best_top):
            best, best_sz, best_top = first, sz, top
    return best[:40] if best else "（無標題）"


def _slide_tables(slide):
    """回傳該投影片所有表格的 (表頭list, 資料列數)。"""
    out = []
    for sh in slide.shapes:
        if not getattr(sh, "has_table", False):
            continue
        tbl = sh.table
        headers = [tbl.cell(0, c).text.strip() for c in range(len(tbl.columns))]
        data_rows = len(tbl.rows) - 1
        out.append((headers, data_rows))
    return out


def dump(module_label, module, db):
    print("\n" + "=" * 78)
    print(f"  {module_label}（module='{module}'）  {YEAR}年{MONTH}月")
    print("=" * 78)
    try:
        buf = _build_repair_pptx(module, YEAR, MONTH, db)
    except Exception:
        print(f"[FAIL] 建立 {module} 失敗：")
        traceback.print_exc()
        return []
    out = os.path.join(backend_dir, f"cmp_{module}_{YEAR}.pptx")
    with open(out, "wb") as fh:
        fh.write(buf.getbuffer())
    prs = Presentation(out)
    titles = []
    print(f"輸出檔：{out}　總投影片數 = {len(prs.slides)}\n")
    for i, s in enumerate(prs.slides):
        title = _slide_title(s)
        titles.append(title)
        line = f"  #{i:02d}  {title}"
        tabs = _slide_tables(s)
        if tabs:
            for headers, n in tabs:
                hdr = " | ".join(h for h in headers if h)
                # 表頭過長時截斷顯示
                if len(hdr) > 90:
                    hdr = hdr[:90] + " …"
                line += f"\n        表頭[{n}列]: {hdr}"
        print(line)
    return titles


db = SessionLocal()
try:
    hotel_titles = dump("hotel/overview", "dazhi", db)
    mall_titles = dump("mall/overview", "luqun", db)

    # ── 區段層級差異摘要（去掉分頁重複、去掉「匯出時間…」雜訊）──────────────────
    def norm(t):
        return t.split("　")[0].split("  ")[0].strip()

    def uniq(seq):
        seen, res = set(), []
        for x in seq:
            k = norm(x)
            if k not in seen:
                seen.add(k)
                res.append(k)
        return res

    hset = uniq(hotel_titles)
    mset = uniq(mall_titles)

    print("\n" + "=" * 78)
    print("  區段差異摘要（依標題去重）")
    print("=" * 78)
    print("\n[只在 hotel/overview 出現]")
    for t in hset:
        if t not in mset:
            print("   +", t)
    print("\n[只在 mall/overview 出現]")
    for t in mset:
        if t not in hset:
            print("   +", t)
    print("\n[兩邊都有]")
    for t in hset:
        if t in mset:
            print("   =", t)

    print("\n判讀：『表頭[N列]』的 N=0 代表該表無資料（可能是真的短缺，或該月尚無排程）。"
          "區段差異多為業務領域不同（飯店有客房/IHG；商場有全棟例行維護），非 bug。")
finally:
    db.close()
    print("\n[INFO] 比對結束。")
