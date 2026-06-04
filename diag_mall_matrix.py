# -*- coding: utf-8 -*-
"""
診斷腳本：商場(luqun) 年度計劃表為何沒出現在 PowerPoint
─────────────────────────────────────────────────────────────
用途：釘死「商場週期保養年度計劃表 / 全棟例行維護年度計劃表」沒出現的真正原因，
      區分是「抓到 0 列資料」還是「render 拋錯被 try/except 吞掉」。

執行方式（在本機，能讀到 C:/portal_data/portal.db 的環境）：
    1. 把本檔放在 portal 專案根目錄或 backend/ 下
    2. 啟動專案用的同一個 Python 環境（有裝 fastapi/pptx/sqlalchemy 那個）
    3. 執行：  python diag_mall_matrix.py
       （可選指定年份：python diag_mall_matrix.py 2026）

不會修改任何資料；只讀 DB，並在 backend/ 產生一份 diag_luqun_<年>.pptx 供檢視。
"""
import os
import sys
import traceback
from collections import Counter

# stdout 強制 UTF-8，避免 Windows console 中文/符號亂碼
try:
    sys.stdout.reconfigure(encoding="utf-8")
except Exception:
    pass


def banner(t: str) -> None:
    print("\n" + "=" * 72 + "\n" + t + "\n" + "=" * 72)


# ── 1. 定位 backend（含 app 套件的目錄）────────────────────────────────────────
HERE = os.path.dirname(os.path.abspath(__file__))
CANDIDATES = [HERE, os.path.join(HERE, "backend"), os.path.dirname(HERE)]
backend_dir = None
for c in CANDIDATES:
    if os.path.isdir(os.path.join(c, "app", "routers")):
        backend_dir = c
        break
if backend_dir is None:
    print("[ERROR] 找不到 backend/app 目錄，請把本腳本放在 portal 根目錄或 backend/ 下再執行")
    sys.exit(1)
sys.path.insert(0, backend_dir)
os.chdir(backend_dir)
print(f"[INFO] backend 目錄 = {backend_dir}")

YEAR = int(sys.argv[1]) if len(sys.argv) > 1 else 2026
MONTH = int(sys.argv[2]) if len(sys.argv) > 2 else 6
print(f"[INFO] 診斷年份 = {YEAR}，月份 = {MONTH}")

# ── 2. 匯入 ────────────────────────────────────────────────────────────────────
from app.core.database import SessionLocal, engine  # noqa: E402
from app.routers.repair_ppt_export import (          # noqa: E402
    _get_annual_matrix_rows, _add_annual_matrix_slide,
    _build_repair_pptx, TEMPLATE_PATH,
)
from pptx import Presentation                         # noqa: E402

print(f"[INFO] DATABASE_URL(engine) = {engine.url}")
print(f"[INFO] TEMPLATE_PATH = {TEMPLATE_PATH}  存在={os.path.exists(TEMPLATE_PATH)}")

db = SessionLocal()
try:
    # ── 3. 原始資料筆數 ─────────────────────────────────────────────────────────
    banner("步驟 1：原始資料筆數")
    try:
        from app.models.mall_pm_schedule import MallPMSchedule
        from app.models.full_bldg_pm_schedule import FullBldgPMSchedule
        from app.routers.mall_periodic_maintenance import _mall_get_latest_batch_items
        from app.routers.full_building_maintenance import _fb_get_latest_batch_items

        mall_items = _mall_get_latest_batch_items(db)
        fb_items = _fb_get_latest_batch_items(db)
        mall_sched = (db.query(MallPMSchedule)
                      .filter(MallPMSchedule.year_month.like(f"{YEAR}/%")).count())
        fb_sched = (db.query(FullBldgPMSchedule)
                    .filter(FullBldgPMSchedule.year_month.like(f"{YEAR}/%")).count())
        print(f"商場週期保養 最新批次項目數 _mall_get_latest_batch_items = {len(mall_items)}")
        print(f"全棟例行維護 最新批次項目數 _fb_get_latest_batch_items   = {len(fb_items)}")
        print(f"MallPMSchedule     {YEAR} 年排程筆數 = {mall_sched}")
        print(f"FullBldgPMSchedule {YEAR} 年排程筆數 = {fb_sched}")
    except Exception:
        print("[ERROR] 取原始資料筆數時拋錯：")
        traceback.print_exc()

    # ── 4. _get_annual_matrix_rows 對 luqun / dazhi 的結果 ───────────────────────
    for module in ("luqun", "dazhi"):
        banner(f"步驟 2：_get_annual_matrix_rows('{module}', {YEAR})")
        try:
            rows = _get_annual_matrix_rows(module, YEAR, db)
        except Exception:
            print(f"[FAIL] _get_annual_matrix_rows('{module}') 直接拋錯（這就是被外層 try 吞掉的原因）：")
            traceback.print_exc()
            continue
        print(f"總列數 = {len(rows)}")
        bysrc = Counter(r["source"] for r in rows)
        for src, n in bysrc.items():
            print(f"  source = {src} : {n} 列")
        for src in bysrc:
            sample = [r for r in rows if r["source"] == src][:3]
            print(f"  -- {src} 前 {len(sample)} 列範例 --")
            for r in sample:
                dist = dict(Counter(c["status"] for c in r.get("cells", [])))
                print(f"     類別={r['category']!r} | 項目={r['task_name']!r} | "
                      f"頻率={r['frequency']!r} | 12格狀態分布={dist}")

    # ── 5. 實際嘗試 render 商場兩張表（完全比照 _build_repair_pptx 的非 IHG 設定）─
    banner("步驟 3：實際嘗試產生商場兩張年度計劃表投影片")
    try:
        from app.routers.hotel_overview import _update_cover_date, _delete_slide
        luqun_rows = _get_annual_matrix_rows("luqun", YEAR, db)
        prs = Presentation(TEMPLATE_PATH)
        _update_cover_date(prs.slides[0], YEAR, MONTH)
        _delete_slide(prs, 1)          # 與正式流程一致：刪 TOC，TMPL=1 為內容母片
        TMPL = 1
        SW, SH = prs.slide_width.inches, prs.slide_height.inches

        for title, srcname in [
            ("商場週期保養年度計劃表", "商場週期保養"),
            ("全棟例行維護年度計劃表", "全棟例行維護"),
        ]:
            sub = [r for r in luqun_rows if r["source"] == srcname]
            print(f"\n[{title}] 來源列數 = {len(sub)}")
            n0 = len(prs.slides)
            try:
                _add_annual_matrix_slide(
                    prs, TMPL,
                    title=title,
                    subtitle=f"{YEAR}年　狀態：✓已完成 ○已排定 ✗逾期 △未排定 ?待排",
                    matrix_rows=sub, now_str="diag", SW=SW, SH=SH,
                )
                print(f"  [OK] 成功，新增 {len(prs.slides) - n0} 張投影片")
            except Exception:
                print("  [FAIL] render 拋錯，完整 traceback：")
                traceback.print_exc()
    except Exception:
        print("[ERROR] 步驟 3 前置設定就失敗：")
        traceback.print_exc()

    # ── 6. 完整 build，列出每張投影片標題與關鍵標記位置 ─────────────────────────
    banner("步驟 4：完整 _build_repair_pptx('luqun') 後的投影片清單")
    try:
        buf = _build_repair_pptx("luqun", YEAR, MONTH, db)
        out = os.path.join(backend_dir, f"diag_luqun_{YEAR}.pptx")
        with open(out, "wb") as fh:
            fh.write(buf.getbuffer())
        prs = Presentation(out)
        print(f"輸出檔：{out}")
        print(f"總投影片數 = {len(prs.slides)}")
        MARKERS = ("商場週期保養年度計劃表", "全棟例行維護年度計劃表", "未完成附表")
        for i, s in enumerate(prs.slides):
            alltext = " ".join(
                sh.text_frame.text for sh in s.shapes
                if getattr(sh, "has_text_frame", False) and sh.text_frame.text.strip()
            ).replace("\n", " ")
            head = alltext.strip()[:46]
            hit = [m for m in MARKERS if m in alltext]
            flag = ("   <<< " + " / ".join(hit)) if hit else ""
            print(f"  #{i:02d}: {head}{flag}")
        print("\n判讀：若上面看不到『商場週期保養年度計劃表』『全棟例行維護年度計劃表』兩個標記，"
              "且步驟 2 列數 > 0、步驟 3 顯示 [OK]，"
              "代表問題出在正式流程的內層 try/except；請把本輸出整段回傳。")
    except Exception:
        print("[FAIL] 完整 build 失敗：")
        traceback.print_exc()
finally:
    db.close()
    print("\n[INFO] 診斷結束。")
